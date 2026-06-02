#!/usr/bin/env python3
"""
table_renderer.py — native PowerPoint table (zebra-style) для slide_type table_native.

Зачем: до v1.8 скилл не умел рисовать настоящие PowerPoint таблицы. Регулярные
таблицы делались либо через donor 53/54 (с PNG-заглушкой), либо как набор плашек
flow_diagram_native (что лишало пользователя редактирования таблиц в PowerPoint).

Решение: native PPT table через `slide.shapes.add_table(rows, cols, ...)`, со
стилем slide 56 шаблона Cloud.ru:
  - Header row: белый фон, текст SemiBold #222222
  - Body rows: чередуются #F2F2F2 (серый) / белый (zebra)
  - Vertical separators между колонками: 0.5pt #C8C8C8 (только справа от ячеек,
    кроме последней колонки в каждом ряду)
  - НЕТ горизонтальных границ между строками — зебра-фон сам создаёт разделение
  - Padding в ячейках: L/R 12px, T/B 8px (canonical: 12/8)
  - Шрифт: SB Sans Display. Размер — дефолт 16pt с авто-уменьшением до 12pt
    (Problem #5), header SemiBold того же размера. Форс: table.font_size.
  - Цвет текста: #222222 (графит)
  - Выравнивание текста: **align=left, vanchor=top** (canonical правило)
  - Первая колонка может быть шире (для row labels / категорий) через
    `first_col_wider: true` (default).

Используется build_v9 при slide_type == "table_native".

Ограничения v1.8:
  - Регулярные таблицы (M cols × N rows, без merged cells).
  - Если в исходнике обнаружены merged cells / irregular layout — Slide Classifier
    должен СПРОСИТЬ пользователя (anti-distortion правило, см.
    feedback_anti_distortion_safety.md) и предложить альтернативу
    (flow_diagram_native).

Config schema (передаётся через plan.json):
  {
    "slide_type": "table_native",
    "dark": false,
    "table": {
      "header": "Заголовок слайда",         # обязательное
      "subtitle": "...",                     # опц., 11pt под header'ом
      "style": "zebra",                      # default (и единственный для v1.8)
      "headers": ["Категория", "Параметр 1", "Параметр 2"],  # row 0, обязательно
      "data": [                              # rows ≥ 1
        ["Row 1 label", "val", "val"],
        ["Row 2 label", "val", "val"]
      ],
      "first_col_wider": true,               # опц., default true — первая шире (1.4x)
      "x": 30, "y": 180,                     # опц., default safe-area
      "w": 1220, "h": null,                  # опц., w=safe-width, h=auto от rows
      "header_height": 50,                   # опц., default 50
      "row_height": null,                    # опц., default auto fill (h / rows)
      "borders": {                           # опц., гибкое управление границами
        "vertical": true,                     # внутренние вертикали между колонками
        "horizontal": false,                  # внутренние горизонтали между строками
        "outer_top": false,                   # внешняя верхняя
        "outer_bottom": false,                # внешняя нижняя
        "outer_left": false,                  # внешняя левая
        "outer_right": false,                 # внешняя правая
        "color": "#434343",                   # цвет (default #434343)
        "width_pt": 1.0                       # толщина в pt (default 1.0)
      }
    }
  }

Default `borders` (по правилу 2026-05-26 — slide 56 zebra style):
  - Только vertical=true (внутренние вертикали)
  - Всё остальное false
"""
import os
import json

from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


EMU = 9525
SLIDE_W_PX = 1280
SLIDE_H_PX = 720

# Safe-area (canonical v1.7+, та же что в flow_renderer)
SAFE_TOP = 140
SAFE_BOTTOM = 660
SAFE_LEFT = 35
SAFE_RIGHT = 1245
SAFE_W = SAFE_RIGHT - SAFE_LEFT     # 1210
SAFE_H = SAFE_BOTTOM - SAFE_TOP     # 520


# ============================================================================
# Палитра (из brand/palette.json)
# ============================================================================
def _load_palette():
    here = os.path.dirname(os.path.abspath(__file__))
    for path in (
        os.path.join(here, "..", "brand", "palette.json"),
        os.path.join(os.getcwd(), "brand", "palette.json"),
        os.path.join(os.getcwd(), "pptx-skill", "brand", "palette.json"),
    ):
        try:
            if os.path.isfile(path):
                with open(path, encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            continue
    return None


_PAL = _load_palette()


def _hex(hx):
    h = hx.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _from_palette(section, name, fallback):
    if _PAL and name in _PAL.get(section, {}):
        return _hex(_PAL[section][name])
    return _hex(fallback)


GRAPHITE = _from_palette("base", "Black", "#222222")
WHITE = _from_palette("base", "White", "#FFFFFF")
GRAY = _from_palette("base", "Gray", "#F2F2F2")
GREEN = _from_palette("base", "Green", "#26D07C")
DARK_FILL = GRAPHITE

# Canonical (правило 2026-05-26): vertical separators между колонками таблицы =
# 1pt #434343. Никаких внешних границ и горизонтальных линий — зебра-фон сам
# создаёт визуальное разделение.
# NB: 0.5pt в PowerPoint/Keynote часто плохо видны (anti-aliasing), поэтому 1pt.
SEPARATOR_COLOR = RGBColor(0x43, 0x43, 0x43)
SEPARATOR_WIDTH_PT = 1.0

# Native PP bullets (canonical 2026-05-26, исправлено 2026-05-29):
# Если строка начинается с одного из этих маркеров — штатный PP «Заполненный
# квадрат» (Wingdings §, цвет #434343). Текст пишется БЕЗ маркера в самой строке.
# § в Wingdings = плоский залитый квадрат, без emoji-фолбэка и «эффекта объёма».
BULLET_PREFIXES = ("▪ ", "■ ", "• ", "- ", "* ")
BULLET_FONT = "Wingdings"
BULLET_CHAR = "§"
BULLET_COLOR = "434343"
BULLET_INDENT_EMU = 228600


def _detect_bullet_table(line):
    """Возвращает (is_bullet, clean_text)."""
    for prefix in BULLET_PREFIXES:
        if line.startswith(prefix):
            return True, line[len(prefix):]
    return False, line


def _apply_bullet_to_paragraph_table(p, char=BULLET_CHAR, indent_emu=BULLET_INDENT_EMU):
    """Добавить штатный PP bullet «Заполненный квадрат» к параграфу ячейки.

    Порядок дочерних элементов pPr по схеме OOXML: buClr → buFont → buChar.
    """
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(indent_emu))
    pPr.set("indent", str(-indent_emu))
    for tag in ("buClr", "buSzPct", "buNone", "buChar", "buAutoNum", "buFont"):
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    srgb = etree.SubElement(buClr, qn("a:srgbClr"))
    srgb.set("val", BULLET_COLOR)
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", BULLET_FONT)
    buFont.set("pitchFamily", "2")
    buFont.set("charset", "2")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)

FONT = "SB Sans Display"
# Полужирное = отдельный font face (встроен в шаблон), НЕ bold-флаг.
# Canonical 2026-05-29 (Problem #3): Bold запрещён — эмфаза только через SemiBold.
FONT_SEMIBOLD = "SB Sans Display Semibold"


def _set_weight(font, semibold):
    """Эмфаза через начертание SemiBold, а не bold-флаг (Problem #3)."""
    font.name = FONT_SEMIBOLD if semibold else FONT
    font.bold = False


# ============================================================================
# Утилиты для XML cell properties
# ============================================================================
def _set_cell_fill(cell, rgb):
    """Установить fill ячейки. None → noFill (прозрачная)."""
    if rgb is None:
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb


def _set_cell_borders(cell, left=False, right=False,
                       top=False, bottom=False,
                       color_rgb=None, w_pt=None):
    """Установить границы ячейки. Все четыре стороны явно — для надёжности.

    Canonical (2026-05-26): vertical separator между колонками = 1pt #434343,
    рисуется НА ОБЕИХ соседних ячейках (lnR cell[c] + lnL cell[c+1]) — это
    защита от z-order issue, когда заливка соседней ячейки перекрывает мой
    границу. Никаких внешних рамок и горизонтальных линий.
    """
    rgb = color_rgb if color_rgb is not None else SEPARATOR_COLOR
    width = w_pt if w_pt is not None else SEPARATOR_WIDTH_PT
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))

    # Удалить существующие lnL/lnR/lnT/lnB
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        for el in tcPr.findall(qn(f"a:{tag}")):
            tcPr.remove(el)

    def _add_border(tag, on):
        ln = etree.SubElement(tcPr, qn(f"a:{tag}"))
        ln.set("w", str(int(width * 12700)))
        # cap=flat и algn=ctr — как в оригинальном slide 56 шаблона.
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        if on:
            solidFill = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
            prstDash = etree.SubElement(ln, qn("a:prstDash"))
            prstDash.set("val", "solid")
            etree.SubElement(ln, qn("a:round"))
        else:
            etree.SubElement(ln, qn("a:noFill"))

    # PowerPoint порядок: lnL, lnR, lnT, lnB.
    _add_border("lnL", left)
    _add_border("lnR", right)
    _add_border("lnT", top)
    _add_border("lnB", bottom)


def _set_cell_margins(cell, left_px=12, right_px=12, top_px=8, bottom_px=8):
    """Поля внутри ячейки (canonical v1.8: L/R 12px, T/B 8px)."""
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
    tcPr.set("marL", str(left_px * EMU))
    tcPr.set("marR", str(right_px * EMU))
    tcPr.set("marT", str(top_px * EMU))
    tcPr.set("marB", str(bottom_px * EMU))


def _set_cell_text(cell, text, size_pt=12, bold=False, color_rgb=None):
    """Записать текст в ячейку с canonical стилями (left + top, font SB Sans).

    Поддержка multi-line + native PP bullets:
    - Если text содержит '\\n' — разбиваем на параграфы.
    - Если строка начинается с '▪ ' / '■ ' / '• ' / '- ' / '* ' → применяем
      native PP bullet (auto-detect). Префикс удаляется из видимого текста.
    """
    tf = cell.text_frame
    tf.clear()
    cell.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True

    text_str = str(text) if text is not None else ""
    lines = text_str.split("\n")

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Auto-detect bullet
        is_bullet, clean_line = _detect_bullet_table(line)
        if is_bullet:
            _apply_bullet_to_paragraph_table(p)
        run = p.add_run()
        run.text = clean_line
        run.font.size = Pt(size_pt)
        _set_weight(run.font, bold)
        run.font.color.rgb = color_rgb if color_rgb is not None else GRAPHITE


# Размер шрифта в таблице (Problem #5, 2026-05-29): дефолт 16pt, авто-уменьшение
# вниз до 12pt (комфортный минимум), 10pt — крайняя перегрузка.
TABLE_FONT_DEFAULT_PT = 16
TABLE_FONT_COMFORT_MIN_PT = 12
TABLE_FONT_HARD_MIN_PT = 10
_CELL_MARGIN_X_PX = 24   # 12 + 12
_CELL_MARGIN_V_PX = 16   # 8 + 8


def _estimate_cell_lines(text, usable_w_px, char_w_px):
    """Оценка числа строк, на которые перенесётся текст в ячейке."""
    if text in (None, ""):
        return 1
    cpl = max(1, int(usable_w_px / char_w_px))  # символов в строку
    lines = 0
    for seg in str(text).split("\n"):
        seg_len = len(seg) if seg else 1
        lines += max(1, -(-seg_len // cpl))  # ceil деление
    return max(1, lines)


def _autofit_table_font(headers, data, col_widths, row_h, header_h,
                        start_pt=TABLE_FONT_DEFAULT_PT,
                        hard_min=TABLE_FONT_HARD_MIN_PT):
    """Подбирает крупнейший размер шрифта (≤ start_pt), при котором текст всех
    ячеек влезает по ширине/высоте. Старт 16pt, шаг вниз до hard_min.

    Эвристика по ширине символа SB Sans (~0.6em) и межстрочному (~1.25).
    """
    rows = [headers] + list(data)
    for fpt in range(int(start_pt), int(hard_min) - 1, -1):
        font_px = fpt * 4.0 / 3.0          # pt → px
        char_w = 0.60 * font_px            # средняя ширина символа
        line_h = 1.25 * font_px            # высота строки с интерлиньяжем
        fits = True
        for r_i, row in enumerate(rows):
            avail_h = (header_h if r_i == 0 else row_h) - _CELL_MARGIN_V_PX
            for c_i, cell in enumerate(row):
                usable_w = col_widths[c_i] - _CELL_MARGIN_X_PX
                lines = _estimate_cell_lines(cell, usable_w, char_w)
                if lines * line_h > avail_h:
                    fits = False
                    break
            if not fits:
                break
        if fits:
            return fpt
    return int(hard_min)


def _strip_default_table_style(table):
    """Заменить дефолтный applied table style (Medium Style 1 / etc.) на
    built-in «No Style, No Grid» (PowerPoint GUID {2D5ABB26-...}).

    Если просто УДАЛИТЬ tableStyleId — PowerPoint остаётся без applied style и
    может игнорировать cell-level borders. «No Style, No Grid» — это пустой
    встроенный стиль: без своих fills/borders, но позволяет cell-level overrides.

    Также удаляем атрибуты firstRow/bandRow и т.п. — мы делаем zebra вручную.
    """
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        # Убираем firstRow / bandRow / etc. — мы делаем стили вручную
        for attr in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
            if tblPr.get(attr):
                del tblPr.attrib[attr]
        # Заменить tableStyleId на «No Style, No Grid»
        for sid in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(sid)
        no_style_id = etree.SubElement(tblPr, qn("a:tableStyleId"))
        no_style_id.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


# ============================================================================
# Header слайда + subtitle (общая часть для table_native)
# ============================================================================
def _add_slide_header(slide, text, dark=False):
    """Заголовок слайда — в штатный TITLE-placeholder шаблона (canonical 35,38 /
    963×54 / 20pt SemiBold CAPS). Problem #6, 2026-05-29.
    Делегирует в общий set_slide_title (с fallback на textbox)."""
    from kpi_renderer import set_slide_title
    return set_slide_title(slide, text, dark=dark)


def _add_top_separator(slide, y=110):
    """Тонкая серая линия под заголовком."""
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(35 * EMU), Emu(y * EMU), Emu(1245 * EMU), Emu(y * EMU)
    )
    conn.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    conn.line.width = Emu(int(0.5 * 12700))


def _add_subtitle(slide, text, dark=False):
    """Subtitle 11pt под header'ом."""
    box = slide.shapes.add_textbox(
        Emu(35 * EMU), Emu(122 * EMU), Emu(1200 * EMU), Emu(22 * EMU)
    )
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.bold = False
    run.font.color.rgb = WHITE if dark else GRAPHITE


# ============================================================================
# Пресет before/after (архетип 8) — рисуется из плашек, не из PPT-таблицы
# ============================================================================
def _ba_rect(slide, x, y, w, h, fill_rgb):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Emu(int(x) * EMU), Emu(int(y) * EMU),
                                  Emu(int(w) * EMU), Emu(int(h) * EMU))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_rgb
    rect.line.fill.background()
    # Снять любые эффекты — и явные, и тематические (effectRef). Плоско, без исключений.
    from effects_util import strip_effects
    strip_effects(rect._element)
    return rect


def _ba_text(slide, x, y, w, h, text, size_pt, bold,
             color_rgb, align="left", anchor="middle"):
    box = slide.shapes.add_textbox(
        Emu(int(x) * EMU), Emu(int(y) * EMU),
        Emu(int(w) * EMU), Emu(int(h) * EMU))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor == "middle" else MSO_ANCHOR.TOP
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    align_enum = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                  "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_enum
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        _set_weight(run.font, bold)
        run.font.color.rgb = color_rgb
    return box


def _ba_multitext(slide, x, y, w, h, parts, align="left", anchor="middle"):
    """Текстбокс с несколькими абзацами, у каждого свой (size, bold, color).
    parts: [(text, size_pt, bold, color_rgb), ...]. Для двухстрочной метрики
    (название SemiBold + деталь regular серым)."""
    box = slide.shapes.add_textbox(
        Emu(int(x) * EMU), Emu(int(y) * EMU),
        Emu(int(w) * EMU), Emu(int(h) * EMU))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor == "middle" else MSO_ANCHOR.TOP
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    align_enum = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                  "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    for i, (text, size_pt, bold, color_rgb) in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_enum
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size_pt)
        _set_weight(run.font, bold)
        run.font.color.rgb = color_rgb
    return box


def render_before_after(slide, cfg, dark=False):
    """Архетип 8 — before/after. Сверено по референсу 2026-06-01 (EVOLUTION ML
    INFERENCE VS ON-PREM): три колонки — метрика | «Было» | «Стало». Заголовки-
    табы ЦВЕТНЫЕ (серый / зелёный) с текстом ПО ЛЕВОМУ краю, на всю ширину колонки.
    Тело — без сплошных плашек; лёгкая серая подложка под колонкой «Было», тонкие
    серые разделители строк через все колонки. Метрика — двухстрочная (название
    SemiBold графит + деталь regular серым, через "\\n").

    cfg["rows"]: [{"metric","before","after"}], cfg["before_label"]/["after_label"],
    cfg["metric_label"] (заголовок колонки метрик, опц.).
    """
    rows = cfg.get("rows", [])
    n = len(rows)
    if n == 0:
        return
    before_label = cfg.get("before_label", "Было")
    after_label = cfg.get("after_label", "Стало")
    metric_label = cfg.get("metric_label", "")

    top = cfg.get("content_top", SAFE_TOP + 6)
    tab_h = cfg.get("tab_h", 44)
    body_top = top + tab_h + 10
    body_bottom = SAFE_BOTTOM
    body_h = body_bottom - body_top

    label_w = int(SAFE_W * cfg.get("label_ratio", 0.30))
    col_gap = cfg.get("col_gap", 4)           # ≤10px, базово 4px (правило 2026-06-01)
    rest = SAFE_W - label_w - col_gap
    before_w = rest // 2
    after_w = rest - before_w
    before_x = SAFE_LEFT + label_w
    after_x = before_x + before_w + col_gap

    gray = DARK_FILL if dark else GRAY
    base_txt = WHITE if dark else GRAPHITE
    body_gray = RGBColor(0xCF, 0xCF, 0xCF) if dark else RGBColor(0x5C, 0x5C, 0x5C)
    divider = RGBColor(0x44, 0x44, 0x44) if dark else RGBColor(0xCC, 0xCC, 0xCC)
    faint = RGBColor(0x2A, 0x2A, 0x2A) if dark else RGBColor(0xF7, 0xF7, 0xF7)
    pad = 14

    # Лёгкая подложка под колонкой «Было» (едва заметная — выделяет столбец).
    _ba_rect(slide, before_x, body_top, before_w, body_h, faint)

    # Заголовки-табы — цветные, текст ПО ЛЕВОМУ краю, во всю ширину колонки.
    if metric_label:
        _ba_text(slide, SAFE_LEFT, top, label_w - pad, tab_h, metric_label,
                 cfg.get("tab_font", 16), True, base_txt, align="left")
    _ba_rect(slide, before_x, top, before_w, tab_h, gray)
    _ba_rect(slide, after_x, top, after_w, tab_h, GREEN)
    _ba_text(slide, before_x + pad, top, before_w - 2 * pad, tab_h, before_label,
             cfg.get("tab_font", 16), True, base_txt, align="left")
    _ba_text(slide, after_x + pad, top, after_w - 2 * pad, tab_h, after_label,
             cfg.get("tab_font", 16), True, GRAPHITE, align="left")

    # Строки.
    band_h = body_h / n
    metric_size = cfg.get("metric_size", 16)
    before_size = cfg.get("before_size", 15)
    after_size = cfg.get("after_size", 15)
    from pptx.enum.shapes import MSO_CONNECTOR
    for i, row in enumerate(rows):
        y = body_top + i * band_h
        # Тонкий серый разделитель сверху строки (кроме первой), через все колонки.
        if i > 0:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(SAFE_LEFT) * EMU), Emu(int(y) * EMU),
                Emu(int(after_x + after_w) * EMU), Emu(int(y) * EMU))
            conn.line.color.rgb = divider
            conn.line.width = Emu(int(0.75 * 12700))
        # Метрика — двухстрочная: название SemiBold + деталь regular серым.
        metric = str(row.get("metric", ""))
        mparts = metric.split("\n", 1)
        parts = [(mparts[0], metric_size, True, base_txt)]
        if len(mparts) > 1:
            parts.append((mparts[1], metric_size - 2, False, body_gray))
        _ba_multitext(slide, SAFE_LEFT, y, label_w - pad, band_h, parts,
                      align="left", anchor="middle")
        # «Было» — серым regular.
        _ba_text(slide, before_x + pad, y, before_w - 2 * pad, band_h,
                 row.get("before", ""), before_size, False, body_gray,
                 align="left")
        # «Стало» — графит regular (читается как «после»).
        _ba_text(slide, after_x + pad, y, after_w - 2 * pad, band_h,
                 row.get("after", ""), after_size, False, base_txt,
                 align="left")


# ============================================================================
# Главная функция
# ============================================================================
def render_table_native(slide, table_config, dark=False):
    """Собирает native PowerPoint таблицу в стиле slide 56 (zebra).

    Предполагается, что slide уже прошёл clean_slide_to_blank().
    """
    if not isinstance(table_config, dict):
        raise ValueError("table_config должен быть dict")

    # Preset-архетип before/after (design-principles-from-decks.md, архетип 8).
    # Это НЕ обычная zebra-таблица — рисуется отдельной композицией из плашек.
    if table_config.get("preset") == "before_after":
        header_text = table_config.get("header", "")
        if header_text:
            _add_slide_header(slide, header_text, dark=dark)
            if table_config.get("top_separator"):
                _add_top_separator(slide)
        subtitle = table_config.get("subtitle")
        if subtitle:
            _add_subtitle(slide, subtitle, dark=dark)
        # Отступ тела под подзаголовок (правило 2026-06-01).
        if "content_top" not in table_config:
            table_config["content_top"] = 172 if subtitle else SAFE_TOP + 6
        render_before_after(slide, table_config, dark=dark)
        return

    headers = table_config.get("headers", [])
    data = table_config.get("data", [])

    if not headers:
        raise ValueError("table_config.headers — пусто (нужна шапка)")
    if not data:
        raise ValueError("table_config.data — пусто (нужно ≥1 строка)")

    n_cols = len(headers)
    n_data_rows = len(data)
    n_rows_total = 1 + n_data_rows  # header + data

    # Проверка: все data rows той же длины что и headers
    for i, row in enumerate(data):
        if len(row) != n_cols:
            raise ValueError(
                f"table_config.data[{i}]: {len(row)} ячеек, ожидалось {n_cols}"
            )

    # 1. Header слайда
    header_text = table_config.get("header", "")
    if header_text:
        _add_slide_header(slide, header_text, dark=dark)
        # Дивайдер под заголовком — только по явному флагу (правило 2026-06-01).
        if table_config.get("top_separator"):
            _add_top_separator(slide)

    subtitle = table_config.get("subtitle")
    if subtitle:
        _add_subtitle(slide, subtitle, dark=dark)

    # 2. Расчёт размеров таблицы
    table_x = table_config.get("x", SAFE_LEFT)
    table_y = table_config.get("y", 170 if subtitle else 160)
    table_w = table_config.get("w", SAFE_W)
    table_h = table_config.get("h")  # может быть None — посчитаем

    header_h = table_config.get("header_height", 50)
    row_h = table_config.get("row_height")

    if table_h is None and row_h is None:
        # Заполнить доступную safe-area
        avail_h = SAFE_BOTTOM - table_y - 20  # 20 px зазор перед footer
        row_h = max(40, (avail_h - header_h) // n_data_rows)
        table_h = header_h + row_h * n_data_rows
    elif table_h is None:
        # row_h задан
        table_h = header_h + row_h * n_data_rows
    elif row_h is None:
        row_h = max(30, (table_h - header_h) // n_data_rows)

    # 3. Распределение ширин колонок
    first_col_wider = table_config.get("first_col_wider", True)
    if first_col_wider and n_cols >= 2:
        # Первая колонка 1.4× от остальных
        # x = first_w + (n-1) * rest_w; first_w = 1.4 * rest_w
        # → rest_w = total / (1.4 + n - 1)
        rest_w = table_w / (1.4 + n_cols - 1)
        first_w = int(round(1.4 * rest_w))
        rest_w = int(round(rest_w))
        col_widths = [first_w] + [rest_w] * (n_cols - 1)
        # Скорректировать чтобы сумма = table_w
        col_widths[-1] = table_w - sum(col_widths[:-1])
    else:
        # Равномерно
        col_widths = [table_w // n_cols] * n_cols
        col_widths[-1] = table_w - sum(col_widths[:-1])

    # 4. Создание таблицы
    table_shape = slide.shapes.add_table(
        n_rows_total, n_cols,
        Emu(table_x * EMU), Emu(table_y * EMU),
        Emu(table_w * EMU), Emu(table_h * EMU),
    )
    table = table_shape.table
    _strip_default_table_style(table)

    # 5. Установить ширины колонок
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w * EMU)

    # 6. Установить высоты строк
    table.rows[0].height = Emu(header_h * EMU)
    for r in range(1, n_rows_total):
        table.rows[r].height = Emu(row_h * EMU)

    # Все cell-level borders ВЫКЛЮЧЕНЫ — границы будем рисовать отдельными
    # connector lines поверх таблицы (см. шаг 9). Это обеспечивает гарантированную
    # visibility во всех приложениях (PowerPoint Mac/Win, Keynote, LibreOffice),
    # независимо от applied table style.
    no_borders = {"left": False, "right": False, "top": False, "bottom": False}

    # Размер шрифта (Problem #5): дефолт 16pt, авто-уменьшение пока не влезает
    # (до 12pt комфортно, 10pt — крайняя перегрузка). Можно форсировать
    # table_config["font_size"]. Header и body — один размер, header SemiBold.
    forced = table_config.get("font_size")
    if forced:
        body_font = int(forced)
    else:
        body_font = _autofit_table_font(headers, data, col_widths, row_h, header_h)
    header_font = body_font

    # 7. Заполнить header row (row 0)
    for c_idx, header_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        _set_cell_fill(cell, None)  # noFill — белая/прозрачная
        _set_cell_margins(cell, left_px=12, right_px=12, top_px=8, bottom_px=8)
        _set_cell_borders(cell, **no_borders)
        _set_cell_text(cell, header_text, size_pt=header_font, bold=True)

    # 8. Заполнить data rows (row 1..N) — zebra
    for r_idx, row_data in enumerate(data):
        row_num = r_idx + 1  # в таблице (с учётом header)
        # zebra: первая data row серая, вторая белая, и т.д.
        is_even_data_row = (r_idx % 2 == 0)
        fill_rgb = GRAY if is_even_data_row else None  # None = noFill = белая

        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_num, c_idx)
            _set_cell_fill(cell, fill_rgb)
            _set_cell_margins(cell, left_px=12, right_px=12, top_px=8, bottom_px=8)
            _set_cell_borders(cell, **no_borders)
            _set_cell_text(cell, cell_text, size_pt=body_font, bold=False)

    # 9. Нарисовать границы как отдельные connector lines поверх таблицы.
    # Это гарантирует visibility во всех приложениях.
    _draw_table_borders_as_lines(
        slide, table_config,
        table_x, table_y, table_w, table_h,
        col_widths, header_h, row_h, n_data_rows
    )

    return table_shape


def _draw_table_borders_as_lines(slide, table_config,
                                  table_x, table_y, table_w, table_h,
                                  col_widths, header_h, row_h, n_data_rows):
    """Нарисовать границы как отдельные line shapes поверх таблицы.

    Управление через `table_config['borders']` dict:
      - vertical: bool — внутренние вертикали между колонками
      - horizontal: bool — внутренние горизонтали между строками
      - outer_top / outer_bottom / outer_left / outer_right: bool — внешние
      - color: hex (#434343)
      - width_pt: float (1.0)

    Default: только vertical=true (canonical правило 2026-05-26).
    """
    from pptx.enum.shapes import MSO_CONNECTOR

    borders = table_config.get("borders", {})

    # Defaults
    vertical = borders.get("vertical", True)
    horizontal = borders.get("horizontal", False)
    outer_top = borders.get("outer_top", False)
    outer_bottom = borders.get("outer_bottom", False)
    outer_left = borders.get("outer_left", False)
    outer_right = borders.get("outer_right", False)
    color_hex = borders.get("color", "#434343")
    width_pt = borders.get("width_pt", 1.0)

    color_rgb = _hex(color_hex) if isinstance(color_hex, str) else color_hex
    line_w_emu = int(width_pt * 12700)

    def _line(x1, y1, x2, y2):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(int(x1) * EMU), Emu(int(y1) * EMU),
            Emu(int(x2) * EMU), Emu(int(y2) * EMU)
        )
        conn.line.color.rgb = color_rgb
        conn.line.width = Emu(line_w_emu)
        return conn

    # X-координаты границ колонок (включая левую и правую внешние)
    col_x = [table_x]
    for w in col_widths:
        col_x.append(col_x[-1] + w)
    # col_x = [left, after_col1, after_col2, ..., right]

    # Y-координаты границ строк (header + data rows)
    row_y = [table_y, table_y + header_h]
    for _ in range(n_data_rows):
        row_y.append(row_y[-1] + row_h)
    # row_y = [top, after_header, after_data_row1, ..., bottom]

    table_bottom = row_y[-1]
    table_right = col_x[-1]

    # Внутренние вертикали (между колонками)
    if vertical:
        for i in range(1, len(col_x) - 1):
            x = col_x[i]
            _line(x, table_y, x, table_bottom)

    # Внутренние горизонтали (между строками)
    if horizontal:
        for i in range(1, len(row_y) - 1):
            y = row_y[i]
            _line(table_x, y, table_right, y)

    # Внешние границы
    if outer_top:
        _line(table_x, table_y, table_right, table_y)
    if outer_bottom:
        _line(table_x, table_bottom, table_right, table_bottom)
    if outer_left:
        _line(table_x, table_y, table_x, table_bottom)
    if outer_right:
        _line(table_right, table_y, table_right, table_bottom)


# ============================================================================
# CLI для standalone-теста
# ============================================================================
def _cli_main():
    """python3 table_renderer.py <config.json> <template.pptx> <out.pptx>"""
    import sys
    from pptx import Presentation
    from pptx.oxml.ns import qn

    if len(sys.argv) < 4:
        print("Usage: table_renderer.py <config.json> <template.pptx> <out.pptx>",
              file=sys.stderr)
        sys.exit(1)

    cfg_path, tpl_path, out_path = sys.argv[1], sys.argv[2], sys.argv[3]
    with open(cfg_path, encoding="utf-8") as f:
        cfg = json.load(f)

    here = os.path.dirname(os.path.abspath(__file__))
    import sys as _sys
    _sys.path.insert(0, here)
    from kpi_renderer import clean_slide_to_blank, BLANK_DONOR_WHITE, BLANK_DONOR_DARK
    from build_v9 import clone_slide

    dark = cfg.get("dark", False)
    table_cfg = cfg.get("table", cfg)

    p = Presentation(tpl_path)
    originals = list(p.slides)
    blank_idx = BLANK_DONOR_DARK if dark else BLANK_DONOR_WHITE
    new_slide = clone_slide(p, originals[blank_idx - 1])

    sldIdLst = p.slides._sldIdLst
    for sid in list(sldIdLst)[:len(originals)]:
        rid = sid.attrib[qn("r:id")]
        try:
            p.part.drop_rel(rid)
        except Exception:
            pass
        sldIdLst.remove(sid)

    clean_slide_to_blank(new_slide)
    render_table_native(new_slide, table_cfg, dark=dark)
    p.save(out_path)
    print(f"Saved {out_path}", file=sys.stderr)


if __name__ == "__main__":
    _cli_main()
