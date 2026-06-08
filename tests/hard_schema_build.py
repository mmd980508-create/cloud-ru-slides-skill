#!/usr/bin/env python3
"""Хард-схема Enterprise Router + Cloud Firewall (v2.5.1) — дисциплина сетки.

Принцип порядка: НЕ координаты «на глаз», а единая система:
  - единые шрифты/высоты/поля для всех однотипных таблиц;
  - две горизонтальные «линии чтения» (центры рядов C1/C2) — таблицы и иконки
    КАЖДОЙ зоны центрируются по ним, поэтому зоны читаются строками;
  - высота таблицы считается ТОЙ ЖЕ формулой, что в flow_renderer → центрирование
    точное;
  - зоны-панели общей высоты/верха/низа (146…594).
Запуск: /usr/bin/python3 tests/hard_schema_build.py
"""
import math
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
SCRIPTS = os.path.join(os.path.dirname(HERE), "scripts")
sys.path.insert(0, SCRIPTS)

from pptx import Presentation                       # noqa: E402
from pptx.oxml.ns import qn                          # noqa: E402
import flow_renderer as fr                           # noqa: E402
from kpi_renderer import (clean_slide_to_blank,      # noqa: E402
                          BLANK_DONOR_WHITE)
from build_v9 import clone_slide                     # noqa: E402
from template_path import resolve_template            # noqa: E402

# ---- палитра + СЛОИ ФРЕЙМОВ (бренд-спека 2026-06-05) -----------------------
# Чередование заливок по вложенности:
#   Background = Gray 100  →  L1 = White  →  L2 = Gray 100 + нижний штрих graphite-50
#   →  L3 = пунктир graphite-50 (единственный пунктирный слой).
GRAY100 = "#F2F2F2"        # фон (Background) и L2-заливка (security group)
WHITE = "#FFFFFF"          # L1-заливка (VPC-зоны)
BAND_GRAY = "#8A8A8A"      # серый бэнд заголовка таблиц
GREEN = "#26D07C"          # акцент: нижний штрих компонент-узлов (ECS/Firewall)
STROKE = "#D3D3D3"         # graphite-50: нижний штрих L2 + пунктир L3
BS_GRAY = {"color": STROKE, "w_pt": 1.0}   # bottom stroke L2
BS_GREEN = {"color": GREEN, "w_pt": 1.0}   # bottom stroke компонента (акцент)

# ---- единая типографика таблиц (консистентность однотипных элементов) ------
TF, HF, DF = 10, 9, 9      # title / header / data font pt — ЕДИНЫЕ везде
TMIN, HMIN, RMIN = 14, 14, 14

# ---- ОСИ выравнивания (воображаемые; в макете их нет) ----------------------
HEAD_Y = 150               # горизонталь заголовков всех зон верхнего ряда
PANEL_TOP, PANEL_BOT = 146, 594
T1, T2 = 246, 486          # ВЕРХ таблиц двух рядов (плашки-заголовки на одной линии)
ARROW1, ARROW2 = 300, 540  # горизонтали трафика (request / response)
# Вертикали — края зон/таблиц (после snap) и центры зазоров. Центр БЕЗ рамки →
# стрелки подходят к краям его route-таблиц (472..800).
L_RIGHT, C_LEFT, C_RIGHT, R_LEFT = 400, 472, 800, 864
GAP_L = (L_RIGHT + C_LEFT) // 2     # центр левого зазора
GAP_R = (C_RIGHT + R_LEFT) // 2     # центр правого зазора

ICON_W, ICON_H = 78, 46

flow = {
    "header": "Enterprise Router + Cloud Firewall",
    "grid_snap": True,
    "auto_align": False,      # multi-zone: выравнивание задано конструкцией сетки
    "groups": [], "tables": [], "blocks": [], "labels": [], "arrows": [],
    "decor": {"enabled": True, "x_start": 1092, "y_start": 632,
              "count": 4, "size": 24, "gap": 9, "w_pt": 1.2},
}


# ---- оценка высоты таблицы (зеркало формулы flow_renderer) ------------------
def _cell_lines(text, width_px, fpt):
    usable = max(1, width_px - 16)
    char_w = 0.56 * (fpt * 4.0 / 3.0)
    cpl = max(1, int(usable / char_w))
    return max(1, math.ceil(len(str(text)) / cpl))


def _row_h(cells, widths, fpt, min_h):
    line_h = 1.18 * (fpt * 4.0 / 3.0)
    ml = 1
    for i, wd in enumerate(widths):
        v = cells[i] if i < len(cells) else ""
        ml = max(ml, _cell_lines(v, wd, fpt))
    return max(min_h, int(ml * line_h + 5))


def _col_widths(w, n, col_fracs):
    if col_fracs and len(col_fracs) == n:
        s = float(sum(col_fracs)) or 1.0
        ws = [int(w * f / s) for f in col_fracs]
    else:
        ws = [w // n] * n
    ws[-1] = w - sum(ws[:-1])
    return ws


def table_h(w, title, headers, rows, col_fracs=None):
    n = max(len(headers) if headers else 0,
            max((len(r) for r in rows), default=1))
    ws = _col_widths(w, n, col_fracs)
    h = 0
    if title:
        h += _row_h([title], [w], TF, TMIN)
    if headers:
        h += _row_h(headers, ws, HF, HMIN)
    for r in rows:
        h += _row_h(r, ws, DF, RMIN)
    return h


# ---- строители: таблица по ВЕРХУ ряда; иконка центрируется по своей таблице --
def add_table(tid, x, top, w, title, headers, rows, col_fracs=None):
    h = table_h(w, title, headers, rows, col_fracs)
    flow["tables"].append({
        "id": tid, "x": x, "y": top, "w": w, "title": title,
        "title_fill": BAND_GRAY, "headers": headers, "rows": rows,
        "col_fracs": col_fracs, "font_pt": DF, "title_font_pt": TF,
        "header_font_pt": HF, "title_h": TMIN, "header_h": HMIN, "row_h": RMIN,
    })
    return top + h // 2      # вертикальный центр таблицы (для иконки)


def add_icon(x, center, lines):
    flow["blocks"].append({
        "x": x, "y": center - ICON_H // 2, "w": ICON_W, "h": ICON_H, "lines": lines,
        "font_sizes": [10] + [9] * (len(lines) - 1),
        "bolds": [True] + [False] * (len(lines) - 1),
        "fill": GRAY100, "align": "left", "vanchor": "middle",
        "bottom_stroke": BS_GREEN,    # компонент-узел: Gray 100 + зелёный нижний штрих
    })


def label(x, y, w, text, size=10, bold=False, align="left", color=None):
    d = {"x": x, "y": y, "w": w, "h": 18, "text": text,
         "font_size": size, "bold": bold, "align": align}
    if color:
        d["color"] = color
    flow["labels"].append(d)


# Слайд белый = Background. Заливки ЧЕРЕДУЮТСЯ от белого: зона Gray100 (L1) →
# security group White (L2) → subnet пунктир (L3). Серый фон-стейдж НЕ нужен
# (давал серые «колонки» в зазорах между картами).

# ============================================================ LEFT (VPC1/VPC2)
PANEL_H = 200   # единая высота левых панелей (VPC1/VPC2)


def vpc_left(panel_y, n, cidr, sub_cidr, ecs, top, rt_title, rt_rows):
    px, pw = 36, 372
    # L1: VPC = Gray 100 + нижний штрих graphite-50 (на белом слайде).
    flow["groups"].append({"style": "panel", "x": px, "y": panel_y, "w": pw,
                           "h": PANEL_H, "fill": GRAY100, "bottom_stroke": BS_GRAY})
    # Заголовок: верх к оси HEAD_Y (panel_y+4), левый край к панели (px).
    label(px, panel_y + 4, pw - 8, "VPC %d–Region A\n%s" % (n, cidr), 11, True)
    sx, sy, sw = px + 12, panel_y + 44, pw - 24
    sh = PANEL_H - 44 - 6
    # L2: security group = White (чередование заливки на сером VPC).
    flow["groups"].append({"style": "panel", "x": sx, "y": sy, "w": sw, "h": sh,
                           "fill": WHITE})
    label(sx + 10, sy + 6, sw - 20, "Security group (general-purpose web server)",
          10, True)
    # L3: subnet = пунктир graphite-50 (единственный пунктирный слой).
    bx, by, bw, bh = sx + 8, sy + 28, sw - 16, sh - 36
    flow["groups"].append({"x": bx, "y": by, "w": bw, "h": bh})
    label(bx + 8, by + 4, bw - 16, "Subnet 1–%s" % sub_cidr, 9)
    # иконка слева (bx+8..+ICON_W), таблица правее с зазором 8 — без наезда
    cy = add_table("rt_vpc%d" % n, bx + 8 + ICON_W + 8, top, bw - 16 - ICON_W - 8,
                   rt_title, ["Destination", "Next Hop"], rt_rows,
                   col_fracs=[0.42, 0.58])
    add_icon(bx + 8, cy, [ecs])


vpc_left(PANEL_TOP, 1, "10.1.0.0/16", "10.1.0.0/24", "ECS 1", T1,
         "VPC 1 route table",
         [["10.1.0.0/24", "Local"], ["0.0.0.0/0", "Enterprise router"]])
vpc_left(394, 2, "10.2.0.0/16", "10.2.0.0/24", "ECS 2", T2,
         "VPC 2 route table",
         [["10.2.0.0/24", "Local"], ["0.0.0.0/0", "Enterprise router"]])

# ============================================================ CENTER (router)
# Enterprise router — НЕ VPC: без серой рамки (как в оригинале). Только заголовок
# + две route-таблицы (у них свои границы). Серые фреймы = только VPC-зоны.
cx, cw = 470, 330
label(cx, HEAD_Y, cw, "Enterprise router–Region A", 11, True)
add_table("rt1", cx, T1, cw, "Route table 1",
          ["Destination", "Next Hop", "Route Type"],
          [["0.0.0.0/0", "VPC 3 attachment", "Static route"],
           ["x.x.x.x/xx", "xxx attachment", "xxx route"],
           ["x.x.x.x/xx", "xxx attachment", "xxx route"]],
          col_fracs=[0.30, 0.40, 0.30])
add_table("rt2", cx, T2, cw, "Route table 2",
          ["Destination", "Next Hop", "Route Type"],
          [["10.1.0.0/16", "VPC 1 attachment", "Propagated route"],
           ["10.2.0.0/16", "VPC 2 attachment", "Propagated route"],
           ["x.x.x.x/xx", "xxx attachment", "xxx route"]],
          col_fracs=[0.30, 0.40, 0.30])

# ============================================================ RIGHT (VPC3)
rx, rw = 860, 384
# L1: VPC3 = Gray 100 + нижний штрих.
flow["groups"].append({"style": "panel", "x": rx, "y": PANEL_TOP, "w": rw,
                       "h": PANEL_BOT - PANEL_TOP, "fill": GRAY100,
                       "bottom_stroke": BS_GRAY})
label(rx + 4, HEAD_Y, rw - 8, "VPC 3–Region A\n192.168.0.0/16", 11, True)
# L2: security group = White (вокруг обоих subnet'ов).
sgx, sgy, sgw = rx + 8, 190, rw - 16
sgh = PANEL_BOT - sgy - 6
flow["groups"].append({"style": "panel", "x": sgx, "y": sgy, "w": sgw, "h": sgh,
                       "fill": WHITE})
label(sgx + 10, sgy + 6, sgw - 20, "Security group (general-purpose web server)",
      10, True)
# L3: subnet 1 = пунктир (ECS 3 + Default route table).
flow["groups"].append({"x": sgx + 8, "y": T1 - 28, "w": sgw - 16, "h": 150})
label(sgx + 16, T1 - 24, sgw - 32, "Subnet 1–192.168.0.0/24", 9)
cy1 = add_table("rt_def", rx + 108, T1, rw - 132, "Default route table",
                ["Destination", "Next Hop"],
                [["192.168.0.0/24", "Local"], ["192.168.0.0/24", "Local"],
                 ["10.1.0.0/16", "ECS 3"], ["10.2.0.0/16", "ECS 3"]],
                col_fracs=[0.44, 0.56])
add_icon(rx + 22, cy1, ["ECS 3", "NIC eth0"])
# L3: subnet 2 = пунктир (Firewall + Custom route table).
flow["groups"].append({"x": sgx + 8, "y": T2 - 28, "w": sgw - 16, "h": 130})
label(sgx + 16, T2 - 24, sgw - 32, "Subnet 2–192.168.1.0/24", 9)
cy2 = add_table("rt_cust", rx + 108, T2, rw - 132, "Custom route table",
                ["Destination", "Next Hop"],
                [["192.168.0.0/24", "Local"], ["192.168.0.0/24", "Local"],
                 ["0.0.0.0/0", "Enterprise router"]],
                col_fracs=[0.44, 0.56])
add_icon(rx + 22, cy2, ["Firewall", "NIC eth1"])

# ВНУТРЕННИЙ поток VPC3 (в оригинале есть — не терять): ECS3(eth0) ↔ Firewall(eth1).
# request вниз (solid) на инспекцию, response вверх (dashed); blocked ✕ на сегменте.
icx = rx + 22 + ICON_W // 2          # центр колонки иконок
flow["arrows"] += [
    {"x1": icx - 6, "y1": cy1 + ICON_H // 2, "x2": icx - 6, "y2": cy2 - ICON_H // 2},
    {"x1": icx + 8, "y1": cy2 - ICON_H // 2, "x2": icx + 8, "y2": cy1 + ICON_H // 2,
     "dashed": True},
]
label(icx - 12, (cy1 + cy2) // 2 - 9, 24, "✕", 12, bold=True, color="#E03B3B")

# ============================================================ Arrows + labels
# Трафик строго между краями зон (L_RIGHT↔C_LEFT, C_RIGHT↔R_LEFT) по осям ARROW1/2.
flow["arrows"] += [
    {"x1": L_RIGHT, "y1": ARROW1, "x2": C_LEFT, "y2": ARROW1},
    {"x1": C_RIGHT, "y1": ARROW1, "x2": R_LEFT, "y2": ARROW1},
    {"x1": C_LEFT, "y1": ARROW2, "x2": L_RIGHT, "y2": ARROW2, "dashed": True},
    {"x1": R_LEFT, "y1": ARROW2, "x2": C_RIGHT, "y2": ARROW2, "dashed": True},
]
# Подписи attachment — отцентрованы в зазорах, единая ширина/текст.
label(GAP_L - 34, ARROW1 - 30, 68, "VPC 1\nattachment", 9, align="center")
label(GAP_L - 34, ARROW2 - 30, 68, "VPC 2\nattachment", 9, align="center")
label(GAP_R - 34, ARROW1 - 30, 68, "VPC 3\nattachment", 9, align="center")

# Associate attachment (зелёные линии без наконечника) — referent для легенды:
# связь VPC route table ↔ route table Enterprise router. Под линией трафика.
A1, A2 = ARROW1 + 16, ARROW2 + 16     # под линиями трафика, единый отступ
flow["arrows"] += [
    {"x1": L_RIGHT, "y1": A1, "x2": C_LEFT, "y2": A1, "with_head": False, "color": GREEN},
    {"x1": C_RIGHT, "y1": A1, "x2": R_LEFT, "y2": A1, "with_head": False, "color": GREEN},
    {"x1": L_RIGHT, "y1": A2, "x2": C_LEFT, "y2": A2, "with_head": False, "color": GREEN},
    {"x1": C_RIGHT, "y1": A2, "x2": R_LEFT, "y2": A2, "with_head": False, "color": GREEN},
]

# ============================================================ Legend
ly = 624
label(36, ly, 170, "Associate attachment", 10)
label(36, ly + 24, 170, "Request traffic", 10)
label(520, ly, 170, "Response traffic (dashed)", 10)
label(520, ly + 24, 170, "Blocked traffic", 10)
label(700, ly + 23, 24, "✕", 12, bold=True, color="#E03B3B")
flow["arrows"] += [
    {"x1": 210, "y1": ly + 8, "x2": 290, "y2": ly + 8,
     "with_head": False, "color": GREEN},
    {"x1": 210, "y1": ly + 32, "x2": 290, "y2": ly + 32},
    {"x1": 700, "y1": ly + 8, "x2": 780, "y2": ly + 8, "dashed": True},
]


# ---------------------------------------------------------------------- build
def main():
    tpl = resolve_template()
    out = os.path.join(HERE, "out_hard_schema.pptx")
    p = Presentation(tpl)
    originals = list(p.slides)
    new_slide = clone_slide(p, originals[BLANK_DONOR_WHITE - 1])
    sldIdLst = p.slides._sldIdLst
    for sid in list(sldIdLst)[:len(originals)]:
        rid = sid.attrib[qn("r:id")]
        try:
            p.part.drop_rel(rid)
        except Exception:
            pass
        sldIdLst.remove(sid)
    clean_slide_to_blank(new_slide)
    fr.render_flow_diagram_slide(new_slide, flow, dark=False)
    p.save(out)
    print("Saved", out)
    _parity_check()


def _parity_check():
    """Фаза 4 анализатора полноты: счётчики по категориям + legend⇒referent."""
    g = flow["groups"]
    panels = [x for x in g if x.get("style") == "panel"]
    dashed = [x for x in g if x.get("style") != "panel"]
    arr = flow["arrows"]
    solid = [a for a in arr if not a.get("dashed") and a.get("color") is None
             and a.get("with_head", True)]
    dash = [a for a in arr if a.get("dashed")]
    assoc = [a for a in arr if a.get("color") == GREEN]
    cells = sum(1 + (len(t.get("headers") or [])) // max(1, len(t.get("headers") or [1]))
                for t in flow["tables"])
    print("--- PARITY (output) ---")
    print("F frames:   panels=%d  dashed(sg/subnet)=%d" % (len(panels), len(dashed)))
    print("N nodes:    icon-boxes=%d" % len(flow["blocks"]))
    print("T tables:   %d" % len(flow["tables"]))
    print("W labels:   %d" % len(flow["labels"]))
    print("P paths:    request(solid)=%d  response(dashed)=%d  associate(green)=%d"
          % (len(solid), len(dash), len(assoc)))
    # legend⇒referent
    ok = len(assoc) >= 1 and any("✕" in l["text"] for l in flow["labels"]) \
        and len(dash) >= 1 and len(solid) >= 1
    print("L legend⇒referent: %s" % ("OK" if ok else "BROKEN — потерян путь!"))


if __name__ == "__main__":
    main()
