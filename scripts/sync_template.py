#!/usr/bin/env python3
"""
sync_template.py — пересборка производных артефактов после обновления шаблона.

Зачем: когда дизайнеры присылают новую версию Cloud.ru_Template_2026.pptx,
простой замены .pptx НЕ достаточно. От шаблона зависят:
  • brand/template-layouts-dump.json  — координаты placeholder'ов (МЕХАНИЧЕСКИ → авто)
  • template/png/*                     — эталонные рендеры для visual diff (авто, --render)
  • brand/template-version.json        — индексы donor-слайдов (СЕМАНТИКА → ручная сверка)
  • brand/donor-slot-map.yaml          — карта слотов доноров (СЕМАНТИКА → ручная сверка)
  • brand/template-analysis.md         — каталог 102 layouts (СЕМАНТИКА → ручная сверка)

Этот скрипт автоматизирует механическую часть и печатает health-report:
что обновилось само и что требует ручной/LLM-проверки.

Использование:
  python3 scripts/sync_template.py            # dry-run: только отчёт, без записи
  python3 scripts/sync_template.py --apply     # + пересобрать layouts-dump (с бэкапом)
  python3 scripts/sync_template.py --apply --render   # + перерендерить template/png/
"""
import os
import sys
import json
import shutil
import subprocess

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)

from template_path import resolve_template  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

DUMP_PATH = os.path.join(ROOT, "brand", "template-layouts-dump.json")
VER_PATH = os.path.join(ROOT, "brand", "template-version.json")
PNG_DIR = os.path.join(ROOT, "..", "template", "png")
EMU_PER_PX = 9525


def _px(emu):
    if emu is None:
        return None
    return round(int(emu) / EMU_PER_PX)


def _ph_type(ph):
    try:
        t = ph.placeholder_format.type
        if t is None:
            return "NONE"
        return f"{t.name} ({int(t)})"
    except Exception:
        return "UNKNOWN"


def build_dump(prs, old_categories):
    """Собирает layouts-dump в текущей схеме. category берётся из старого по имени."""
    layouts = []
    idx = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            phs = []
            for ph in layout.placeholders:
                phs.append({
                    "idx": ph.placeholder_format.idx,
                    "type": _ph_type(ph),
                    "name": ph.name,
                    "left_px": _px(ph.left),
                    "top_px": _px(ph.top),
                    "width_px": _px(ph.width),
                    "height_px": _px(ph.height),
                })
            layouts.append({
                "index": idx,
                "name": layout.name,
                "category": old_categories.get(layout.name, "UNKNOWN"),
                "placeholders_count": len(phs),
                "placeholders": phs,
            })
            idx += 1
    return {
        "slide_size": {
            "width_emu": prs.slide_width,
            "height_emu": prs.slide_height,
            "width_px_at96": _px(prs.slide_width),
            "height_px_at96": _px(prs.slide_height),
        },
        "layouts_count": len(layouts),
        "layouts": layouts,
    }


def main():
    apply = "--apply" in sys.argv
    render = "--render" in sys.argv

    tpl = resolve_template()
    print(f"Шаблон: {tpl}")
    prs = Presentation(tpl)
    n_slides = len(prs.slides)

    # старый дамп для сравнения + сохранения категорий
    old = json.load(open(DUMP_PATH, encoding="utf-8")) if os.path.isfile(DUMP_PATH) else {"layouts": []}
    old_categories = {l["name"]: l.get("category", "UNKNOWN") for l in old.get("layouts", [])}
    old_names = [l["name"] for l in old.get("layouts", [])]

    new_dump = build_dump(prs, old_categories)
    new_names = [l["name"] for l in new_dump["layouts"]]

    added = [n for n in new_names if n not in old_names]
    removed = [n for n in old_names if n not in new_names]
    unknown_cat = [l["name"] for l in new_dump["layouts"] if l["category"] == "UNKNOWN"]

    print("\n=== LAYOUTS ===")
    print(f"  было: {len(old_names)}   стало: {len(new_names)}")
    if added:
        print(f"  НОВЫЕ ({len(added)}): " + ", ".join(added))
    if removed:
        print(f"  УДАЛЕНЫ ({len(removed)}): " + ", ".join(removed))
    if not added and not removed:
        print("  имена layouts не изменились ✓")

    # сверка donor-индексов из template-version.json
    print("\n=== DONOR-СЛАЙДЫ (template-version.json) ===")
    ver = json.load(open(VER_PATH, encoding="utf-8")) if os.path.isfile(VER_PATH) else {}
    donor_groups = {**ver.get("blank_donors", {}), **ver.get("guide_donors", {})}
    donor_groups = {k: v for k, v in donor_groups.items() if isinstance(v, int)}
    print(f"  всего слайдов в шаблоне: {n_slides}")
    for name, num in donor_groups.items():
        if 1 <= num <= n_slides:
            sl = list(prs.slides)[num - 1]
            lname = sl.slide_layout.name
            print(f"  {name:14s} = слайд {num:>3} → layout '{lname}'")
        else:
            print(f"  {name:14s} = слайд {num:>3} → ВНЕ ДИАПАЗОНА (1..{n_slides}) ⚠")

    # запись
    print("\n=== ДЕЙСТВИЯ ===")
    if apply:
        if os.path.isfile(DUMP_PATH):
            shutil.copy2(DUMP_PATH, DUMP_PATH + ".bak")
            print(f"  бэкап старого дампа → {os.path.basename(DUMP_PATH)}.bak")
        with open(DUMP_PATH, "w", encoding="utf-8") as f:
            json.dump(new_dump, f, ensure_ascii=False, indent=2)
        print(f"  пересобран {os.path.relpath(DUMP_PATH, ROOT)} ✓")
    else:
        print("  (dry-run) дамп НЕ записан — добавь --apply чтобы применить")

    if render:
        out = os.path.abspath(PNG_DIR)
        tmp = os.path.join(out, "_tmp_render")
        print(f"  рендер PNG → {out} (может занять пару минут)...")
        if os.path.isdir(tmp):
            shutil.rmtree(tmp)
        os.makedirs(tmp, exist_ok=True)
        r = subprocess.run(
            [sys.executable, os.path.join(HERE, "render_slides.py"), tpl, tmp],
            capture_output=True, text=True)
        if r.returncode != 0:
            print("  рендер: FAIL ⚠\n" + r.stderr[-500:])
        else:
            # render_slides пишет slide-NN.png → переименовываем в Слайд{N}.png
            # (имя, по которому агент-верификатор ищет эталоны)
            import glob
            import re
            for old in glob.glob(os.path.join(out, "Слайд*.png")):
                os.remove(old)
            n = 0
            for f in sorted(glob.glob(os.path.join(tmp, "slide-*.png"))):
                m = re.search(r"slide-(\d+)\.png$", f)
                if m:
                    num = int(m.group(1))
                    shutil.move(f, os.path.join(out, f"Слайд{num}.png"))
                    n += 1
            shutil.rmtree(tmp, ignore_errors=True)
            print(f"  рендер: OK ✓ — {n} эталонов как Слайд{{N}}.png (старые удалены)")
    elif apply:
        print("  PNG НЕ перерендерены — добавь --render (нужен для visual diff)")

    # что требует ручной проверки
    print("\n=== ТРЕБУЕТ РУЧНОЙ / LLM ПРОВЕРКИ ===")
    review = []
    if added or removed:
        review.append("• Layouts изменились → сверь donor-индексы в template-version.json "
                       "(blank white/dark, kpi, table) и пересмотри donor-slot-map.yaml")
    if unknown_cat:
        review.append(f"• Новые layouts без категории ({len(unknown_cat)}): "
                      + ", ".join(unknown_cat[:8])
                      + " → проставь category в дампе / template-analysis.md")
    review.append("• Глазами проверь, что donor-слайды выше всё ещё blank/guide того типа "
                  "(белый/тёмный фон, KPI-гайд, таблица) — индексы могли «съехать»")
    review.append("• Если менялись цвета/шрифты бренда → обнови brand/palette.json и brand-rules.md")
    for r in review:
        print("  " + r)

    print("\nГотово." + ("" if apply else "  Это был dry-run."))


if __name__ == "__main__":
    main()
