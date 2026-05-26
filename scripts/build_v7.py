#!/usr/bin/env python3
"""
build_v7.py — build_v6 + поддержка ДУБЛЕЙ donor.

Ключевое улучшение vs v6:
- Если donor 13 нужен 3 раза → клонируется 3 раза (не один и не игнорируется)
- Через XML deepcopy slide part-а в presentation
- Сохраняет правильный порядок слайдов

Plan:
{
  "slides": [
    {"clone_from_slide": 13, "slots": {...}},
    {"clone_from_slide": 12, "slots": {...}},
    {"clone_from_slide": 13, "slots": {...}}  ← ДУБЛЬ donor 13!
  ]
}

Usage:
    python3 build_v7.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]
"""
import sys
import json
import os
import copy
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from build_v5 import (
    load_donor_map, get_text_frame_by_shape_idx, replace_text_with_style,
    clear_text_frame
)

EMU_PER_PX = 9525


def clone_slide(prs, src_slide):
    """Глубоко копирует slide-part и регистрирует его в presentation.
    Возвращает новый Slide (последний в prs.slides)."""
    from pptx.opc.constants import CONTENT_TYPE as CT
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.packuri import PackURI
    from pptx.parts.slide import SlidePart

    src_part = src_slide.part
    src_xml = src_part.blob

    # Подбираем уникальное имя slideN.xml
    package = prs.part.package
    existing_partnames = {str(p.partname) for p in package.iter_parts()}
    next_idx = 1
    while f"/ppt/slides/slide{next_idx}.xml" in existing_partnames:
        next_idx += 1
    new_partname = PackURI(f"/ppt/slides/slide{next_idx}.xml")

    # Создаём новый part — используем тот же content_type как у src
    new_part = SlidePart.load(
        partname=new_partname,
        content_type=src_part.content_type,
        blob=src_xml,
        package=package,
    )

    # Копируем relationships от source slide в new slide
    for rel in src_part.rels.values():
        if rel.is_external:
            new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_part.relate_to(rel.target_part, rel.reltype)

    # Регистрируем slide в presentation через relationship
    rId = prs.part.relate_to(new_part, RT.SLIDE)

    # Добавляем sldId в sldIdLst
    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(el.attrib["id"]) for el in sldIdLst if "id" in el.attrib]
    next_id = max(existing_ids) + 1 if existing_ids else 256
    new_sldId = etree.SubElement(sldIdLst, qn("p:sldId"))
    new_sldId.set("id", str(next_id))
    new_sldId.set(qn("r:id"), rId)

    return prs.slides[-1]


def build(plan_path, template_path, output_path, donor_map_path):
    plan = json.load(open(plan_path, encoding="utf-8"))
    p = Presentation(template_path)
    donors = load_donor_map(donor_map_path)

    # === STEP 1: Собираем нужные donor_nums и клонируем все слайды plan-а ===
    # Сохраняем reference на оригинальные donor slides ДО любых модификаций
    original_slides = list(p.slides)
    donor_originals = {}  # {donor_num: original_slide}
    for ps in plan["slides"]:
        n = ps.get("clone_from_slide")
        if n and n not in donor_originals:
            if 1 <= n <= len(original_slides):
                donor_originals[n] = original_slides[n - 1]
            else:
                print(f"WARN: donor {n} вне диапазона (1..{len(original_slides)})", file=sys.stderr)

    # Клонируем КАЖДЫЙ слайд из plan (включая дубли) → новые slides в конце
    cloned_for_plan = []
    for ps in plan["slides"]:
        n = ps.get("clone_from_slide")
        if not n or n not in donor_originals:
            cloned_for_plan.append(None)
            continue
        new_slide = clone_slide(p, donor_originals[n])
        cloned_for_plan.append(new_slide)

    # === STEP 2: Удаляем все ОРИГИНАЛЬНЫЕ слайды (101+ template слайдов), оставляем только клоны ===
    sldIdLst = p.slides._sldIdLst
    n_originals = len(original_slides)
    # Первые n_originals элементов — это оригиналы. Удаляем их.
    all_sldIds = list(sldIdLst)
    for sldId in all_sldIds[:n_originals]:
        rId = sldId.attrib[qn('r:id')]
        try:
            p.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)

    # === STEP 3: Заполняем text + pictures для каждого clone ===
    pictures_inserted = 0
    for plan_slide, actual in zip(plan["slides"], cloned_for_plan):
        if actual is None:
            continue
        src_num = plan_slide["clone_from_slide"]
        donor_def = donors.get(src_num)

        # TEXT slots
        if donor_def is not None:
            slot_defs = donor_def.get("slots", {})
            slots_filled = plan_slide.get("slots", {})
            styles_override = plan_slide.get("slot_styles_override", {})

            for slot_name, new_text in slots_filled.items():
                if slot_name not in slot_defs:
                    print(f"WARN: slot '{slot_name}' undefined for donor {src_num}", file=sys.stderr)
                    continue
                shape_idx = slot_defs[slot_name]["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is None:
                    continue
                override = styles_override.get(slot_name)
                replace_text_with_style(tf, new_text, override)

            # Очистить незаполненные обязательные слоты
            for slot_name, slot_def in slot_defs.items():
                if slot_name in slots_filled:
                    continue
                if slot_def.get("optional"):
                    continue
                shape_idx = slot_def["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is not None:
                    clear_text_frame(tf)

        # PICTURES (вставляются ПОВЕРХ donor shapes)
        for pic in plan_slide.get("pictures", []):
            file_path = pic.get("file")
            if not file_path or not os.path.exists(file_path):
                print(f"WARN: image not found: {file_path}", file=sys.stderr)
                continue
            try:
                actual.shapes.add_picture(
                    file_path,
                    Emu(pic.get("left_px", 0) * EMU_PER_PX),
                    Emu(pic.get("top_px", 0) * EMU_PER_PX),
                    Emu(pic.get("width_px", 100) * EMU_PER_PX),
                    Emu(pic.get("height_px", 100) * EMU_PER_PX),
                )
                pictures_inserted += 1
            except Exception as e:
                print(f"WARN: insert_picture failed: {e}", file=sys.stderr)

        # TABLES (новое в v7)
        table_data = plan_slide.get("table_data")
        if table_data:
            try:
                rows = len(table_data)
                cols = max(len(r) for r in table_data) if rows else 1
                # Размещаем под title (top=120) с margin 35 по бокам
                left = Emu(35 * EMU_PER_PX)
                top = Emu(120 * EMU_PER_PX)
                width = Emu(1210 * EMU_PER_PX)
                height = Emu(min(550, rows * 50) * EMU_PER_PX)
                tbl_shape = actual.shapes.add_table(rows, cols, left, top, width, height)
                tbl = tbl_shape.table
                for r_idx, row_data in enumerate(table_data):
                    for c_idx, cell_text in enumerate(row_data):
                        if c_idx >= cols:
                            continue
                        tbl.cell(r_idx, c_idx).text = str(cell_text)
            except Exception as e:
                print(f"WARN: add_table failed: {e}", file=sys.stderr)

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides, {pictures_inserted} pictures inserted",
          file=sys.stderr)


def main():
    if len(sys.argv) < 4:
        print("Usage: build_v7.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]",
              file=sys.stderr)
        sys.exit(1)
    plan_p = sys.argv[1]
    tpl_p = sys.argv[2]
    out_p = sys.argv[3]
    donor_p = sys.argv[4] if len(sys.argv) > 4 else "pptx-skill/brand/donor-slot-map.yaml"
    build(plan_p, tpl_p, out_p, donor_p)


if __name__ == "__main__":
    main()
