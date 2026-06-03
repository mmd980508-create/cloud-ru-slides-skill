#!/usr/bin/env python3
"""
brand_guardian.py — валидация .pptx против Cloud.ru 2.0 brand rules.

Проверяет:
- Цвета text-run в палитре §4 brand-rules.md
- Шрифты SB Sans Display / Verdana (fallback) / theme fonts (+mj-lt, +mn-lt)
- Размеры ≥ 10pt
- Overflow heuristic (chars vs frame width)
- Эмодзи отсутствуют
- Сочетание базовой и расширенной палитр (warning)

Usage:
    python3 brand_guardian.py <input.pptx> [output_report.json]

Returns exit code:
    0 — PASS (нет нарушений)
    1 — WARN (есть warnings, но критических нарушений нет)
    2 — FAIL (есть критические нарушения)
"""
import os
import sys
import json
import unicodedata
import math
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# Палитра — единый источник истины: brand/palette.json
# Hardcoded fallback ниже используется ТОЛЬКО если файл недоступен.
# ---------------------------------------------------------------------------

_PALETTE_FALLBACK = {
    "base": {
        "Green": "#26D07C", "Yellow": "#CFF500", "Purple": "#A068FF",
        "Blue": "#C0E0FC", "Black": "#222222", "White": "#FFFFFF", "Gray": "#F2F2F2",
    },
    "base_alts": {
        "Black-pure-alt": "#000000",
        "Green-alt": "#00D97B", "Black-soft-alt-1": "#181818",
        "Black-soft-alt-2": "#2B2B2B", "Green-dark-alt": "#1AB066",
    },
    "extended": {
        "Aquamarine": "#18F4CF", "Aquamarine2": "#C9F2EA",
        "Ultramarine": "#0063FF", "Ultramarine5": "#C9D9F2",
        "Magenta": "#FF00FF", "Magenta3": "#C067C0",
        "Carrot": "#FF4517", "Carrot3": "#DD7D64",
        "Coral": "#FF0642", "Coral3": "#E25B7C",
    },
    "text_neutral": [
        "#222222", "#000000", "#181818", "#2B2B2B",
        "#F2F2F2", "#C8C8C8", "#888888", "#BDBDBD", "#666666",
    ],
    "text_dark_extra": ["#FFFFFF"],  # 2026-05-29 (Problem #2): зелёный текст запрещён везде
    "green_for_arrows": ["#26D07C", "#00D97B", "#1AB066", "#1ABF6F"],
}


def _load_palette():
    """Грузит brand/palette.json относительно корня pptx-skill/, иначе fallback."""
    here = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(here, "..", "brand", "palette.json"),
        os.path.join(os.getcwd(), "brand", "palette.json"),
        os.path.join(os.getcwd(), "pptx-skill", "brand", "palette.json"),
    ]
    for path in candidates:
        try:
            if os.path.isfile(path):
                with open(path, encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            continue
    return _PALETTE_FALLBACK


_PAL = _load_palette()

# Backwards-compatible: PALETTE_BASE/EXTENDED как и были — {hex: name}.
PALETTE_BASE = {hx: name for name, hx in {**_PAL.get("base", {}),
                                           **_PAL.get("base_alts", {})}.items()}
PALETTE_EXTENDED = {hx: name for name, hx in _PAL.get("extended", {}).items()}

ALLOWED_FONTS_PREFIXES = (
    "SB Sans",  # SB Sans Display / SB Sans Text / SB Sans Interface / SB Sans Display Semibold / Regular
    "Verdana",
    "Pingfang", "PingFang",
    "SimSun",
    "+mj-",  # theme major
    "+mn-",  # theme minor
)


def is_font_allowed(font_name):
    if not font_name:
        return True  # inherit from theme
    for prefix in ALLOWED_FONTS_PREFIXES:
        if font_name.startswith(prefix):
            return True
    return False

MIN_FONT_PT = 10           # жёсткий пол: < 10pt → FAIL (нечитаемо даже в перегрузке)
COMFORTABLE_MIN_PT = 12    # комфортный минимум (Problem #5): < 12pt → WARN (только перегрузка)
STANDARD_FONT_PT = 16      # стандартный размер контента (Problem #5, 2026-05-29)
MAX_GREEN_AREA_PCT = 30  # зелёный должен быть акцентом, ≤ 30% площади
EMU_PER_PT = 12700


# Текст должен быть нейтральным: графит или серые. Цветной текст — WARN.
# Источник: brand/palette.json (text_neutral, text_dark_extra).
TEXT_NEUTRAL_HEXES = tuple(_PAL.get("text_neutral",
                                     _PALETTE_FALLBACK["text_neutral"]))
# На тёмных слайдах дополнительно разрешён только белый текст.
# Canonical 2026-05-29 (Problem #2): зелёный как ЦВЕТ ТЕКСТА запрещён везде
# (вкл. тёмные слайды и крупные KPI/divider-цифры). Зелёный — только как
# цветной ЭЛЕМЕНТ (плашка/divider/заливка), не буквы/цифры.
TEXT_DARK_BG_EXTRA_HEXES = tuple(_PAL.get("text_dark_extra",
                                           _PALETTE_FALLBACK["text_dark_extra"]))
TEXT_COLOR_TOLERANCE = 20

DARK_LAYOUT_KEYWORDS = ("тёмн", "темн", "dark", "graphite", "ночь", "black")


def _is_dark_slide(slide):
    """Эвристика: тёмный фон — по layout name (содержит ключевое слово) или по
    background fill с luminance < 128. Используется для разрешения белого/зелёного текста."""
    try:
        layout = (slide.slide_layout.name or "").lower()
        for kw in DARK_LAYOUT_KEYWORDS:
            if kw in layout:
                return True
    except Exception:
        pass
    try:
        bg = slide.background.fill
        if bg.type is not None:
            hx = hex_color(bg.fore_color.rgb)
            if hx and len(hx) == 7:
                r = int(hx[1:3], 16)
                g = int(hx[3:5], 16)
                b = int(hx[5:7], 16)
                lum = 0.299 * r + 0.587 * g + 0.114 * b
                if lum < 128:
                    return True
    except Exception:
        pass
    return False


def _text_color_violates(hx, is_dark):
    """True если цвет текста НЕ из whitelist — это значит цветной текст (WARN)."""
    if not hx:
        return False
    allowed = list(TEXT_NEUTRAL_HEXES)
    if is_dark:
        allowed += list(TEXT_DARK_BG_EXTRA_HEXES)
    for ok_hex in allowed:
        try:
            if color_distance(hx, ok_hex) <= TEXT_COLOR_TOLERANCE:
                return False
        except Exception:
            continue
    return True


def has_emoji(text):
    """Эвристика: emoji в Unicode block ≥ 0x1F300.

    Whitelist для типографических символов категории 'So' но не emoji:
    - № (U+2116 Numero sign) — стандартный русский символ
    - © ® ™ — copyright/trademark
    - ° — degree
    - ✓ ✗ ★ ☆ — бытовые символы
    """
    if not text:
        return False
    SYMBOL_WHITELIST = {
        "№", "©", "®", "™", "°", "§", "¶",
        "✓", "✗", "✔", "✘", "★", "☆",
    }
    for c in text:
        if c in SYMBOL_WHITELIST:
            continue
        # Real emoji block (Misc Symbols and Pictographs+)
        if ord(c) >= 0x1F300:
            return True
        # Emoticons block specifically
        if 0x2600 <= ord(c) <= 0x27BF:
            return True
    return False


def hex_color(rgb):
    """python-pptx RGBColor → '#RRGGBB' uppercase."""
    if rgb is None:
        return None
    s = str(rgb).upper()
    if len(s) == 6:
        return "#" + s
    return s


def color_distance(hex1, hex2):
    """Euclidean RGB distance, 0-441."""
    def to_rgb(h):
        h = h.lstrip("#")
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r1, g1, b1 = to_rgb(hex1)
    r2, g2, b2 = to_rgb(hex2)
    return ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5


def closest_palette_color(hex_in, max_distance=35):
    """Find nearest palette color. Returns (name, distance) or (None, None) if too far."""
    best = (None, None, float("inf"))
    for hx, name in {**PALETTE_BASE, **PALETTE_EXTENDED}.items():
        d = color_distance(hex_in, hx)
        if d < best[2]:
            best = (name, hx, d)
    if best[2] <= max_distance:
        return best[0], best[2]
    return None, None


def check_text_frame(tf, slide_width_emu=12192000, slide_height_emu=6858000, is_dark=False):
    """Проверить text_frame на нарушения. Возвращает {violations, warnings, base_colors, ext_colors}.
    is_dark=True разрешает белый и фирменный зелёный для текста (тёмные слайды)."""
    issues = {"violations": [], "warnings": [], "base_colors": set(), "ext_colors": set()}
    if tf is None:
        return issues

    # Размеры frame
    frame = tf._txBody.getparent()  # sp
    spPr = None
    for child in frame:
        if child.tag.endswith("}spPr"):
            spPr = child
            break

    text = tf.text
    if has_emoji(text):
        issues["warnings"].append({
            "type": "emoji",
            "text_preview": text[:50],
            "msg": f"Эмодзи в тексте: {text[:30]}"
        })

    # Run-level checks
    for para_idx, para in enumerate(tf.paragraphs):
        for run in para.runs:
            font = run.font
            text_run = run.text or ""

            # 1. Шрифт
            font_name = font.name
            if not is_font_allowed(font_name):
                issues["violations"].append({
                    "type": "font",
                    "text_preview": text_run[:30],
                    "msg": f"Шрифт '{font_name}' не в allowed (SB Sans*, Verdana, Pingfang, theme-fonts)"
                })

            # 2. Размер (Problem #5: стандарт 16pt, комфортный минимум 12pt).
            if font.size:
                pt = font.size.pt
                if pt < MIN_FONT_PT:
                    issues["violations"].append({
                        "type": "size_too_small",
                        "text_preview": text_run[:30],
                        "msg": f"Размер {pt}pt < жёсткого минимума {MIN_FONT_PT}pt — нечитаемо"
                    })
                elif pt < COMFORTABLE_MIN_PT:
                    issues["warnings"].append({
                        "type": "size_below_comfortable",
                        "text_preview": text_run[:30],
                        "msg": (f"Размер {pt}pt < {COMFORTABLE_MIN_PT}pt. Стандарт {STANDARD_FONT_PT}pt; "
                                f"меньше {COMFORTABLE_MIN_PT}pt — только при сильной перегрузке слайда.")
                    })

            # 2b. Bold-флаг запрещён (Problem #3 2026-05-29): эмфаза только через
            # начертание «SB Sans Display Semibold», не через bold (b=1).
            if font.bold:
                issues["warnings"].append({
                    "type": "bold_flag",
                    "text_preview": text_run[:30],
                    "msg": "Bold-флаг (b=1). По бренду эмфаза только через начертание «SB Sans Display Semibold», не bold.",
                })

            # 2c. Italic / underline запрещены (brand-rules §2).
            if font.italic:
                issues["warnings"].append({
                    "type": "italic",
                    "text_preview": text_run[:30],
                    "msg": "Italic (i=1) — по бренду без курсива.",
                })
            if font.underline:
                issues["warnings"].append({
                    "type": "underline",
                    "text_preview": text_run[:30],
                    "msg": "Underline — по бренду без подчёркивания.",
                })

            # 3. Цвет (с tolerance для близких к палитре)
            try:
                if font.color and font.color.type is not None:
                    hx = hex_color(font.color.rgb)
                    if hx:
                        if hx in PALETTE_BASE:
                            issues["base_colors"].add(hx)
                        elif hx in PALETTE_EXTENDED:
                            issues["ext_colors"].add(hx)
                        else:
                            close_name, close_d = closest_palette_color(hx, max_distance=50)
                            if close_name:
                                # Близкий — WARN, не FAIL
                                issues["warnings"].append({
                                    "type": "color_near_palette",
                                    "text_preview": text_run[:30],
                                    "msg": f"Цвет {hx} ≈ {close_name} (distance {close_d:.0f})"
                                })
                            else:
                                issues["violations"].append({
                                    "type": "color_off_palette",
                                    "text_preview": text_run[:30],
                                    "msg": f"Цвет {hx} не в палитре (base/extended)"
                                })

                        # Правило B: текст должен быть нейтральным (графит/серый;
                        # +белый только на тёмных). Зелёный/цветной текст — WARN везде
                        # (Problem #2 2026-05-29: зелёный — только цветной элемент, не текст).
                        if _text_color_violates(hx, is_dark):
                            bg_hint = "тёмном" if is_dark else "светлом"
                            issues["warnings"].append({
                                "type": "colored_text",
                                "text_preview": text_run[:30],
                                "msg": (
                                    f"Цветной текст {hx} на {bg_hint} фоне. По бренду текст "
                                    f"всегда графит (#222222) или серый"
                                    + (", на тёмном — белый. Зелёный/цвет — только как плашка/элемент."
                                       if is_dark else
                                       ". Зелёный/белый как текст запрещён — акцент делать цветным элементом.")
                                ),
                            })
            except Exception:
                pass

    return issues


def estimate_overflow(tf):
    """Грубая эвристика: текст в frame не влезает.
    avg_char_width ≈ 0.55 × size_pt, line_height ≈ 1.0 × size_pt."""
    if tf is None:
        return None

    text = tf.text
    if not text.strip():
        return None

    sp = tf._txBody.getparent()
    spPr = None
    for child in sp:
        if child.tag.endswith("}spPr"):
            spPr = child
            break
    if spPr is None:
        return None
    # find xfrm
    xfrm = None
    for child in spPr:
        if child.tag.endswith("}xfrm"):
            xfrm = child
            break
    if xfrm is None:
        return None
    ext = None
    for child in xfrm:
        if child.tag.endswith("}ext"):
            ext = child
            break
    if ext is None:
        return None

    width_emu = int(ext.get("cx", "0"))
    height_emu = int(ext.get("cy", "0"))
    width_pt = width_emu / EMU_PER_PT
    height_pt = height_emu / EMU_PER_PT

    # First run size for estimation
    size_pt = 14
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size:
                size_pt = run.font.size.pt
                break
        if size_pt != 14:
            break

    avg_char_w = size_pt * 0.55
    line_h = size_pt * 1.1
    chars_per_line = max(1, width_pt / avg_char_w)
    max_lines = max(1, height_pt / line_h)
    capacity = chars_per_line * max_lines

    # Считаем актуальную длину (учёт \n как новой строки)
    lines = text.split("\n")
    needed_lines = sum(max(1, math.ceil(len(line) / chars_per_line)) for line in lines)
    needed_chars = len(text)

    overflow_ratio = max(needed_chars / capacity, needed_lines / max_lines)
    # Tolerance для больших шрифтов: они часто намеренно переполняют frame (design feature)
    # Threshold: 100pt+ → tolerance 2x, 50-100pt → tolerance 1.5x, < 50pt → tolerance 1.1x
    if size_pt >= 100:
        threshold = 2.0
    elif size_pt >= 50:
        threshold = 1.5
    else:
        threshold = 1.1

    return {
        "frame_size_pt": (round(width_pt, 1), round(height_pt, 1)),
        "size_pt": size_pt,
        "capacity_chars": round(capacity),
        "needed_chars": needed_chars,
        "max_lines": round(max_lines, 1),
        "needed_lines": needed_lines,
        "overflow_ratio": round(overflow_ratio, 2),
        "threshold": threshold,
        "overflow": overflow_ratio > threshold,
    }


EMU_PER_PX = 9525  # at 96 DPI


# Canonical §3 (slide 66): стрелки серые 1pt, не зелёные.
# Источник: brand/palette.json (green_for_arrows).
GREEN_HEXES_FOR_ARROWS = tuple(_PAL.get("green_for_arrows",
                                          _PALETTE_FALLBACK["green_for_arrows"]))


def _shape_line_hex(shape):
    """Hex line color у фигуры или None."""
    try:
        line = shape.line
        if line is None:
            return None
        # У connector / line color может быть в line.color (solid) или fill
        if line.color and line.color.type is not None:
            return hex_color(line.color.rgb)
    except Exception:
        pass
    return None


def _shape_fill_hex(shape):
    """Hex fill color у фигуры или None (для arrow-block с заливкой)."""
    try:
        fill = shape.fill
        if fill is None or fill.type is None:
            return None
        return hex_color(fill.fore_color.rgb)
    except Exception:
        return None


def _hex_is_dark(hx):
    """True если цвет тёмный (luminance < 128) — фон, на котором белый текст ОК."""
    if not hx or len(hx) != 7:
        return False
    try:
        r = int(hx[1:3], 16)
        g = int(hx[3:5], 16)
        b = int(hx[5:7], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) < 128
    except Exception:
        return False


def _line_has_arrowhead(shape):
    """True если у line/connector есть <a:tailEnd> или <a:headEnd> с type != 'none'.
    Декоративные линии (L-уголки, разделители) — без наконечника, не считаются стрелками."""
    try:
        spPr = shape._element.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            return False
        for tag in ("tailEnd", "headEnd"):
            end = ln.find(qn(f"a:{tag}"))
            if end is not None:
                etype = end.get("type")
                # Если type явно "none" — наконечника нет; иначе (triangle/arrow/...) — есть
                if etype and etype.lower() != "none":
                    return True
        return False
    except Exception:
        return False


def _is_arrow_shape(shape):
    """Стрелка = MSO_SHAPE_TYPE.LINE/CONNECTOR с реальным arrowhead ИЛИ
    AUTO_SHAPE с ARROW в auto_shape_type.

    Простые линии без наконечника (декор, разделители) — НЕ стрелки.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return _line_has_arrowhead(shape)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            ast = shape.auto_shape_type
            if ast is not None and "ARROW" in str(ast).upper():
                return True
    except Exception:
        pass
    return False


def check_arrow_colors(slide):
    """Canonical §3: стрелки серые, не зелёные. Возвращает список violations."""
    violations = []
    for shape_idx, shape in enumerate(slide.shapes):
        if not _is_arrow_shape(shape):
            continue
        line_color = _shape_line_hex(shape)
        fill_color = _shape_fill_hex(shape)
        for color in (c for c in (line_color, fill_color) if c):
            try:
                if any(color_distance(color, g) < 30 for g in GREEN_HEXES_FOR_ARROWS):
                    violations.append({
                        "type": "green_arrow",
                        "shape_idx": shape_idx,
                        "color": color,
                        "msg": (
                            f"Стрелка зелёная ({color}) — на проекторе сольётся с фоном. "
                            f"Canonical §3: стрелки серые 1pt."
                        ),
                    })
                    break
            except Exception:
                continue
    return violations


# Форма (brand-rules §2): прямые углы, без теней/glow.
ROUNDED_PRESETS = ("roundRect", "round1Rect", "round2SameRect", "round2DiagRect",
                   "snip1Rect", "snip2SameRect", "snip2DiagRect", "snipRoundRect")
EFFECT_TAGS = ("effectLst", "effectDag", "outerShdw", "innerShdw", "prstShdw",
               "glow", "reflection", "softEdge")


def _iter_all_shapes(shapes):
    """Рекурсивно все фигуры (вкл. содержимое групп)."""
    for sh in shapes:
        yield sh
        try:
            if sh.shape_type == 6:  # GROUP
                yield from _iter_all_shapes(sh.shapes)
        except Exception:
            pass


def check_shape_geometry(slide):
    """Форма: скругления (warning — допустимо только как метафора) и эффекты-тени
    (violation — запрещены). Возвращает {'violations':[], 'warnings':[]}."""
    out = {"violations": [], "warnings": []}
    for idx, sh in enumerate(_iter_all_shapes(slide.shapes)):
        spPr = sh._element.find(qn("p:spPr"))
        if spPr is None:
            continue
        geom = spPr.find(qn("a:prstGeom"))
        if geom is not None and geom.get("prst") in ROUNDED_PRESETS:
            out["warnings"].append({
                "type": "rounded_shape", "shape_idx": idx,
                "msg": f"Скруглённая форма '{geom.get('prst')}' — по бренду прямые углы (искл. метафора).",
            })
        for tag in EFFECT_TAGS:
            for e in spPr.findall(qn(f"a:{tag}")):
                if tag in ("effectLst", "effectDag") and len(e) == 0:
                    continue  # пустой effectLst = эффектов нет
                out["violations"].append({
                    "type": "shape_effect", "shape_idx": idx,
                    "msg": f"Эффект <{tag}> на фигуре — по бренду без теней/glow.",
                })
                break
    return out


# Сопоставление типов issue → оси отчёта (Color / Typography / Shape).
AXIS_OF = {
    "color_off_palette": "Color", "color_near_palette": "Color", "colored_text": "Color",
    "mixed_palettes": "Color", "extended_without_green": "Color", "green_arrow": "Color",
    "composition_subtitle_green": "Color",
    "font": "Typography", "bold_flag": "Typography", "italic": "Typography",
    "underline": "Typography", "size_too_small": "Typography", "size_below_comfortable": "Typography",
    "rounded_shape": "Shape", "shape_effect": "Shape",
}


def axis_report(pptx_path):
    """Прогоняет validate_slide по всем слайдам и группирует issues по осям
    Color/Typography/Shape (token-diff формат MATCH/DIVERGE). Всё на базе того же
    brand_guardian — единый источник чеков. → dict для scorecard/печати."""
    p = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(p.slides, 1):
        r = validate_slide(slide, i)
        axes = {"Color": [], "Typography": [], "Shape": []}
        for it in r["violations"] + r["warnings"]:
            ax = AXIS_OF.get(it["type"])
            if ax:
                axes[ax].append(it["msg"])
        verdict = "PASS" if all(not v for v in axes.values()) else "DIVERGE"
        slides.append({"num": i, "axes": axes, "verdict": verdict, "score": r["score"]})
    summary = {ax: (sum(1 for s in slides if not s["axes"][ax]), len(slides))
               for ax in ("Color", "Typography", "Shape")}
    overall = "PASS" if all(s["verdict"] == "PASS" for s in slides) else "DIVERGE"
    return {"slides": slides, "axes": summary, "overall": overall}


def composition_checks(slide):
    """Композиционные проверки: top padding, subtitle color, KPI size consistency."""
    issues = []
    title_candidates = []  # (top_px, size_pt, text)
    subtitle_candidates = []  # (size_pt, color, text)
    kpi_numbers = []  # (size_pt, color)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip():
            continue
        if not shape.top or not shape.left:
            continue
        top_px = shape.top / EMU_PER_PX
        text = shape.text_frame.text.strip()

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.font.size:
                    continue
                pt = run.font.size.pt
                color = None
                try:
                    if run.font.color and run.font.color.type is not None:
                        color = hex_color(run.font.color.rgb)
                except Exception:
                    pass

                # Title candidate: 40-100pt только (исключаем KPI 199pt и dividers 200+pt)
                # И только если в верхней трети слайда (top < 200px)
                if 40 <= pt <= 100 and top_px < 200:
                    title_candidates.append((top_px, pt, text[:40], color))
                # Subtitle candidate: 18-22pt
                if 18 <= pt <= 22:
                    subtitle_candidates.append((pt, color, text[:30]))
                # KPI number candidate: ≥150pt
                if pt >= 150:
                    kpi_numbers.append((pt, color, text[:8]))
                break  # check only first run of paragraph
            break  # check only first paragraph of shape

    # Check 1: title pressed to top (top < 25px) — реальная проблема прижатия к верху
    # composition_title_low отключен — слишком много false positives для divider/центральных layouts
    for top_px, pt, txt, color in title_candidates:
        if top_px < 25 and pt >= 60:
            issues.append({
                "type": "composition_title_pressed",
                "msg": f"Title '{txt}' ({pt}pt) at top={top_px:.0f}px — прижат к верхней границе слайда"
            })

    # Check 2: subtitles in multicolumn — должны быть чёрные #222222 (canonical)
    if len(subtitle_candidates) >= 3:
        green_subs = [s for s in subtitle_candidates if s[1] and s[1].upper() in ("#26D07C", "#00D97B", "#1AB066")]
        if green_subs:
            issues.append({
                "type": "composition_subtitle_green",
                "msg": (
                    f"Multicolumn ({len(subtitle_candidates)} subtitles) с {len(green_subs)} зелёными подзаголовками — "
                    f"canonical violation (должны быть чёрные Bold)"
                )
            })

    # Check 3: KPI numbers consistency
    if len(kpi_numbers) >= 2:
        sizes = [k[0] for k in kpi_numbers]
        if max(sizes) - min(sizes) > 80:
            issues.append({
                "type": "composition_kpi_inconsistent",
                "msg": (
                    f"KPI цифры разных размеров: {sorted(sizes)} (разброс {max(sizes)-min(sizes):.0f}pt). "
                    f"Если variant=equal_size — выровнять все до 199pt"
                )
            })

    # Check 4: Peer font consistency (Mck Peer Harmonize)
    # Shapes на одной Y-координате (с tolerance ±10px) должны иметь одинаковый font size.
    # Не флагать title vs subtitle (разные намеренные иерархии) — только same-row группы 3+ shapes.
    Y_TOLERANCE_PX = 10
    peer_groups = {}  # quantized_y → [(size_pt, text)]
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        if not shape.top:
            continue
        # Quantize y to nearest 10px grid
        y_bucket = round((shape.top / EMU_PER_PX) / Y_TOLERANCE_PX) * Y_TOLERANCE_PX
        # Get first run size
        first_size = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    first_size = run.font.size.pt
                    break
            if first_size:
                break
        if first_size:
            peer_groups.setdefault(y_bucket, []).append((first_size, shape.text_frame.text.strip()[:30]))

    for y_bucket, members in peer_groups.items():
        if len(members) < 3:
            continue  # 1-2 элемента — не peer-row, может быть title+subtitle
        sizes = [m[0] for m in members]
        if len(set(sizes)) > 1:
            issues.append({
                "type": "composition_peer_inconsistent",
                "msg": (
                    f"Peer-row at y≈{y_bucket}px: {len(members)} shapes с font sizes {sorted(set(sizes))}. "
                    f"Mck rule: same-Y shapes → same font size. Harmonize to min={min(sizes)}pt"
                )
            })

    return issues


def validate_slide(slide, slide_num):
    """Полная валидация одного слайда."""
    result = {
        "num": slide_num,
        "layout_name": slide.slide_layout.name,
        "violations": [],
        "warnings": [],
        "stats": {
            "shapes": len(slide.shapes),
            "text_frames": 0,
            "base_colors_used": [],
            "ext_colors_used": [],
        }
    }

    # Композиционные проверки (все становятся warnings — не критичные FAIL)
    for ci in composition_checks(slide):
        result["warnings"].append(ci)

    # Зелёные стрелки — violation (canonical §3, slide 66)
    for av in check_arrow_colors(slide):
        result["violations"].append(av)

    # Форма: скругления (warning) + эффекты-тени (violation) — brand-rules §2
    geom = check_shape_geometry(slide)
    result["violations"].extend(geom["violations"])
    result["warnings"].extend(geom["warnings"])

    is_dark = _is_dark_slide(slide)
    result["stats"]["is_dark"] = is_dark

    base_colors_all = set()
    ext_colors_all = set()

    for shape_idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if not tf.text.strip():
            continue

        result["stats"]["text_frames"] += 1
        # Контекст «тёмный» = тёмный слайд ИЛИ тёмная плашка под текстом.
        # Белый текст на тёмной плашке (fill #222222/чёрный) — легитимен,
        # не помечать как colored_text (Problem #2 false-positive guard).
        shape_is_dark = is_dark or _hex_is_dark(_shape_fill_hex(shape))
        issues = check_text_frame(tf, is_dark=shape_is_dark)

        # Префиксируем shape_idx в каждый violation
        for v in issues["violations"]:
            v["shape_idx"] = shape_idx
            result["violations"].append(v)
        for w in issues["warnings"]:
            w["shape_idx"] = shape_idx
            result["warnings"].append(w)

        base_colors_all |= issues["base_colors"]
        ext_colors_all |= issues["ext_colors"]

        # Overflow
        ov = estimate_overflow(tf)
        if ov and ov["overflow"]:
            result["warnings"].append({
                "type": "overflow",
                "shape_idx": shape_idx,
                "msg": (
                    f"Возможный overflow: текст {ov['needed_chars']} chars при capacity ~{ov['capacity_chars']} "
                    f"(ratio {ov['overflow_ratio']}x). Рекомендация: смотреть memory/feedback_overflow_strategy.md"
                ),
                "details": ov,
            })

    result["stats"]["base_colors_used"] = sorted(base_colors_all)
    result["stats"]["ext_colors_used"] = sorted(ext_colors_all)

    # Mixed palettes warning
    if base_colors_all and ext_colors_all:
        result["warnings"].append({
            "type": "mixed_palettes",
            "msg": (
                f"Базовая ({sorted(base_colors_all)}) + расширенная ({sorted(ext_colors_all)}) палитры "
                f"в одном слайде. По брендбуку — не сочетать."
            )
        })

    # Правило A (canonical §2 slide 82): доп. цвета только вместе с фирменным зелёным.
    # Если на слайде есть extended-цвета и НЕТ зелёного (ни в base, ни декоре) — WARN.
    GREEN_BASE_HEXES = {"#26D07C", "#00D97B", "#1AB066"}
    if ext_colors_all and not (base_colors_all & GREEN_BASE_HEXES):
        result["warnings"].append({
            "type": "extended_without_green",
            "msg": (
                f"Расширенная палитра ({sorted(ext_colors_all)}) использована без фирменного "
                f"зелёного. Canonical §2: доп. цвета — только вместе с зелёным."
            )
        })

    # Score: 100 - 20*violations - 5*warnings
    score = 100 - 20 * len(result["violations"]) - 5 * len(result["warnings"])
    result["score"] = max(0, score)

    return result


def main():
    if len(sys.argv) < 2:
        print("Usage: brand_guardian.py <input.pptx> [report.json]", file=sys.stderr)
        sys.exit(2)

    pptx_path = sys.argv[1]
    report_path = sys.argv[2] if len(sys.argv) > 2 else None

    p = Presentation(pptx_path)
    report = {
        "input": pptx_path,
        "n_slides": len(p.slides),
        "slides": [],
        "summary": {"violations_total": 0, "warnings_total": 0, "score_avg": 0},
    }

    for i, slide in enumerate(p.slides, 1):
        r = validate_slide(slide, i)
        report["slides"].append(r)
        report["summary"]["violations_total"] += len(r["violations"])
        report["summary"]["warnings_total"] += len(r["warnings"])

    if report["slides"]:
        report["summary"]["score_avg"] = round(
            sum(s["score"] for s in report["slides"]) / len(report["slides"]), 1
        )

    # Verdict
    if report["summary"]["violations_total"] > 0:
        report["verdict"] = "FAIL"
        exit_code = 2
    elif report["summary"]["warnings_total"] > 0:
        report["verdict"] = "WARN"
        exit_code = 1
    else:
        report["verdict"] = "PASS"
        exit_code = 0

    # Print human-readable summary
    print(f"\n{'=' * 60}", file=sys.stderr)
    print(f"Brand Guardian Report — {pptx_path}", file=sys.stderr)
    print(f"{'=' * 60}", file=sys.stderr)
    print(f"Verdict: {report['verdict']}  |  Score avg: {report['summary']['score_avg']}/100", file=sys.stderr)
    print(f"Slides: {report['n_slides']}  |  Violations: {report['summary']['violations_total']}  |  Warnings: {report['summary']['warnings_total']}", file=sys.stderr)
    print(file=sys.stderr)

    for s in report["slides"]:
        if s["violations"] or s["warnings"]:
            print(f"--- Slide {s['num']} '{s['layout_name'][:30]}' (score {s['score']}) ---", file=sys.stderr)
            for v in s["violations"]:
                print(f"  ❌ FAIL [{v.get('shape_idx', '?')}] {v['type']}: {v['msg']}", file=sys.stderr)
            for w in s["warnings"]:
                print(f"  ⚠️  WARN [{w.get('shape_idx', '?')}] {w['type']}: {w['msg']}", file=sys.stderr)

    if report_path:
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(report, f, ensure_ascii=False, indent=2)
        print(f"\nFull report → {report_path}", file=sys.stderr)

    sys.exit(exit_code)


if __name__ == "__main__":
    main()
