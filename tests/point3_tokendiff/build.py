#!/usr/bin/env python3
"""
Point 3 test — визуальный scorecard token-diff (сверка с бренд-контрактом).
Гоняет token_diff по point2 before/after и рисует Cloud.ru-слайд с вердиктами
по осям (Color/Typography/Shape): БЫЛО (DIVERGE) vs СТАЛО (PASS).
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
SCRIPTS = os.path.join(HERE, "..", "..", "scripts")
sys.path.insert(0, SCRIPTS)

from design_tokens import load_tokens          # noqa: E402
import brand_guardian                           # noqa: E402
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu, Pt                    # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE           # noqa: E402
from pptx.oxml.ns import qn                      # noqa: E402

T = load_tokens()
EMU = T.EMU
LOGO = os.path.join(HERE, "..", "point1_tokens", "logo.png")


def emu(px):
    return Emu(int(round(px * EMU)))


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)
    sp.fill.solid(); sp.fill.fore_color.rgb = T.rgb(color)
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, runs, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (s, role, col, align) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); spec = T.role(role)
        r.text = s.upper() if spec.get("caps") else s
        r.font.name = T.font_face(role); r.font.size = Pt(spec["size"]); r.font.bold = False
        r.font.color.rgb = T.rgb(col or spec.get("color") or "Black")
    return tb


def check_chip(slide, x, y, ok):
    """Зелёный чип с галочкой (PASS) или серый с крестом (DIVERGE)."""
    size = 28
    rect(slide, x, y, size, size, "Green" if ok else "Gray")
    sp = slide.shapes[-1]
    tf = sp.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "✓" if ok else "—"
    r.font.name = T.family; r.font.size = Pt(14); r.font.bold = False
    r.font.color.rgb = T.rgb("Black")


def add_logo(slide):
    lg = T.place("logo")["content"]
    if os.path.isfile(LOGO):
        slide.shapes.add_picture(LOGO, emu(lg["x"]), emu(lg["y"]), emu(lg["w"]), emu(lg["h"]))


def add_copyright(slide):
    cp = T.place("copyright")
    tb = slide.shapes.add_textbox(emu(cp["x"]), emu(cp["y"]), emu(cp["w"]), emu(cp["h"]))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = cp["text"]
    r.font.name = T.family; r.font.size = Pt(cp["size"]); r.font.bold = False
    r.font.color.rgb = T.rgb(cp["color"])


def card(slide, x, w, label, report):
    cy, ch = 200, 380
    rect(slide, x, cy, w, ch, "Gray")
    pad = 28
    verdict = report["overall"]
    text(slide, x + pad, cy + pad, w - 2 * pad, 30,
         [(label, "header", "Black", PP_ALIGN.LEFT)])
    text(slide, x + pad, cy + pad + 30, w - 2 * pad, 24,
         [("ИТОГ: " + verdict, "caption", "text_gray" if verdict == "PASS" else "Black", PP_ALIGN.LEFT)])
    # оси
    summ = report["axes"]
    ry = cy + pad + 78
    for axis in ("Color", "Typography", "Shape"):
        ok_count, total = summ[axis]
        ok = ok_count == total
        check_chip(slide, x + pad, ry, ok)
        text(slide, x + pad + 40, ry + 2, w - 2 * pad - 40, 24,
             [(axis, "body", "Black", PP_ALIGN.LEFT)])
        # детали дивергенции (с первого слайда)
        issues = report["slides"][0]["axes"][axis] if report["slides"] else []
        if issues:
            det = "; ".join(issues)[:80]
            text(slide, x + pad + 40, ry + 28, w - 2 * pad - 40, 40,
                 [(det, "caption", "text_gray", PP_ALIGN.LEFT)])
        ry += 86


def build():
    before = brand_guardian.axis_report(os.path.join(HERE, "..", "point2_normalizer", "before.pptx"))
    after = brand_guardian.axis_report(os.path.join(HERE, "..", "point2_normalizer", "after.pptx"))

    prs = Presentation()
    prs.slide_width = emu(1280); prs.slide_height = emu(720)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 1280, 720, "White")
    add_logo(s); add_copyright(s)

    safe = T.safe
    text(s, 35, 38, safe["right"] - 35 - 170, 40,
         [("Token-diff | сверка с бренд-контрактом", "header", "Black", PP_ALIGN.LEFT)])
    text(s, 35, 76, safe["right"] - 35, 24,
         [("Оси: Color | Typography | Shape. Зелёная галочка = MATCH, серый прочерк = DIVERGE",
           "caption", "text_gray", PP_ALIGN.LEFT)])

    gap = 30
    w = (safe["right"] - safe["left"] - gap) / 2
    card(s, safe["left"], w, "БЫЛО | чужой скин", before)
    card(s, safe["left"] + w + gap, w, "СТАЛО | нормализовано", after)

    out = os.path.join(HERE, "point3_tokendiff.pptx")
    prs.save(out); print("saved:", out)
    return out


if __name__ == "__main__":
    build()
