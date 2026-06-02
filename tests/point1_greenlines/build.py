#!/usr/bin/env python3
"""
Визуальная проверка (user 2026-06-02): зелёные ДЕТАЛИ сохранены, убрана только
«висящая» декор-черта под заголовком; и — ГЛАВНОЕ — эффектов/теней нет НИГДЕ.
Рендерит 3 архетипа НАСТОЯЩИМИ рендерами скилла на blank-канвасе:
  1) KPI         — зелёный акцент-маркер над цифрой (привязан к объекту) — ОК.
  2) hero        — зелёная плашка + зелёные offset-рамки (композиция) — ОК.
  3) numbered    — крупные ЗЕЛЁНЫЕ номера — допустимы; + бренд-стрелки ↗ зелёные.
Везде: плоско, без эффектов (effects_util снимает их со всех фигур).
Каждый слайд с лого (1096,38) и копирайтом — как константы.
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
SCRIPTS = os.path.join(HERE, "..", "..", "scripts")
sys.path.insert(0, SCRIPTS)

from design_tokens import load_tokens  # noqa: E402
import kpi_renderer                     # noqa: E402
import flow_renderer                    # noqa: E402
from pptx import Presentation           # noqa: E402
from pptx.util import Emu, Pt           # noqa: E402
from pptx.enum.text import PP_ALIGN     # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.oxml.ns import qn             # noqa: E402

T = load_tokens()
EMU = T.EMU
LOGO_PNG = os.path.join(HERE, "..", "point1_tokens", "logo.png")


def emu(px):
    return Emu(int(round(px * EMU)))


def white_bg(slide):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(0), emu(0),
                                emu(T.meta["canvas"]["w"]), emu(T.meta["canvas"]["h"]))
    style = sp._element.find(qn("p:style"))
    if style is not None:
        sp._element.remove(style)
    sp.fill.solid(); sp.fill.fore_color.rgb = T.rgb("White")
    sp.line.fill.background(); sp.shadow.inherit = False
    # отправить назад
    sp._element.getparent().remove(sp._element)
    slide.shapes._spTree.insert(2, sp._element)


def add_logo(slide):
    lg = T.place("logo")["content"]
    if os.path.isfile(LOGO_PNG):
        slide.shapes.add_picture(LOGO_PNG, emu(lg["x"]), emu(lg["y"]),
                                 emu(lg["w"]), emu(lg["h"]))


def add_copyright(slide):
    cp = T.place("copyright")
    tb = slide.shapes.add_textbox(emu(cp["x"]), emu(cp["y"]), emu(cp["w"]), emu(cp["h"]))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = cp["text"]
    r.font.name = T.family; r.font.size = Pt(cp["size"]); r.font.bold = False
    r.font.color.rgb = T.rgb(cp["color"])


def label(slide, text):
    tb = slide.shapes.add_textbox(emu(35), emu(38), emu(963), emu(40))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text.upper()
    r.font.name = T.semibold_face; r.font.size = Pt(20); r.font.bold = False
    r.font.color.rgb = T.rgb("Black")


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(s); add_logo(s); add_copyright(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width = emu(T.meta["canvas"]["w"])
    prs.slide_height = emu(T.meta["canvas"]["h"])

    # 1) KPI — зелёный акцент-маркер над цифрой (привязан к объекту) — ОК
    s = new_slide(prs)
    kpi_renderer.render_kpi(s, {
        "title": "KPI | зелёный акцент-маркер над цифрой — ок",
        "numbers": [
            {"value": "12", "desc": "AI-инициатив", "accent": False, "pct": False},
            {"value": "84", "desc": "вовлечённость", "accent": False, "pct": True},
            {"value": "5",  "desc": "инструментов (accent)", "accent": True, "pct": False},
        ]})

    # 2) hero — зелёная заливка + зелёные offset-рамки (композиция)
    s = new_slide(prs)
    label(s, "Hero | зелёная плашка и рамки — композиция")
    flow_renderer.render_hero_statement(s, {
        "statement": "Устойчивая основа для движения вперёд",
        "support": "Зелёная плашка-заливка и offset-рамки держат композицию.",
        "content_top": 150})
    # бренд-стрелки ↗ — зелёные (дефолт; допустимы и серые) — в свободном углу
    flow_renderer.add_decor_diagonals(s, count=3, x_start=1070, y_start=600, size=46, gap=10)

    # 3) numbered_columns — крупные ЗЕЛЁНЫЕ номера (допустимы)
    s = new_slide(prs)
    label(s, "Numbered | зелёные номера — допустимы")
    flow_renderer.render_numbered_columns(s, {
        "columns": [
            {"title": "Контекст", "text": "Короткое пояснение к первому тезису колонки.", "number": "01"},
            {"title": "Проблема", "text": "Короткое пояснение ко второму тезису колонки.", "number": "02"},
            {"title": "Решение", "text": "Короткое пояснение к третьему тезису колонки.", "number": "03"},
        ]})

    # Финальный sweep эффектов по всем слайдам — как enforce_canonical в проде
    # (гарантия: теней нет НИГДЕ, на любом канвасе).
    import effects_util
    for s in prs.slides:
        effects_util.strip_effects_recursive(s.shapes)

    out = os.path.join(HERE, "point1_greenlines.pptx")
    prs.save(out)
    print("saved:", out)
    return out


if __name__ == "__main__":
    build()
