#!/usr/bin/env python3
"""
Point 6 demo (наглядно) — letter-spacing (spc, сотые пункта) на крупном кегле:
  • CAPS-заголовок: положительный трекинг 0 / 300 / 600;
  • крупная цифра: отрицательный трекинг 0 / -400 / -800 (брендбук §6: большие
    цифры с отрицательным трекингом).
Разброс намеренно большой, чтобы разница была очевидна.
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "..", "scripts"))
from design_tokens import load_tokens          # noqa: E402
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu, Pt                    # noqa: E402
from pptx.enum.text import PP_ALIGN             # noqa: E402
from pptx.enum.shapes import MSO_SHAPE          # noqa: E402
from pptx.oxml.ns import qn                     # noqa: E402
from pptx.dml.color import RGBColor             # noqa: E402

T = load_tokens()
EMU = T.EMU


def emu(px):
    return Emu(int(round(px * EMU)))


def rect(slide, x, y, w, h, hexc):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(hexc.lstrip("#"))
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def run_text(slide, x, y, w, h, text, size, spc, semibold=True, color="222222"):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.name = T.semibold_face if semibold else T.family
    r.font.size = Pt(size); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string(color)
    rPr = r._r.find(qn("a:rPr"))
    if spc:
        rPr.set("spc", str(spc))   # letter-spacing, 1/100 pt (может быть < 0)


def lbl(slide, x, y, s):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(420), emu(22))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = s
    r.font.name = T.family; r.font.size = Pt(12); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string("5C5C5C")


def build():
    prs = Presentation(); prs.slide_width = emu(1280); prs.slide_height = emu(720)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 1280, 720, "FFFFFF")

    run_text(s, 35, 30, 1200, 44, "ПУНКТ 6 · ТРЕКИНГ НАГЛЯДНО", 22, 0)

    # --- CAPS-заголовок: положительный трекинг ---
    lbl(s, 35, 92, "CAPS-заголовок (положительный трекинг):")
    for i, spc in enumerate((0, 300, 600)):
        run_text(s, 35, 120 + i * 70, 1200, 60, "ЗАГОЛОВОК СЛАЙДА", 30, spc)
        lbl(s, 720, 132 + i * 70, f"spc {spc}")

    # --- Крупная цифра: отрицательный трекинг (брендбук §6) ---
    lbl(s, 35, 360, "Крупная цифра (отрицательный трекинг, брендбук §6):")
    for i, spc in enumerate((0, -400, -800)):
        run_text(s, 35, 390 + i * 100, 600, 90, "1 280", 72, spc)
        lbl(s, 360, 420 + i * 100, f"spc {spc}")

    out = os.path.join(HERE, "point6_typography.pptx")
    prs.save(out); print("saved:", out)
    return out


if __name__ == "__main__":
    build()
