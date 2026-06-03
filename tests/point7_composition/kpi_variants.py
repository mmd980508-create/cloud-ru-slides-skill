#!/usr/bin/env python3
"""
Point 7 — более УВЕРЕННЫЕ KPI-композиции (3 варианта на выбор). Цель: убрать
«робкость» — явная иерархия, структура, уверенное использование холста.
Цифры — Regular, графит, трекинг §6, % меньшим кеглем в строке (без переноса).
Зелёный — карточкой/чипом (не за голой цифрой).
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "..", "scripts"))
from design_tokens import load_tokens          # noqa: E402
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu, Pt                    # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE          # noqa: E402
from pptx.oxml.ns import qn                     # noqa: E402
from pptx.dml.color import RGBColor             # noqa: E402

T = load_tokens()
EMU = T.EMU
G_GRAPHITE, G_GREEN, G_GRAY, G_TGRAY = "222222", "26D07C", "F2F2F2", "5C5C5C"


def emu(px):
    return Emu(int(round(px * EMU)))


def rect(slide, x, y, w, h, hexc):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(hexc)
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, s, size, hexc=G_GRAPHITE, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, semibold=False, caps=False):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s.upper() if caps else s
    r.font.name = T.semibold_face if semibold else T.family
    r.font.size = Pt(size); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string(hexc)
    return tb


def bignum(slide, x, y, w, h, val, size, hexc=G_GRAPHITE,
           anchor=MSO_ANCHOR.BOTTOM, align=PP_ALIGN.LEFT):
    """Крупная цифра + опц. % меньшим кеглем в строке, БЕЗ переноса, трекинг §6."""
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    pct = val.endswith("%")
    num = val[:-1] if pct else val
    r = p.add_run(); r.text = num
    r.font.name = T.family; r.font.size = Pt(size); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string(hexc)
    r._r.get_or_add_rPr().set("spc", "-400")
    if pct:
        rp = p.add_run(); rp.text = "%"
        rp.font.name = T.family; rp.font.size = Pt(int(size * 0.5)); rp.font.bold = False
        rp.font.color.rgb = RGBColor.from_string(hexc)
    return tb


def hdr(slide, s):
    txt(slide, 35, 36, 1200, 40, s, 20, G_GRAPHITE, semibold=True, caps=True)


def variant1(slide):
    """Герой-цифра: одна доминантная (масса/размер), две поддержки справа."""
    hdr(slide, "Вариант 1 · герой-цифра + поддержка")
    rect(slide, 35, 150, 46, 46, G_GREEN)                       # бренд-чип у героя
    bignum(slide, 35, 150, 680, 330, "30", 230, anchor=MSO_ANCHOR.TOP)
    txt(slide, 40, 478, 660, 50, "рабочих инструментов в одном окне", 22, G_TGRAY)
    rx = 780
    txt(slide, rx, 165, 440, 26, "вовлечённость", 14, G_TGRAY)
    bignum(slide, rx, 190, 440, 120, "84%", 84, anchor=MSO_ANCHOR.TOP)
    rect(slide, rx, 322, 430, 2, "D9D9D9")
    txt(slide, rx, 345, 440, 26, "сократили инициатив", 14, G_TGRAY)
    bignum(slide, rx, 370, 440, 120, "12", 84, anchor=MSO_ANCHOR.TOP)


def variant2(slide):
    """Карточки во всю высоту: ярлык-таб сверху, цифра у низа. Акцент-карточка зелёная."""
    hdr(slide, "Вариант 2 · карточки во всю высоту, акцент зелёный")
    data = [("12", "AI-инициатив", False), ("84%", "вовлечённость", False),
            ("30", "инструментов", True)]
    L, R, top, h, gap = 35, 1245, 110, 545, 18
    cw = (R - L - gap * 2) / 3
    for i, (val, lab, acc) in enumerate(data):
        x = L + i * (cw + gap)
        rect(slide, x, top, cw, h, G_GREEN if acc else G_GRAY)
        txt(slide, x + 28, top + 26, cw - 56, 26, lab, 15,
            G_GRAPHITE if acc else G_TGRAY, caps=True)
        bignum(slide, x + 24, top + h - 215, cw - 48, 190, val, 130,
               anchor=MSO_ANCHOR.BOTTOM)


def variant3(slide):
    """Базовая линия: крупные цифры стоят НА графитовой линии (уверенный ритм)."""
    hdr(slide, "Вариант 3 · цифры на базовой линии")
    base_y = 472
    rect(slide, 35, base_y, 1210, 3, G_GRAPHITE)               # базовая линия
    data = [("12", "AI-инициатив", False), ("84%", "вовлечённость", False),
            ("30", "инструментов", True)]
    slot = 1210 / 3
    for i, (val, lab, acc) in enumerate(data):
        x = 35 + i * slot
        size = 168 if acc else 134                              # акцент массой/размером
        if acc:
            rect(slide, x, base_y - int(size * 1.0) - 60, 44, 44, G_GREEN)
        bignum(slide, x, base_y - int(size * 1.15), slot - 24, int(size * 1.15),
               val, size, anchor=MSO_ANCHOR.BOTTOM)
        txt(slide, x, base_y + 14, slot - 24, 30, lab, 15, G_TGRAY)


def build():
    prs = Presentation(); prs.slide_width = emu(1280); prs.slide_height = emu(720)
    for fn in (variant1, variant2, variant3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        rect(s, 0, 0, 1280, 720, "FFFFFF")
        fn(s)
    out = os.path.join(HERE, "kpi_variants.pptx")
    prs.save(out); print("saved:", out)
    return out


if __name__ == "__main__":
    build()
