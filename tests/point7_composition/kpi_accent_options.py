#!/usr/bin/env python3
"""
Point 7 — варианты выделения главной KPI-цифры (выбрать акцент, который не слабый
и не «в лоб»). 5 трактовок числа "30" рядом, с подписями.
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "..", "scripts"))
from design_tokens import load_tokens          # noqa: E402
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu, Pt                    # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE          # noqa: E402
from pptx.oxml.ns import qn                     # noqa: E402
from pptx.dml.color import RGBColor             # noqa: E402

T = load_tokens()
EMU = T.EMU


def emu(px):
    return Emu(int(round(px * EMU)))


def rect(slide, x, y, w, h, hexc):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(hexc.lstrip("#"))
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def number(slide, x, y, w, h, text, size, hexc="222222", spc=-400):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = T.semibold_face; r.font.size = Pt(size); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string(hexc)
    if spc:
        r._r.get_or_add_rPr().set("spc", str(spc))


def lbl(slide, x, y, w, s, semibold=False):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(40))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s
    r.font.name = T.semibold_face if semibold else T.family
    r.font.size = Pt(13); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string("5C5C5C")


def headline(slide, x, y, w, s):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(40))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = s.upper()
    r.font.name = T.semibold_face; r.font.size = Pt(20); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string("222222")


def build():
    prs = Presentation(); prs.slide_width = emu(1280); prs.slide_height = emu(720)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 1280, 720, "FFFFFF")
    headline(s, 35, 38, 1200, "KPI-акцент · варианты выделения главной цифры")
    lbl(s, 35, 78, 700, "Слева — без акцента. Какой вариант для главной цифры? (1=в лоб … 5=тонко)")

    cols = 5
    cw = (1245 - 35) / cols
    base_x = 35
    ny, nh = 200, 200

    # 1 — без акцента
    cx = base_x + 0 * cw
    number(s, cx, ny, cw, nh, "30", 96)
    lbl(s, cx, ny + nh + 10, cw, "0 · без акцента\n(графит)")

    # 2 — зелёная плашка (в лоб)
    cx = base_x + 1 * cw
    rect(s, cx + 20, ny, cw - 40, nh, "26D07C")
    number(s, cx, ny, cw, nh, "30", 96)
    lbl(s, cx, ny + nh + 10, cw, "1 · зелёная плашка\n(в лоб)")

    # 3 — зелёный маркер тесно (highlight за глифами)
    cx = base_x + 2 * cw
    rect(s, cx + cw / 2 - 70, ny + 55, 140, 95, "26D07C")
    number(s, cx, ny, cw, nh, "30", 96)
    lbl(s, cx, ny + nh + 10, cw, "2 · зелёный маркер\n(тесно за цифрой)")

    # 4 — крупнее размером (контраст массой, без цвета)
    cx = base_x + 3 * cw
    number(s, cx, ny - 20, cw, nh + 40, "30", 120)
    lbl(s, cx, ny + nh + 10, cw, "3 · крупнее размером\n(контраст массой)")

    # 5 — графит + тонкая зелёная подложка-кромка снизу (subtle)
    cx = base_x + 4 * cw
    rect(s, cx + cw / 2 - 60, ny + 150, 120, 14, "26D07C")
    number(s, cx, ny, cw, nh, "30", 96)
    lbl(s, cx, ny + nh + 10, cw, "4 · цифра + узкая\nзелёная подложка снизу")

    out = os.path.join(HERE, "kpi_accent_options.pptx")
    prs.save(out); print("saved:", out)
    return out


if __name__ == "__main__":
    build()
