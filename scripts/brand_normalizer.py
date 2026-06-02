#!/usr/bin/env python3
"""
brand_normalizer.py — снап ЛЮБОГО (импортированного/чужого) оформления к токенам
Cloud.ru. Ядро extraction-layer (Point 2): из референса берём СТРУКТУРУ/контент,
а «скин» прогоняем сюда — на выходе бренд.

Что делает normalize_slide(slide):
  • geometry  — скругления (roundRect/snip…) → прямой прямоугольник (rounded.none=0);
  • effects   — тени/glow/reflection → снять (effects_util; плоско);
  • fill      — градиент → сплошной; произвольный srgb-цвет → БЛИЖАЙШИЙ цвет палитры;
  • line      — произвольный srgb-цвет линии → ближайший нейтраль;
  • text      — чужой шрифт → SB Sans Display; bold-флаг → начертание Semibold;
                italic/underline → снять; цвет текста → графит/белый (ближайший).
Цвета палитры/таргеты тянутся из design_tokens (palette.json). Идемпотентно:
повторный прогон бренд-слайда ничего не портит.

CLI:  python3 brand_normalizer.py in.pptx out.pptx
"""
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor

from design_tokens import load_tokens
from effects_util import strip_effects

T = load_tokens()

# ---- целевые палитры для снапа ---------------------------------------------
_FILL_NAMES = ("Green", "Yellow", "Purple", "Blue", "Black", "White", "Gray")
FILL_TARGETS = {n: T.hex(n) for n in _FILL_NAMES}
TEXT_TARGETS = {"Black": T.hex("Black"), "White": T.hex("White")}
NEUTRAL_LINE_TARGETS = {                       # линии → нейтрали
    "Black": T.hex("Black"),
    "arrow": T.hex("arrow"),
    "separator": T.hex("separator"),
    "Gray": T.hex("Gray"),
}
BRAND_FONT = T.family                          # "SB Sans Display"
BRAND_SEMIBOLD = T.semibold_face
_BRAND_FONT_PREFIXES = ("SB Sans", "Verdana", "Pingfang", "PingFang", "SimSun")

_ROUNDED_PRESETS = {
    "roundRect", "round1Rect", "round2SameRect", "round2DiagRect",
    "snip1Rect", "snip2SameRect", "snip2DiagRect", "snipRoundRect",
}


def _rgb_tuple(hex6):
    h = hex6.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _nearest(hex6, targets):
    """Ближайший по RGB-расстоянию цвет из targets {name:hex}. → hex."""
    r, g, b = _rgb_tuple(hex6)
    best, bestd = None, 1e18
    for _, thex in targets.items():
        tr, tg, tb = _rgb_tuple(thex)
        d = (r - tr) ** 2 + (g - tg) ** 2 + (b - tb) ** 2
        if d < bestd:
            best, bestd = thex, d
    return best


# Шкала нейтралей (из брендбука: White + Grey1/2/3). Нейтрали доминируют по массе.
NEUTRAL_TARGETS = {
    "White":  "#FFFFFF",
    "Grey1":  "#F2F2F2",
    "Grey2":  "#D9D9D9",
    "Grey3":  "#434343",
    "Graphite": "#222222",
}
ADDITIONAL_TARGETS = {"Yellow": T.hex("Yellow"), "Purple": T.hex("Purple"), "Blue": T.hex("Blue")}
GREEN_HEX = T.hex("Green")
GREY1 = "#F2F2F2"


def _dist(h1, h2):
    a, b = _rgb_tuple(h1), _rgb_tuple(h2)
    return ((a[0] - b[0]) ** 2 + (a[1] - b[1]) ** 2 + (a[2] - b[2]) ** 2) ** 0.5


def _sat(hex6):
    r, g, b = _rgb_tuple(hex6)
    return max(r, g, b) - min(r, g, b)


def _brand_snap_fill(hex6):
    """Снап заливки с УЧЁТОМ масс-иерархии бренда: нейтрали доминируют, зелёный —
    главный акцент, дополнительные цвета — редко. Чужой насыщенный цвет, не
    совпадающий явно с green/additional, → НЕЙТРАЛЬ (серый), а не случайный акцент."""
    if _sat(hex6) <= 36:                       # малонасыщенный → нейтраль (по светлоте)
        return _nearest(hex6, NEUTRAL_TARGETS)
    if _dist(hex6, GREEN_HEX) <= 90:           # зеленоватый → green (главный акцент)
        return GREEN_HEX
    for _, thex in ADDITIONAL_TARGETS.items():  # только ПОЧТИ точное совпадение
        if _dist(hex6, thex) <= 40:
            return thex
    return GREY1                                # иначе чужой акцент → нейтраль


def _is_brand_font(name):
    return bool(name) and any(name.startswith(p) for p in _BRAND_FONT_PREFIXES)


def _spPr(shape):
    return shape._element.find(qn("p:spPr"))


def _shape_fill_hex(shape):
    """Текущий solidFill srgbClr фигуры → '#RRGGBB' или None."""
    spPr = _spPr(shape)
    if spPr is None:
        return None
    sf = spPr.find(qn("a:solidFill"))
    if sf is None:
        return None
    srgb = sf.find(qn("a:srgbClr"))
    return "#" + srgb.get("val") if srgb is not None else None


# ---- geometry: убрать скругления -------------------------------------------
def _deround(shape):
    spPr = _spPr(shape)
    if spPr is None:
        return 0
    geom = spPr.find(qn("a:prstGeom"))
    if geom is not None and geom.get("prst") in _ROUNDED_PRESETS:
        geom.set("prst", "rect")
        av = geom.find(qn("a:avLst"))
        if av is not None:
            for c in list(av):
                av.remove(c)
        return 1
    return 0


# ---- fill: градиент→сплошной; произвольный srgb→ближайший палитровый --------
def _normalize_fill(shape):
    spPr = _spPr(shape)
    if spPr is None:
        return 0
    grad = spPr.find(qn("a:gradFill"))
    if grad is not None:
        first = grad.find(".//" + qn("a:srgbClr"))
        src = "#" + first.get("val") if first is not None else "#222222"
        snapped = _brand_snap_fill(src)          # масс-иерархия: чужое → нейтраль/green
        shape.fill.solid()                       # заменяет gradFill на solidFill
        shape.fill.fore_color.rgb = RGBColor(*_rgb_tuple(snapped))
        return 1
    sf = spPr.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None:
            snapped = _brand_snap_fill("#" + srgb.get("val"))
            srgb.set("val", snapped.lstrip("#"))
            return 1
    return 0


def _cap_additional_le_green(slide):
    """Масс-правило: суммарная площадь дополнительных цветов (Yellow/Purple/Blue)
    ≤ площади green. Если больше — крупнейшие доп-фигуры демотим в нейтраль (Grey1).
    Возвращает число демотированных фигур."""
    add_hexes = {h.upper() for h in ADDITIONAL_TARGETS.values()}
    green_area, addl = 0.0, []
    for sh in _iter_all(slide.shapes):
        hx = _shape_fill_hex(sh)
        if not hx:
            continue
        try:
            a = (sh.width / 9525.0) * (sh.height / 9525.0)
        except Exception:
            a = 0.0
        if hx.upper() == GREEN_HEX.upper():
            green_area += a
        elif hx.upper() in add_hexes:
            addl.append((a, sh))
    addl_area = sum(a for a, _ in addl)
    addl.sort(key=lambda t: -t[0])
    demoted, i = 0, 0
    while addl_area > green_area + 1e-6 and i < len(addl):
        a, sh = addl[i]
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        addl_area -= a
        demoted += 1
        i += 1
    return demoted


# ---- line: произвольный srgb линии → ближайший нейтраль ---------------------
def _normalize_line(shape):
    spPr = _spPr(shape)
    if spPr is None:
        return 0
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return 0
    sf = ln.find(qn("a:solidFill"))
    if sf is None:
        return 0
    srgb = sf.find(qn("a:srgbClr"))
    if srgb is None:
        return 0
    cur = "#" + srgb.get("val")
    # зелёный/палитровый цвет линии оставляем (бренд допускает зелёные линии-
    # композиции/стрелки); снапим только ЯВНО чужие к ближайшему нейтралю,
    # если цвет далёк от любого палитрового.
    if _nearest(cur, FILL_TARGETS) == cur.upper():
        return 0
    srgb.set("val", _nearest(cur, NEUTRAL_LINE_TARGETS).lstrip("#"))
    return 1


# ---- text runs --------------------------------------------------------------
def _normalize_runs(shape):
    n = 0
    if not shape.has_text_frame:
        return 0
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            rPr = run._r.find(qn("a:rPr"))
            if rPr is None:
                continue
            # italic / underline → снять
            if rPr.get("i") == "1":
                rPr.set("i", "0"); n += 1
            if rPr.get("u") not in (None, "none"):
                rPr.set("u", "none"); n += 1
            # bold-флаг → начертание Semibold (эмфаза через face)
            semibold = rPr.get("b") == "1"
            if semibold:
                rPr.set("b", "0"); n += 1
            # шрифт → бренд (если чужой)
            for tag in ("latin", "ea", "cs"):
                el = rPr.find(qn(f"a:{tag}"))
                if el is not None and not _is_brand_font(el.get("typeface")):
                    el.set("typeface", BRAND_SEMIBOLD if semibold else BRAND_FONT)
                    n += 1
                elif el is not None and semibold:
                    el.set("typeface", BRAND_SEMIBOLD)
            # цвет текста → графит/белый (ближайший)
            sf = rPr.find(qn("a:solidFill"))
            if sf is not None:
                srgb = sf.find(qn("a:srgbClr"))
                if srgb is not None:
                    srgb.set("val", _nearest("#" + srgb.get("val"), TEXT_TARGETS).lstrip("#"))
                    n += 1
    return n


def _iter_all(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # GROUP
            try:
                yield from _iter_all(sh.shapes)
            except Exception:
                pass


def normalize_slide(slide, dark=False):
    """Привести оформление всех фигур слайда к токенам Cloud.ru. → dict со счётчиками.

    Двухслойно: (1) снап «скина» (геометрия/эффекты/заливка/линия/шрифт/начертание/
    цвет текста к графит-белый); (2) fill-aware КОНТРАСТ через enforce_canonical —
    белый-на-светлом/зелёном → графит (брендовое правило: текст всегда #222222,
    белый только на тёмном). Так читаемость и канон гарантированы."""
    st = {"deround": 0, "effects": 0, "fill": 0, "line": 0, "text": 0, "contrast": 0}
    for sh in _iter_all(slide.shapes):
        st["deround"] += _deround(sh)
        if strip_effects(sh._element):
            st["effects"] += 1
        try:
            st["fill"] += _normalize_fill(sh)
        except Exception:
            pass
        try:
            st["line"] += _normalize_line(sh)
        except Exception:
            pass
        st["text"] += _normalize_runs(sh)
    # Масс-правило: дополнительные цвета ≤ green (избыток → нейтраль)
    st["demoted_additional"] = _cap_additional_le_green(slide)
    # Слой 2: fill-aware контраст/канон (white-on-green→графит, semibold, мин. размер)
    try:
        import enforce_canonical
        cs = enforce_canonical.enforce_canonical_slide(slide, dark=dark)
        st["contrast"] = cs.get("green_text", 0) + cs.get("white_on_light", 0)
    except Exception:
        pass
    return st


def normalize_pptx(in_path, out_path):
    prs = Presentation(in_path)
    total = {"deround": 0, "effects": 0, "fill": 0, "line": 0, "text": 0}
    for slide in prs.slides:
        st = normalize_slide(slide)
        for k in total:
            total[k] += st[k]
    prs.save(out_path)
    return total


def main():
    if len(sys.argv) != 3:
        print("Usage: brand_normalizer.py in.pptx out.pptx", file=sys.stderr)
        sys.exit(1)
    st = normalize_pptx(sys.argv[1], sys.argv[2])
    print("normalized:", st)
    print("saved:", sys.argv[2])


if __name__ == "__main__":
    main()
