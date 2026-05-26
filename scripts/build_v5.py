#!/usr/bin/env python3
"""
build_v5.py — slot-based + style override (size/color/bold/font).

v0.5 улучшения над v0.4:
- Поддержка `slot_styles_override` в plan.json — можно задать кастомный размер/цвет/bold для конкретного слота
- Default-стиль читается из donor-slot-map.yaml (если override не указан, берём оттуда)
- Поддержка multi-line через нормальный \n (не \x0b как в исходных placeholders шаблона)
- Очистка vertical tab \x0b в исходниках донора при clear_text_frame

Usage:
    python3 build_v5.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]

plan.json:
{
  "slides": [
    {
      "clone_from_slide": 5,
      "slots": {
        "title": "Заголовок"
      },
      "slot_styles_override": {
        "title": {"size_pt": 44, "color": "#26D07C", "bold": true}
      },
      "_donor_category": "title_white_with_3d"
    }
  ]
}
"""
import sys, json
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
from lxml import etree

try:
    import yaml
except ImportError:
    print("ERROR: PyYAML required. pip3 install --user pyyaml", file=sys.stderr)
    sys.exit(1)


def load_donor_map(path):
    with open(path, encoding="utf-8") as f:
        data = yaml.safe_load(f)
    return data["donors"]


def get_text_frame_by_shape_idx(slide, shape_idx):
    shapes = list(slide.shapes)
    if shape_idx >= len(shapes):
        return None
    sh = shapes[shape_idx]
    if not sh.has_text_frame:
        return None
    return sh.text_frame


def hex_to_rgb(hex_color):
    """#RRGGBB → RGBColor"""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def replace_text_with_style(text_frame, new_text, style_override=None):
    """Заменяет текст в text_frame, СОХРАНЯЯ XML-уровневые run properties (rPr) первого run.
    Опционально применяет style_override (size_pt, color, bold, font).

    Алгоритм:
      1. Берём первый <a:p> и его первый <a:r>
      2. Сохраняем XML-копию <a:rPr> + <a:pPr>
      3. Стираем все existing <a:p>
      4. Для каждой строки (split по \n) создаём <a:p> с одним <a:r>, копируем rPr
      5. Если есть style_override — мутируем атрибуты rPr перед вставкой
    """
    if text_frame is None:
        return False

    txBody = text_frame._txBody

    p_elements = txBody.findall(qn("a:p"))
    if not p_elements:
        return False
    first_p = p_elements[0]
    first_r = first_p.find(qn("a:r"))

    saved_rPr = None
    if first_r is not None:
        rPr_el = first_r.find(qn("a:rPr"))
        if rPr_el is not None:
            saved_rPr = deepcopy(rPr_el)

    saved_pPr = first_p.find(qn("a:pPr"))
    saved_pPr_copy = deepcopy(saved_pPr) if saved_pPr is not None else None

    # Apply style override to saved_rPr
    if style_override and saved_rPr is not None:
        # size_pt → sz attribute (in hundredths of points)
        if "size_pt" in style_override and style_override["size_pt"] is not None:
            saved_rPr.set("sz", str(int(style_override["size_pt"] * 100)))
        # bold → b attribute
        if "bold" in style_override and style_override["bold"] is not None:
            saved_rPr.set("b", "1" if style_override["bold"] else "0")
        # color → <a:solidFill><a:srgbClr val="..."/></a:solidFill>
        if "color" in style_override and style_override["color"] is not None:
            color_hex = style_override["color"].lstrip("#")
            # Remove existing fill
            for fill_tag in ("solidFill", "noFill", "gradFill", "blipFill", "pattFill", "grpFill"):
                el = saved_rPr.find(qn(f"a:{fill_tag}"))
                if el is not None:
                    saved_rPr.remove(el)
            # Add solidFill with srgbClr
            solidFill = etree.SubElement(saved_rPr, qn("a:solidFill"))
            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgbClr.set("val", color_hex.upper())
        # font → <a:latin typeface="..."/>
        if "font" in style_override and style_override["font"] is not None:
            for tag in ("latin", "ea", "cs"):
                el = saved_rPr.find(qn(f"a:{tag}"))
                if el is not None:
                    saved_rPr.remove(el)
            latin = etree.SubElement(saved_rPr, qn("a:latin"))
            latin.set("typeface", style_override["font"])

    # Стираем все existing <a:p>
    for p_el in p_elements:
        txBody.remove(p_el)

    # Конвертируем \x0b (vertical tab) в \n
    new_text = new_text.replace("\v", "\n").replace("\x0b", "\n")
    lines = new_text.split("\n")
    if not lines:
        lines = [""]

    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))
        if saved_pPr_copy is not None:
            new_p.append(deepcopy(saved_pPr_copy))
        new_r = etree.SubElement(new_p, qn("a:r"))
        if saved_rPr is not None:
            new_r.append(deepcopy(saved_rPr))
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line

    return True


def clear_text_frame(text_frame):
    return replace_text_with_style(text_frame, "")


def reorder_and_filter_slides(prs, source_slide_nums):
    sldIdLst = prs.slides._sldIdLst
    all_ids = list(sldIdLst)
    num_to_sldId = {i + 1: el for i, el in enumerate(all_ids)}

    used_nums = []
    seen = set()
    for n in source_slide_nums:
        if n is None:
            continue
        if n in seen:
            print(f"WARN: дубли донор-слайдов не поддерживаются (slide {n} повторяется)", file=sys.stderr)
            continue
        if n not in num_to_sldId:
            print(f"WARN: слайд {n} не найден (1..{len(all_ids)})", file=sys.stderr)
            continue
        used_nums.append(n)
        seen.add(n)

    keep_set = set(used_nums)
    for n, sldId in list(num_to_sldId.items()):
        if n not in keep_set:
            rId = sldId.attrib[qn('r:id')]
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)

    remaining = list(sldIdLst)
    rem_map = {}
    for el in remaining:
        rId = el.attrib[qn('r:id')]
        rem_map[rId] = el
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for n in used_nums:
        sld = num_to_sldId[n]
        rId = sld.attrib[qn('r:id')]
        sldIdLst.append(rem_map[rId])


def build(plan_path, template_path, output_path, donor_map_path):
    plan = json.load(open(plan_path, encoding="utf-8"))
    p = Presentation(template_path)
    donors = load_donor_map(donor_map_path)

    cloned_nums = [s.get("clone_from_slide") for s in plan["slides"]]
    reorder_and_filter_slides(p, cloned_nums)

    cloned_iter = iter(p.slides)
    cloned_slide_for = {}
    for plan_slide, actual in zip(
        [s for s in plan["slides"] if s.get("clone_from_slide")],
        list(p.slides)
    ):
        cloned_slide_for[plan_slide["clone_from_slide"]] = actual

    for plan_slide in plan["slides"]:
        if not plan_slide.get("clone_from_slide"):
            continue
        src_num = plan_slide["clone_from_slide"]
        actual = cloned_slide_for[src_num]
        donor_def = donors.get(src_num)
        if donor_def is None:
            print(f"WARN: нет donor-slot-map для slide {src_num}", file=sys.stderr)
            continue

        slot_defs = donor_def.get("slots", {})
        slots_filled = plan_slide.get("slots", {})
        styles_override = plan_slide.get("slot_styles_override", {})

        # Заполнить слоты с применением style override (если задан)
        for slot_name, new_text in slots_filled.items():
            if slot_name not in slot_defs:
                print(f"WARN: slot '{slot_name}' не определён для donor {src_num}", file=sys.stderr)
                continue
            shape_idx = slot_defs[slot_name]["shape_idx"]
            tf = get_text_frame_by_shape_idx(actual, shape_idx)
            if tf is None:
                print(f"WARN: shape_idx {shape_idx} не найден на slide {src_num}", file=sys.stderr)
                continue
            override = styles_override.get(slot_name)
            replace_text_with_style(tf, new_text, override)

        # Очистить НЕзаполненные слоты
        for slot_name, slot_def in slot_defs.items():
            if slot_name in slots_filled:
                continue
            if slot_def.get("optional"):
                continue
            shape_idx = slot_def["shape_idx"]
            tf = get_text_frame_by_shape_idx(actual, shape_idx)
            if tf is not None:
                clear_text_frame(tf)

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides", file=sys.stderr)


def main():
    if len(sys.argv) < 4:
        print("Usage: build_v5.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]", file=sys.stderr)
        sys.exit(1)
    plan_p = sys.argv[1]
    tpl_p = sys.argv[2]
    out_p = sys.argv[3]
    donor_p = sys.argv[4] if len(sys.argv) > 4 else "pptx-skill/brand/donor-slot-map.yaml"
    build(plan_p, tpl_p, out_p, donor_p)


if __name__ == "__main__":
    main()
