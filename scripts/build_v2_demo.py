#!/usr/bin/env python3
"""
build_v2_demo.py — реальная вёрстка для ai test 2.pptx → 8 слайдов Cloud.ru.

В отличие от v0.1 build_pptx.py, здесь:
- Explicit стили (font.name, size, color) на runs
- Multi-slide split (5 → 8)
- Shapes для KPI больших цифр
- Brand colors из палитры
- Запасные значения если SB Sans недоступен

Usage:
    python3 build_v2_demo.py <template.pptx> <output.pptx>
"""
import sys, os, shutil, copy
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# === Brand palette ===
GREEN = RGBColor(0x26, 0xD0, 0x7C)
BLACK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xF2, 0xF2, 0xF2)
FONT = "SB Sans Display"
FONT_TEXT = "SB Sans Text"


def style_run(run, font_name=FONT, size_pt=18, bold=False, color=BLACK):
    """Применить explicit стили к run."""
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def fill_placeholder(slide, ph_idx, text, font_name=FONT, size_pt=18, bold=False, color=BLACK):
    """Найти placeholder и заполнить с явными стилями."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            tf.text = ""  # очищаем
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            style_run(run, font_name, size_pt, bold, color)
            return ph
    return None


def add_text_box(slide, left, top, width, height, text, font_name=FONT,
                 size_pt=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Создать text box с явными стилями (для shapes-based layout)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    style_run(run, font_name, size_pt, bold, color)
    return tb


def add_filled_rect(slide, left, top, width, height, fill, line_color=None):
    """Прямоугольник с заливкой (без скруглений — брендбук)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    sh.shadow.inherit = False  # без теней
    return sh


def remove_existing_slides(prs):
    """Удаление существующих слайдов из шаблона корректно (с физическими файлами)."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def build(template_path, output_path):
    p = Presentation(template_path)
    remove_existing_slides(p)

    sw, sh = p.slide_width, p.slide_height  # EMU
    # 1280x720 px ≈ 12192000 x 6858000 EMU; 1 px = 9525 EMU

    # ===== Slide 1: Title (idx 52 "Титул / три строки_2") =====
    layout = p.slide_layouts[52]
    s1 = p.slides.add_slide(layout)
    # placeholders idx 52: 3 ph (мы используем title + subtitle/speaker)
    # Просто заполним по факту — поищем placeholder по типу
    title_text = "Как мы перестали бояться факапов и начали учиться на\u00a0них"
    sub_text = "История построения AI-first команды в\u00a0СТО Cloud.ru\nНадежда Погина · Руководитель проектов"
    # idx 52 имеет TITLE и BODY/SUBTITLE
    placeholders = list(s1.placeholders)
    if len(placeholders) >= 1:
        fill_placeholder(s1, placeholders[0].placeholder_format.idx, title_text,
                         FONT, 40, bold=True, color=BLACK)
    if len(placeholders) >= 2:
        fill_placeholder(s1, placeholders[1].placeholder_format.idx, sub_text,
                         FONT_TEXT, 18, bold=False, color=BLACK)

    # ===== Slide 2: Divider (idx 8 "Раздел / Зеленый 1") =====
    layout = p.slide_layouts[8]
    s2 = p.slides.add_slide(layout)
    placeholders = list(s2.placeholders)
    if placeholders:
        fill_placeholder(s2, placeholders[0].placeholder_format.idx,
                         "01\nКОНТЕКСТ", FONT, 100, bold=True, color=BLACK)

    # ===== Slide 3: 3-колонки факты (idx 31) =====
    layout = p.slide_layouts[31]
    s3 = p.slides.add_slide(layout)
    fill_placeholder(s3, 0, "Кто мы и\u00a0что такое\u00a0«СТО»", FONT, 32, bold=True, color=BLACK)
    # idx 31: ph 18,19,20 = headings; ph 12,13,16 = bodies
    headings = [
        (18, "Что такое СТО"),
        (19, "Точка отсчёта"),
        (20, "Цель"),
    ]
    bodies = [
        (12, "Центр эксплуатации и\u00a0развития инфраструктуры Cloud.ru"),
        (13, "6\u00a0месяцев назад: уровень зрелости AI\u00a0—\u00a02,2 из\u00a05.\n40+ разрозненных инициатив,\n0\u00a0системности"),
        (16, "Не\u00a0«внедрить ИИ»,\nа\u00a0построить AI-first команду"),
    ]
    for idx, txt in headings:
        fill_placeholder(s3, idx, txt, FONT, 18, bold=True, color=GREEN)
    for idx, txt in bodies:
        fill_placeholder(s3, idx, txt, FONT_TEXT, 14, bold=False, color=BLACK)

    # ===== Slide 4: Callout (idx 23 "Белый / Важная информация") — цитата =====
    layout = p.slide_layouts[23]
    s4 = p.slides.add_slide(layout)
    placeholders = list(s4.placeholders)
    if len(placeholders) >= 1:
        fill_placeholder(s4, placeholders[0].placeholder_format.idx,
                         "Не\u00a0«внедрить ИИ»,\nа\u00a0построить\nAI-first команду",
                         FONT, 60, bold=True, color=BLACK)
    if len(placeholders) >= 2:
        fill_placeholder(s4, placeholders[1].placeholder_format.idx,
                         "— Цель проекта", FONT_TEXT, 18, bold=False, color=GREEN)

    # ===== Slide 5: Divider тёмный (idx 54) =====
    layout = p.slide_layouts[54]
    s5 = p.slides.add_slide(layout)
    placeholders = list(s5.placeholders)
    if placeholders:
        fill_placeholder(s5, placeholders[0].placeholder_format.idx,
                         "02\nРЕЗУЛЬТАТЫ", FONT, 100, bold=True, color=GREEN)

    # ===== Slide 6: KPI большие цифры через shapes (на base layout idx 10) =====
    layout = p.slide_layouts[10]  # Контент / Белый 1
    s6 = p.slides.add_slide(layout)
    fill_placeholder(s6, 0, "Что получилось через 6\u00a0месяцев", FONT, 32, bold=True, color=BLACK)
    # 3 KPI вручную через shapes
    px = 9525  # 1 px in EMU
    margin_left = 60 * px
    margin_top = 200 * px
    col_w = (1280 - 60*2) // 3 * px  # ширина колонки
    gap = 0
    for i, (num, label) in enumerate([
        ("2,2 → 4", "Уровень зрелости AI\nиз\u00a05 баллов"),
        ("84%", "Вовлечённости\nкоманды"),
        ("5", "Рабочих инструментов\nкомандой освоено"),
    ]):
        x = margin_left + i * (col_w + gap)
        # Большая цифра
        add_text_box(s6, x, margin_top, col_w, 200*px, num,
                     FONT, 80, bold=True, color=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        # Подпись
        add_text_box(s6, x, margin_top + 220*px, col_w, 80*px, label,
                     FONT_TEXT, 16, bold=False, color=BLACK, align=PP_ALIGN.LEFT)
        # Вертикальная линия-разделитель (кроме последней)
        if i < 2:
            line_x = x + col_w
            add_filled_rect(s6, line_x, margin_top, 1*px, 280*px, GRAY)

    # ===== Slide 7: 4 блока механизмов (idx 28) =====
    layout = p.slide_layouts[28]
    s7 = p.slides.add_slide(layout)
    fill_placeholder(s7, 0, "Механизмы интеграции и\u00a0синергия", FONT, 32, bold=True, color=BLACK)
    # idx 28 — 4 текстовых блока, ph idx неявно — пробежим по всем доступным
    mechanisms = [
        "Масштабируем\nлучшие практики",
        "Находим точки синергии\nвместо дублирования",
        "Систематизируем AI-активности\nчерез единого представителя",
        "Делимся\nэкспертизой",
    ]
    body_phs = [ph for ph in s7.placeholders if ph.placeholder_format.idx not in (0, 11)]
    for ph, txt in zip(body_phs, mechanisms):
        fill_placeholder(s7, ph.placeholder_format.idx, txt, FONT, 18, bold=False, color=BLACK)

    # ===== Slide 8: Финальный логотип (idx 95 "Зеленый / Логотип") =====
    layout = p.slide_layouts[95]
    p.slides.add_slide(layout)

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides", file=sys.stderr)


def main():
    if len(sys.argv) != 3:
        print("Usage: build_v2_demo.py <template.pptx> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    build(sys.argv[1], sys.argv[2])


if __name__ == "__main__":
    main()
