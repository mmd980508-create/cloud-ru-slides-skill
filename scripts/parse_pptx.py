#!/usr/bin/env python3
"""
parse_pptx.py — извлекает структуру .pptx в JSON для агентов.

Usage:
    python3 parse_pptx.py <input.pptx> [output.json]

Output:
    JSON со структурой:
    {
      "file": "...",
      "slide_count": N,
      "slide_size": {"width_emu": ..., "height_emu": ...},
      "slides": [
        {
          "num": 1,
          "layout_name": "...",
          "layout_idx_in_master": <int|null>,
          "title": "...",
          "body": ["...", "..."],
          "text_runs": [...],
          "images": [{"name": "...", "left_emu": ..., "top_emu": ..., "width_emu": ..., "height_emu": ...}],
          "shapes_count": N,
          "tables_count": N
        }
      ]
    }
"""
import sys, json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def parse(input_path):
    p = Presentation(input_path)
    result = {
        "file": input_path,
        "slide_count": len(p.slides),
        "slide_size": {"width_emu": p.slide_width, "height_emu": p.slide_height,
                       "width_px_at96": int(p.slide_width / 9525), "height_px_at96": int(p.slide_height / 9525)},
        "slides": [],
    }

    # Build layout idx lookup (SlideLayout is not hashable — use id())
    layout_to_idx = {id(lay): i for i, lay in enumerate(p.slide_layouts)}

    for snum, slide in enumerate(p.slides, start=1):
        sdata = {
            "num": snum,
            "layout_name": slide.slide_layout.name,
            "layout_idx_in_master": layout_to_idx.get(id(slide.slide_layout)),
            "title": None,
            "body": [],
            "text_runs": [],
            "images": [],
            "shapes_count": 0,
            "tables_count": 0,
        }

        for shape in slide.shapes:
            sdata["shapes_count"] += 1

            # Title placeholder
            if shape.has_text_frame and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                txt = shape.text_frame.text.strip()
                if not txt:
                    continue
                # Title placeholder type values: TITLE=13/CTR_TITLE=15
                if str(ph_type) in ("TITLE (13)", "CENTER_TITLE (15)") or "TITLE" in str(ph_type).upper():
                    if not sdata["title"]:
                        sdata["title"] = txt
                    else:
                        sdata["body"].append(txt)
                else:
                    sdata["body"].append(txt)
                continue

            # Generic text frames (not placeholder)
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    sdata["text_runs"].append(txt)
                    if not sdata["title"] and len(txt) < 100:
                        sdata["title"] = txt
                    else:
                        sdata["body"].append(txt)

            # Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sdata["images"].append({
                    "name": shape.name,
                    "left_emu": shape.left, "top_emu": shape.top,
                    "width_emu": shape.width, "height_emu": shape.height,
                })

            # Tables
            if shape.has_table:
                sdata["tables_count"] += 1

        result["slides"].append(sdata)

    return result


def main():
    if len(sys.argv) < 2:
        print("Usage: parse_pptx.py <input.pptx> [output.json]", file=sys.stderr)
        sys.exit(1)
    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    data = parse(input_path)
    out = json.dumps(data, ensure_ascii=False, indent=2)

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(out)
        print(f"Wrote {output_path} ({data['slide_count']} slides)", file=sys.stderr)
    else:
        print(out)


if __name__ == "__main__":
    main()
