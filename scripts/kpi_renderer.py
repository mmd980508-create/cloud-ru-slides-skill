#!/usr/bin/env python3
"""
kpi_renderer.py — canonical KPI rendering для build_v9.

Зачем: donor 43/44 в шаблоне = guide-слайды (показывают примеры размеров),
не доноры для произвольных KPI. Их placeholders расположены обучающе
(199pt + 199pt + 44pt) — не для real data.

Решение: рисуем KPI shapes с нуля на blank donor с CLEAN GRID:
- 1 number → hero center
- 2 numbers → 50/50 grid
- 3 numbers → 33/33/33 grid (все одного размера)
- Title 32pt SemiBold top
- Numbers 130pt SemiBold (помещается без overflow)
- Descriptions 16pt Regular under numbers
- Accent: 1 of N number = #26D07C (canonical §12)

Layout (1280×720 slide):
  Title:           (35, 38, 1209, 53) px       32pt SemiBold #222222
  Number row:      y=220, height=280
    1 number:      x=440, w=400 (centered)
    2 numbers:     x=140, w=400 + x=740, w=400
    3 numbers:     x=35, w=400 + x=440, w=400 + x=840, w=400
  Desc row:        y=530, height=100
    Same x columns
  Optional %:      next to number, top-right corner of number box
"""
import os
import sys
import json
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def _load_template_version():
    """Грузит brand/template-version.json (маппинг slide-индексов). Fallback — None."""
    here = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(here, "..", "brand", "template-version.json"),
        os.path.join(os.getcwd(), "brand", "template-version.json"),
        os.path.join(os.getcwd(), "pptx-skill", "brand", "template-version.json"),
    ]
    for path in candidates:
        try:
            if os.path.isfile(path):
                with open(path, encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            continue
    return None


_TPL_VER = _load_template_version()


GRAPHITE = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x26, 0xD0, 0x7C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xF2, 0xF2, 0xF2)

FONT = "SB Sans Display"
# Полужирное = отдельный font face (встроен в шаблон), НЕ bold-флаг.
# Canonical 2026-05-29 (Problem #3): Bold запрещён — эмфаза только через SemiBold.
FONT_SEMIBOLD = "SB Sans Display Semibold"


# Canonical §4 (slide 43): фрейм 199pt помещает максимум 2 значащих цифры.
KPI_HERO_DIGIT_LIMIT = 2

# Брендбук §6: у крупных цифр отрицательный трекинг. spc в 1/100 pt (user 2026-06-02).
BIG_NUMBER_SPC = -400


def _count_significant_digits(value):
    """Считает значащие цифры в KPI value, игнорируя разделители и unit-suffix.

    Примеры:
      '12'    → 2 цифры
      '99%'   → 2 (% игнорируется — pct идёт отдельным флагом)
      '1.5'   → 2 (точка игнорируется, но 1+5=2)
      '1.5K'  → 2 (K — буква-сокращение, не считается)
      '150'   → 3 (overflow!)
      '~150'  → 3 (тильда игнорируется)
      '12,5'  → 3 (русская десятичная — 3 цифры в frame не помещаются)
    """
    if value is None:
        return 0
    s = str(value)
    digits = "".join(c for c in s if c.isdigit())
    return len(digits)


def _add_text_box(slide, left_px, top_px, width_px, height_px, text,
                  font_size_pt=14, font_name=None,
                  bold=False, color=GRAPHITE, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, spc=None):
    """Add text box with canonical font/size/color.

    Эмфаза (bold=True) реализуется через начертание SemiBold, а не bold-флаг
    (Problem #3 2026-05-29). font_name=None → авто: Regular или Semibold по bold.
    """
    EMU = 9525
    box = slide.shapes.add_textbox(
        Emu(left_px * EMU), Emu(top_px * EMU),
        Emu(width_px * EMU), Emu(height_px * EMU)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    # SemiBold через face, Bold-флаг не используем
    if font_name is None:
        run.font.name = FONT_SEMIBOLD if bold else FONT
    else:
        run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = False
    run.font.color.rgb = color
    # letter-spacing (1/100 pt, может быть отрицательным) — для крупных цифр
    # брендбук §6 требует отрицательный трекинг.
    if spc is not None:
        run._r.get_or_add_rPr().set("spc", str(int(spc)))
    return box


def _add_number_box(slide, x, y, w, h, value, pct=False, size=130, color=GRAPHITE,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Крупная цифра (Regular, отрицательный трекинг §6) + опц. знак % в строке
    (меньшим кеглем). Выравнивание задаётся align/anchor — по умолчанию верх-лево
    (редакторский край-стиль, не центр)."""
    EMU = 9525
    box = slide.shapes.add_textbox(Emu(int(x) * EMU), Emu(int(y) * EMU),
                                   Emu(int(w) * EMU), Emu(int(h) * EMU))
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(value)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = False
    r.font.color.rgb = color
    r._r.get_or_add_rPr().set("spc", str(int(BIG_NUMBER_SPC)))
    if pct:
        rp = p.add_run()
        rp.text = "%"
        rp.font.name = FONT
        rp.font.size = Pt(max(28, size // 3))
        rp.font.bold = False
        rp.font.color.rgb = color
    return box


def _add_accent_bar(slide, left_px, top_px, width_px, height_px, color=GREEN):
    """Зелёная подчёркивающая плашка-акцент под главной KPI-цифрой.

    Canonical 2026-05-29 (Problem #2): акцент НЕ цветом текста, а отдельным
    цветным элементом. Сама цифра остаётся #222222 (светлый) / белой (тёмный),
    а «главный показатель» помечается этой зелёной плашкой.
    """
    EMU = 9525
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(left_px * EMU), Emu(top_px * EMU),
        Emu(width_px * EMU), Emu(height_px * EMU),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # без рамки
    # Снять любые эффекты — и явные, и тематические (effectRef). Плоско, без исключений.
    from effects_util import strip_effects
    strip_effects(shape._element)
    return shape


def _is_title_placeholder(sp):
    """sp (<p:sp>) — это title/ctrTitle placeholder шаблона?"""
    nv = sp.find(qn("p:nvSpPr"))
    if nv is None:
        return False
    nvpr = nv.find(qn("p:nvPr"))
    if nvpr is None:
        return False
    ph = nvpr.find(qn("p:ph"))
    if ph is None:
        return False
    return ph.get("type") in ("title", "ctrTitle")


def clean_slide_to_blank(slide, keep_title=True):
    """Удалить все shapes на slide, КРОМЕ title-placeholder шаблона (если
    keep_title). Title-placeholder сохраняется и очищается от донорского текста,
    чтобы заголовок вписывался в штатное место шаблона (Problem #6, 2026-05-29).
    Layout-inherited (logo, footer) останутся."""
    spTree = slide.shapes._spTree
    to_remove = []
    for child in list(spTree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'cxnSp', 'grpSp'):
            if keep_title and tag == 'sp' and _is_title_placeholder(child):
                # сохранить placeholder, но очистить донорский текст
                txBody = child.find(qn('p:txBody'))
                if txBody is not None:
                    for p_el in txBody.findall(qn('a:p')):
                        for r_el in p_el.findall(qn('a:r')):
                            p_el.remove(r_el)
                continue
            to_remove.append(child)
    for el in to_remove:
        spTree.remove(el)


def set_slide_title(slide, text, dark=False):
    """Вписать заголовок слайда в штатный TITLE-placeholder шаблона.

    Canonical (Problem #6, 2026-05-29): единая позиция/размер из шаблона
    (35,38 / 963×54 / 20pt SemiBold CAPS, графит/белый). Если placeholder на
    слайде нет — fallback на канонический textbox той же геометрии.
    """
    if not text:
        return None
    color = WHITE if dark else GRAPHITE
    title_ph = None
    try:
        title_ph = slide.shapes.title
    except Exception:
        title_ph = None
    if title_ph is None:
        return _add_text_box(slide, 35, 38, 963, 54, text.upper(),
                             font_size_pt=20, bold=True, color=color,
                             anchor=MSO_ANCHOR.MIDDLE)
    tf = title_ph.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = text.upper()
    run.font.name = FONT_SEMIBOLD
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = color
    return title_ph


def render_kpi(slide, kpi_config, dark=False):
    """
    Render KPI block on slide.

    kpi_config = {
        "title": "Что получилось",
        "numbers": [
            {"value": "12", "desc": "сократили AI-инициатив", "accent": False, "pct": False},
            {"value": "84", "desc": "вовлечённость", "accent": False, "pct": True},
            {"value": "5",  "desc": "инструментов", "accent": True, "pct": False}
        ]
    }

    dark: True if slide is dark (text white instead of graphite)
    """
    text_color = WHITE if dark else GRAPHITE
    n = len(kpi_config.get("numbers", []))

    if n == 0 or n > 3:
        raise ValueError(f"KPI supports 1-3 numbers, got {n}")

    # Card-layout (user-decision 2026-06-02): все цифры в отдельных карточках,
    # акцентная — зелёная. ТОЛЬКО в этом варианте допустима зелёная акцент-плашка
    # (у «голых» крупных цифр акцента нет).
    if kpi_config.get("layout") == "cards":
        return render_kpi_cards(slide, kpi_config, dark=dark)

    # Title — в штатный placeholder шаблона (Problem #6)
    title = kpi_config.get("title", "")
    if title:
        set_slide_title(slide, title, dark=dark)

    # Number boxes layout
    if n == 1:
        x_positions = [440]
        block_width = 400
    elif n == 2:
        x_positions = [140, 740]
        block_width = 400
    else:  # n == 3
        x_positions = [35, 440, 845]
        block_width = 400

    NUMBER_TOP = 200
    NUMBER_HEIGHT = 260
    DESC_TOP = 470
    DESC_HEIGHT = 120
    NUMBER_FONT = 130 if n >= 2 else 199  # single hero number может быть 199pt

    # Reduce font for 3 numbers if they are wide
    if n == 3:
        max_chars = max(len(x["value"]) for x in kpi_config["numbers"])
        if max_chars > 3:
            NUMBER_FONT = 100  # 4+ digits — smaller

    for i, num in enumerate(kpi_config["numbers"]):
        x = x_positions[i]
        is_accent = num.get("accent", False)
        # Canonical 2026-05-29 (Problem #2): цифра ВСЕГДА #222222 (светлый) /
        # белая (тёмный) — НЕ зелёная. Акцент главного показателя выносится в
        # отдельную зелёную плашку-подчёркивание под цифрой (см. ниже).
        color = text_color
        value = num["value"]
        has_pct = num.get("pct", False)

        # Canonical §4: 199pt frame ⇒ max 2 значащих цифры. Иначе overflow.
        # Применяем только к hero (199pt single number); для 130/100pt запас больше.
        if NUMBER_FONT >= 199:
            digits = _count_significant_digits(value)
            if digits > KPI_HERO_DIGIT_LIMIT:
                print(
                    f"WARN: KPI value '{value}' содержит {digits} значащих цифр "
                    f"при 199pt frame (canonical max {KPI_HERO_DIGIT_LIMIT}). "
                    f"Используй сокращение ('1.5K', '~{value[:2]}') или сменить layout (donor 47).",
                    file=sys.stderr,
                )

        # АКЦЕНТ у «голых» крупных цифр НЕ ставим (user-decision 2026-06-02):
        # большие цифры на слайде — без акцента (ни плашки, ни полосы). Зелёная
        # акцентная плашка допустима ТОЛЬКО когда все цифры помещены в отдельные
        # карточки (card-layout) — это отдельный вариант, не этот «голый» рендер.
        # Начертание крупной цифры — Regular (bold=False), брендово.
        _ = is_accent

        # Цифра + % одним боксом, выравнивание по ЛЕВОМУ краю сетки (не центр) —
        # редакторский край-стиль (user-decision 2026-06-02). Трекинг §6.
        _add_number_box(slide, x, NUMBER_TOP, block_width, NUMBER_HEIGHT,
                        value, pct=has_pct, size=NUMBER_FONT, color=color,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        # Description below — тоже по левому краю
        desc = num.get("desc", "")
        if desc:
            _add_text_box(slide, x, DESC_TOP, block_width, DESC_HEIGHT,
                          desc, font_size_pt=16, bold=False, color=text_color,
                          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


def render_kpi_cards(slide, kpi_config, dark=False):
    """KPI в карточках: каждая цифра в отдельной карточке #F2F2F2, акцентная —
    ЗЕЛЁНАЯ (графитовая цифра на зелёном; white-on-green запрещён). Только в этом
    варианте допустима зелёная акцент-плашка (user-decision 2026-06-02). Цифра —
    Regular-начертание (bold=False), отрицательный трекинг (§6)."""
    numbers = kpi_config.get("numbers", [])
    n = len(numbers)
    title = kpi_config.get("title", "")
    if title:
        set_slide_title(slide, title, dark=dark)

    SAFE_L, SAFE_R = 35, 1245
    gap = 20
    cw = int((SAFE_R - SAFE_L - gap * (n - 1)) / n)
    CARD_Y, CARD_H = 175, 430
    NUMBER_FONT = 130 if n >= 2 else 199
    if n == 3 and max(len(x["value"]) for x in numbers) > 3:
        NUMBER_FONT = 100

    PAD = 28
    for i, num in enumerate(numbers):
        x = SAFE_L + i * (cw + gap)
        is_accent = num.get("accent", False)
        # карточка (filled rect, без эффектов) — рисуется ПЕРВОЙ, контент поверх
        _add_accent_bar(slide, x, CARD_Y, cw, CARD_H, color=(GREEN if is_accent else GRAY))
        # Выравнивание по КРАЯМ карточки (не центр): цифра — верх-лево, подпись —
        # низ-лево. Цифра графит Regular (читаемо на сером и зелёном), трекинг §6.
        _add_number_box(slide, x + PAD, CARD_Y + PAD, cw - 2 * PAD, 230,
                        num["value"], pct=num.get("pct", False),
                        size=NUMBER_FONT, color=GRAPHITE,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        desc = num.get("desc", "")
        if desc:
            _add_text_box(slide, x + PAD, CARD_Y + CARD_H - PAD - 64, cw - 2 * PAD, 64,
                          desc, font_size_pt=16, bold=False, color=GRAPHITE,
                          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)


# Constants for blank donor selection.
# Источник: brand/template-version.json → blank_donors (white/dark).
# Hardcoded fallback — на случай отсутствия файла.
# dark=51 (layout 'Контент / темный 1', bg #222222); было 22 (белый layout — баг).
_BLANK_DEFAULTS = {"white": 30, "dark": 51}
_blank_cfg = (_TPL_VER or {}).get("blank_donors", {}) if _TPL_VER else {}
BLANK_DONOR_WHITE = int(_blank_cfg.get("white", _BLANK_DEFAULTS["white"]))
BLANK_DONOR_DARK = int(_blank_cfg.get("dark", _BLANK_DEFAULTS["dark"]))


# Standalone test
if __name__ == "__main__":
    from pptx import Presentation
    from template_path import resolve_template

    p = Presentation(resolve_template())
    # Use slide 30 as clean blank donor — but FULLY clean shapes
    slide = list(p.slides)[BLANK_DONOR_WHITE - 1]
    clean_slide_to_blank(slide)

    # Render KPI
    render_kpi(slide, {
        "title": "ЧТО ПОЛУЧИЛОСЬ ЧЕРЕЗ 6 МЕСЯЦЕВ",
        "numbers": [
            {"value": "12", "desc": "сократили AI-инициатив", "pct": True},
            {"value": "84", "desc": "вовлечённость команды", "pct": True},
            {"value": "5",  "desc": "рабочих инструментов", "accent": True}
        ]
    }, dark=False)

    out = "pptx-skill/output/kpi_renderer_test.pptx"
    p.save(out)
    print(f"Saved {out}, used slide 30 as blank donor")
