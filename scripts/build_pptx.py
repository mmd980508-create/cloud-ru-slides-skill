#!/usr/bin/env python3
"""
build_pptx.py — собирает .pptx из плана агентов и шаблона Cloud.ru.

Usage:
    python3 build_pptx.py <plan.json> <template.pptx> <output.pptx>

plan.json формат:
    {
      "slides": [
        {
          "num": 1,
          "layout_idx": 0,
          "placeholder_assignments": [
            {"ph_idx": 0, "content": "Заголовок"}
          ],
          "shapes": []  // optional: from Infographic Maker
        }
      ]
    }

ВАЖНО: использует layout_idx из master шаблона. Все 102 layouts шаблона Cloud.ru
доступны через template.slide_layouts[i].
"""
import sys, json
from pptx import Presentation
from pptx.util import Emu


def build(plan_path, template_path, output_path):
    plan = json.load(open(plan_path, "r", encoding="utf-8"))
    p = Presentation(template_path)

    # Удалить существующие слайды шаблона (оставив только мастер с layouts)
    # ВАЖНО: python-pptx не имеет прямого API для удаления слайдов; используем XML hack
    xml_slides = p.slides._sldIdLst
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)

    # Создаём новые слайды из плана
    for sdata in plan["slides"]:
        layout_idx = sdata["layout_idx"]
        if layout_idx >= len(p.slide_layouts):
            print(f"WARN: layout_idx {layout_idx} >= {len(p.slide_layouts)}, fallback to 10", file=sys.stderr)
            layout_idx = 10
        layout = p.slide_layouts[layout_idx]
        slide = p.slides.add_slide(layout)

        # Заполняем placeholder'ы
        for assign in sdata.get("placeholder_assignments", []):
            ph_idx = assign["ph_idx"]
            content = assign["content"]
            ph = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == ph_idx:
                    ph = placeholder
                    break
            if ph is None:
                print(f"WARN: slide {sdata['num']} placeholder idx {ph_idx} not found in layout {layout_idx}", file=sys.stderr)
                continue
            if ph.has_text_frame:
                ph.text_frame.text = content
            else:
                print(f"WARN: slide {sdata['num']} ph {ph_idx} has no text_frame", file=sys.stderr)

        # TODO: Infographic Maker shapes — добавить в следующей итерации

    p.save(output_path)
    print(f"Saved {output_path}: {len(plan['slides'])} slides", file=sys.stderr)


def main():
    if len(sys.argv) != 4:
        print("Usage: build_pptx.py <plan.json> <template.pptx> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    build(sys.argv[1], sys.argv[2], sys.argv[3])


if __name__ == "__main__":
    main()
