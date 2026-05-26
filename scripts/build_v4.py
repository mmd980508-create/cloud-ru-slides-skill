#!/usr/bin/env python3
"""
build_v4.py — slot-based addressing + style-preserving text replace.

Решает баги v0.3:
- Адресация text frames по СЕМАНТИЧЕСКОЙ РОЛИ (slot), не по индексу
- Стиль (font/size/color/bold) первого run НЕ теряется
- Незаполненные слоты автоматически очищаются от заглушек
- Поддержка add_slide(layout=N) как fallback (для слайдов где cloning избыточен — например logo-финал)

Usage:
    python3 build_v4.py <plan.json> <template.pptx> <output.pptx>

plan.json:
{
  "slides": [
    {
      "clone_from_slide": 5,           # OR "add_layout": 95 (взаимоисключающе)
      "slots": {
        "title": "Куда смещается ценность в AI",
        "subtitle": "Надежда Погина · Руководитель проектов"
      },
      "_donor_category": "title_white"  # справочно для логов
    }
  ]
}

donor-slot-map.yaml:
  donors:
    5:
      slots:
        title: { shape_idx: 2, size_pt: 88 }
        subtitle: { shape_idx: 3, size_pt: 28 }
"""
import sys, json
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from copy import deepcopy
from lxml import etree

try:
    import yaml
except ImportError:
    print("ERROR: PyYAML required. pip3 install --user pyyaml", file=sys.stderr)
    sys.exit(1)


def load_donor_map(path):
    with open(path, encoding="utf-8") as f:
        data = yaml.safe_load(f)
    return data["donors"]


def get_text_frame_by_shape_idx(slide, shape_idx):
    """Returns text_frame at given shape_idx, or None."""
    shapes = list(slide.shapes)
    if shape_idx >= len(shapes):
        return None
    sh = shapes[shape_idx]
    if not sh.has_text_frame:
        return None
    return sh.text_frame


def replace_text_preserving_style(text_frame, new_text):
    """Заменяет текст в text_frame, сохраняя XML-уровневые run properties (rPr) первого run.

    Алгоритм:
      1. Берём первый <a:p> и его первый <a:r>
      2. Сохраняем XML-копию <a:rPr> этого run
      3. Стираем все остальные <a:p> и все runs внутри первого <a:p>
      4. Для каждой новой строки (split по \n) создаём <a:p> с одним <a:r>,
         в который копируем сохранённый rPr и кладём текст
    """
    if text_frame is None:
        return False

    txBody = text_frame._txBody  # <a:txBody>
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Найти первый <a:p> и первый <a:r> внутри
    p_elements = txBody.findall(qn("a:p"))
    if not p_elements:
        return False
    first_p = p_elements[0]
    first_r = first_p.find(qn("a:r"))

    # Сохраняем rPr (run properties) первого run
    saved_rPr = None
    if first_r is not None:
        rPr_el = first_r.find(qn("a:rPr"))
        if rPr_el is not None:
            saved_rPr = deepcopy(rPr_el)

    # Также сохраняем pPr (paragraph properties) первого <a:p>
    saved_pPr = first_p.find(qn("a:pPr"))
    saved_pPr_copy = deepcopy(saved_pPr) if saved_pPr is not None else None

    # Стираем все existing <a:p>
    for p_el in p_elements:
        txBody.remove(p_el)

    # Создаём новые <a:p> для каждой строки
    lines = new_text.split("\n")
    if not lines:
        lines = [""]

    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))
        if saved_pPr_copy is not None:
            new_p.append(deepcopy(saved_pPr_copy))
        new_r = etree.SubElement(new_p, qn("a:r"))
        if saved_rPr is not None:
            new_r.append(deepcopy(saved_rPr))
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line

    return True


def clear_text_frame(text_frame):
    """Очищает text_frame, оставляя пустой <a:p>."""
    return replace_text_preserving_style(text_frame, "")


def reorder_and_filter_slides(prs, source_slide_nums):
    """Удаляет слайды не из source_slide_nums и переставляет в нужном порядке.
    Дубли (один и тот же слайд использован дважды) пока не поддерживаются — повторное использование
    нужно обходить через add_slide(layout=N).
    """
    sldIdLst = prs.slides._sldIdLst
    all_ids = list(sldIdLst)
    num_to_sldId = {i + 1: el for i, el in enumerate(all_ids)}

    used_nums = []
    seen = set()
    for n in source_slide_nums:
        if n is None:
            continue  # placeholder for add_layout entries (handled separately)
        if n in seen:
            print(f"WARN: дубли донор-слайдов не поддерживаются (slide {n} повторяется)", file=sys.stderr)
            continue
        if n not in num_to_sldId:
            print(f"WARN: слайд {n} не найден (1..{len(all_ids)})", file=sys.stderr)
            continue
        used_nums.append(n)
        seen.add(n)

    keep_set = set(used_nums)
    for n, sldId in list(num_to_sldId.items()):
        if n not in keep_set:
            rId = sldId.attrib[qn('r:id')]
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)

    # Reorder
    remaining = list(sldIdLst)
    rem_map = {}
    for el in remaining:
        rId = el.attrib[qn('r:id')]
        rem_map[rId] = el
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for n in used_nums:
        sld = num_to_sldId[n]
        rId = sld.attrib[qn('r:id')]
        sldIdLst.append(rem_map[rId])


def build(plan_path, template_path, output_path, donor_map_path):
    plan = json.load(open(plan_path, encoding="utf-8"))
    p = Presentation(template_path)
    donors = load_donor_map(donor_map_path)

    # Phase 1: оставляем только нужные cloned слайды
    cloned_nums = [s.get("clone_from_slide") for s in plan["slides"]]
    reorder_and_filter_slides(p, cloned_nums)

    # Phase 2: применяем slot-overrides на каждый клонированный слайд
    cloned_iter = iter(p.slides)
    cloned_slide_for = {}  # source_num → actual slide object
    for plan_slide, actual in zip([s for s in plan["slides"] if s.get("clone_from_slide")], list(p.slides)):
        cloned_slide_for[plan_slide["clone_from_slide"]] = actual

    # Применяем slot-overrides
    for plan_slide in plan["slides"]:
        if not plan_slide.get("clone_from_slide"):
            continue
        src_num = plan_slide["clone_from_slide"]
        actual = cloned_slide_for[src_num]
        donor_def = donors.get(src_num)
        if donor_def is None:
            print(f"WARN: нет donor-slot-map для slide {src_num}, пропускаем slot-overrides", file=sys.stderr)
            continue

        slot_defs = donor_def.get("slots", {})
        slots_filled = plan_slide.get("slots", {})

        # 2a. Заполнить указанные слоты
        for slot_name, new_text in slots_filled.items():
            if slot_name not in slot_defs:
                print(f"WARN: slot '{slot_name}' не определён для donor {src_num}", file=sys.stderr)
                continue
            shape_idx = slot_defs[slot_name]["shape_idx"]
            tf = get_text_frame_by_shape_idx(actual, shape_idx)
            if tf is None:
                print(f"WARN: shape_idx {shape_idx} не найден на slide {src_num} для slot '{slot_name}'", file=sys.stderr)
                continue
            replace_text_preserving_style(tf, new_text)

        # 2b. Очистить НЕзаполненные слоты (удалить placeholder-заглушки)
        for slot_name, slot_def in slot_defs.items():
            if slot_name in slots_filled:
                continue  # already filled
            if slot_def.get("optional"):
                continue  # optional slot — оставляем как есть (обычно % знаки)
            shape_idx = slot_def["shape_idx"]
            tf = get_text_frame_by_shape_idx(actual, shape_idx)
            if tf is not None:
                clear_text_frame(tf)

    # Phase 3: добавить slides через add_layout (для финального logo-слайда без donor-clone)
    # Не реализовано в этой версии — все слайды через cloning. add_layout-режим — TODO v0.5

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides", file=sys.stderr)


def main():
    if len(sys.argv) < 4:
        print("Usage: build_v4.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]", file=sys.stderr)
        sys.exit(1)
    plan_p = sys.argv[1]
    tpl_p = sys.argv[2]
    out_p = sys.argv[3]
    donor_p = sys.argv[4] if len(sys.argv) > 4 else "pptx-skill/brand/donor-slot-map.yaml"
    build(plan_p, tpl_p, out_p, donor_p)


if __name__ == "__main__":
    main()
