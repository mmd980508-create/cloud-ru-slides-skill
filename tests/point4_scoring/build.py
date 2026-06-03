#!/usr/bin/env python3
"""
Point 4 test — визуальный delta-scorecard взвешенного скоринга.
Считает brand_guardian.score_delta(before, after) по point2 before/after и рисует
Cloud.ru-слайд: ИТОГО было→стало + разбивка по категориям с дельтами.
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
SCRIPTS = os.path.join(HERE, "..", "..", "scripts")
sys.path.insert(0, SCRIPTS)

from design_tokens import load_tokens          # noqa: E402
import brand_guardian                           # noqa: E402
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu, Pt                    # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE           # noqa: E402
from pptx.oxml.ns import qn                      # noqa: E402

T = load_tokens()
EMU = T.EMU
LOGO = os.path.join(HERE, "..", "point1_tokens", "logo.png")


def emu(px):
    return Emu(int(round(px * EMU)))


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)
    sp.fill.solid(); sp.fill.fore_color.rgb = T.rgb(color)
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, s, size, color, semibold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, caps=False):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s.upper() if caps else s
    r.font.name = T.semibold_face if semibold else T.family
    r.font.size = Pt(size); r.font.bold = False
    r.font.color.rgb = T.rgb(color)
    return tb


def delta_chip(slide, x, y, value):
    """Зелёный чип с +N (или серый с 0)."""
    pos = value > 0
    w, h = 56, 28
    rect(slide, x, y, w, h, "Green" if pos else "Gray")
    sp = slide.shapes[-1]; tf = sp.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ("+%g" % value) if pos else "0"
    r.font.name = T.family; r.font.size = Pt(13); r.font.bold = False
    r.font.color.rgb = T.rgb("Black")


def add_logo(slide):
    lg = T.place("logo")["content"]
    if os.path.isfile(LOGO):
        slide.shapes.add_picture(LOGO, emu(lg["x"]), emu(lg["y"]), emu(lg["w"]), emu(lg["h"]))


def add_copyright(slide):
    cp = T.place("copyright")
    tb = slide.shapes.add_textbox(emu(cp["x"]), emu(cp["y"]), emu(cp["w"]), emu(cp["h"]))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = cp["text"]
    r.font.name = T.family; r.font.size = Pt(cp["size"]); r.font.bold = False
    r.font.color.rgb = T.rgb(cp["color"])


def build():
    d = brand_guardian.score_delta(
        os.path.join(HERE, "..", "point2_normalizer", "before.pptx"),
        os.path.join(HERE, "..", "point2_normalizer", "after.pptx"))

    prs = Presentation(); prs.slide_width = emu(1280); prs.slide_height = emu(720)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 1280, 720, "White")
    add_logo(s); add_copyright(s)
    safe = T.safe

    txt(s, 35, 38, safe["right"] - 35 - 170, 40,
        "Взвешенный скоринг | было → стало", 20, "Black", semibold=True, caps=True)
    txt(s, 35, 76, safe["right"] - 35, 24,
        "brand_guardian: 0–100 по категориям с весами + дельта для verify-петли",
        12, "text_gray")

    # ИТОГО hero
    txt(s, 35, 150, 200, 24, "ИТОГО", 14, "text_gray", semibold=True)
    txt(s, 35, 175, 520, 90,
        "%g  →  %g" % (d["before"], d["after"]), 64, "Black", semibold=True)
    delta_chip(s, 470, 200, d["delta"])

    # Категории — серая карточка
    cy, ch = 300, 320
    rect(s, safe["left"], cy, safe["right"] - safe["left"], ch, "Gray")
    pad = 28
    cats = list(d["categories"].items())
    row_h = (ch - 2 * pad) / len(cats)
    for i, (cat, v) in enumerate(cats):
        ry = cy + pad + i * row_h
        wt = brand_guardian.CATEGORY_WEIGHTS[cat]
        txt(s, safe["left"] + pad, ry, 300, 28, f"{cat}", 18, "Black", semibold=True)
        txt(s, safe["left"] + pad + 320, ry + 2, 260, 26,
            "%g → %g  / %d" % (v["before"], v["after"], wt), 16, "Black")
        delta_chip(s, safe["left"] + pad + 620, ry, v["delta"])

    out = os.path.join(HERE, "point4_scoring.pptx")
    prs.save(out); print("saved:", out)
    return out


if __name__ == "__main__":
    build()
