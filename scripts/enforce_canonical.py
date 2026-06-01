#!/usr/bin/env python3
"""
enforce_canonical.py — пост-обработка слайдов: приводит ЛЮБОЙ слайд (включая
clone-based из донора) к canonical-правилам Cloud.ru, которые раньше работали
только в native-рендерерах.

Зачем: общие презентации собираются через `clone_from_slide` (клон донора +
подстановка текста). Этот путь сохранял форматирование донора и LLM-оверрайды,
поэтому фиксы #2 (контраст), #3 (SemiBold), #5 (размер) до него не доходили —
оставались зелёный текст на белом, белый на зелёном, bold-флаги, мелкий шрифт.

Что делает `enforce_canonical_slide(slide, dark=False)`:
  - Цвет текста: зелёный → #222222; белый → #222222 (КРОМЕ тёмного контекста —
    тёмная заливка фигуры или тёмный фон слайда, там белый остаётся).
  - Жирность: bold-флаг (b=1) → начертание «SB Sans Display Semibold», b=0.
  - Размер: < min_pt (по умолч. 12) → поднять до min_pt. Опц. bump_to: поднять
    «мелкий body» (>= bump_from) до bump_to (например 16) — НЕ трогая крупные.

Обходит фигуры рекурсивно: группы (grpSp) и таблицы тоже.

CLI:  python3 enforce_canonical.py in.pptx out.pptx [--min 12] [--bump-to 16]
"""
import sys
import copy
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree


EMU = 9525
FONT_SEMIBOLD = "SB Sans Display Semibold"
GRAPHITE_HEX = "222222"
# Канонический header контент-слайда (как в шаблоне): (35,38), 20pt SemiBold CAPS.
HEADER_X, HEADER_Y = 35, 38
HEADER_PT = 20
# Bump до 16pt применяем только к боксам шире этого (px) — body-блоки. Узкие
# ноды диаграмм/мелкие лейблы оставляем как есть (иначе текст вылезает).
BUMP_MIN_WIDTH_PX = 300
# Слайды, у которых header НЕ выравниваем (титулы/разделы — там заголовок крупный/по центру).
NON_CONTENT_LAYOUT_HINTS = ("титул", "раздел", "divider", "логотип", "фотофон", "фото фон", "обложк")

# Все «зелёные» оттенки Cloud.ru (текст ими быть не должен — только элементы).
GREEN_HEXES = {"26D07C", "00D97B", "1AB066", "1ABF6F", "22C993", "2DD27D"}
WHITE_HEXES = {"FFFFFF", "FEFEFE"}
# scheme-цвета, считающиеся «светлыми» (= белый текст) и «зелёными».
WHITE_SCHEMES = {"bg1", "lt1", "bg2", "lt2", "light1", "light2"}
GREEN_SCHEMES = {"accent1"}
DARK_SCHEMES = {"tx1", "dk1", "tx2", "dk2", "dark1", "dark2"}


def _lum(hex6):
    try:
        r = int(hex6[0:2], 16); g = int(hex6[2:4], 16); b = int(hex6[4:6], 16)
        return 0.299 * r + 0.587 * g + 0.114 * b
    except Exception:
        return None


def _fill_info(shape):
    """(is_dark, is_green) по заливке фигуры/ячейки. Без заливки → (False, False)."""
    try:
        fill = shape.fill
        if fill is None or fill.type is None:
            return (False, False)
        # solid?
        try:
            rgb = fill.fore_color.rgb
            if rgb is not None:
                hx = str(rgb).upper()
                l = _lum(hx)
                return ((l is not None and l < 128), hx in GREEN_HEXES)
        except Exception:
            pass
        # scheme fill
        spPr = shape._element.find(qn("p:spPr")) if shape._element.find(qn("p:spPr")) is not None else None
        if spPr is not None:
            sc = spPr.find(qn("a:solidFill") + "/" + qn("a:schemeClr"))
            if sc is not None:
                v = sc.get("val")
                return (v in DARK_SCHEMES, v in GREEN_SCHEMES)
    except Exception:
        pass
    return (False, False)


def _run_color(rPr):
    """('green'|'white'|'dark'|'other'|None) по цвету run'а."""
    if rPr is None:
        return None
    sf = rPr.find(qn("a:solidFill"))
    if sf is None:
        return None
    srgb = sf.find(qn("a:srgbClr"))
    if srgb is not None:
        hx = (srgb.get("val") or "").upper()
        if hx in GREEN_HEXES:
            return "green"
        if hx in WHITE_HEXES:
            return "white"
        l = _lum(hx)
        return "dark" if (l is not None and l < 128) else "other"
    sc = sf.find(qn("a:schemeClr"))
    if sc is not None:
        v = sc.get("val")
        if v in GREEN_SCHEMES:
            return "green"
        if v in WHITE_SCHEMES:
            return "white"
        if v in DARK_SCHEMES:
            return "dark"
    return "other"


def _set_run_color(rPr, hex6):
    for tag in ("solidFill", "noFill", "gradFill", "blipFill", "pattFill", "grpFill"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is not None:
            rPr.remove(el)
    # solidFill должен идти ПОСЛЕ ln, но ДО latin/ea/cs по схеме CT_TextCharacterProperties.
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    srgb = etree.SubElement(sf, qn("a:srgbClr"))
    srgb.set("val", hex6)
    # порядок: переставим solidFill перед latin/ea/cs/sym/hlink, если они есть
    _reorder_rpr_children(rPr)


_RPR_ORDER = ["a:ln", "a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
              "a:pattFill", "a:grpFill", "a:effectLst", "a:effectDag",
              "a:highlight", "a:uLnTx", "a:uLn", "a:uFillTx", "a:uFill",
              "a:latin", "a:ea", "a:cs", "a:sym", "a:hlinkClick", "a:hlinkMouseOver",
              "a:rtl", "a:extLst"]


def _reorder_rpr_children(rPr):
    order = {qn(t): i for i, t in enumerate(_RPR_ORDER)}
    children = list(rPr)
    children.sort(key=lambda el: order.get(el.tag, 999))
    for ch in children:
        rPr.remove(ch)
    for ch in children:
        rPr.append(ch)


def _set_semibold(rPr):
    rPr.set("b", "0")
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is not None:
            el.set("typeface", FONT_SEMIBOLD)
    if rPr.find(qn("a:latin")) is None:
        lat = etree.SubElement(rPr, qn("a:latin"))
        lat.set("typeface", FONT_SEMIBOLD)
        _reorder_rpr_children(rPr)


def _iter_text_shapes(shapes):
    """Рекурсивно: (shape, [runs]) для всех текстовых фигур, групп и ячеек таблиц."""
    for sh in shapes:
        st = sh.shape_type
        if st == 6:  # GROUP
            try:
                yield from _iter_text_shapes(sh.shapes)
            except Exception:
                pass
            continue
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    runs = [r for p in cell.text_frame.paragraphs for r in p.runs]
                    if runs:
                        yield (cell, runs)
            continue
        if sh.has_text_frame:
            runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
            if runs:
                yield (sh, runs)


def _shape_max_font_pt(shape):
    mx = 0
    try:
        for para in shape.text_frame.paragraphs:
            for r in para.runs:
                rPr = r._r.find(qn("a:rPr"))
                if rPr is not None and rPr.get("sz"):
                    mx = max(mx, int(rPr.get("sz")) / 100.0)
    except Exception:
        pass
    return mx


def _is_content_layout(slide):
    lay = (slide.slide_layout.name or "").lower()
    if any(h in lay for h in NON_CONTENT_LAYOUT_HINTS):
        return False
    return ("контент" in lay) or ("content" in lay)


def _find_header_textbox(slide):
    """СТРОГО ищет текстбокс-заголовок: широкий горизонтальный, у верха, 16-28pt,
    короткий текст, не повёрнут. Узкие/вертикальные сайдбары, огромный body и
    логотип справа — исключаются (чтобы не повторить поломку слайда 6)."""
    best = None
    for sh in slide.shapes:
        try:
            if sh.is_placeholder:
                continue
            if not sh.has_text_frame:
                continue
            if sh.top is None or sh.left is None or sh.width is None:
                continue
            top = sh.top / EMU
            left = sh.left / EMU
            w = sh.width / EMU
            rot = getattr(sh, "rotation", 0) or 0
            if top >= 95 or left >= 200 or w < 400 or abs(rot) > 1:
                continue
            mf = _shape_max_font_pt(sh)
            if mf < 16 or mf > 28:
                continue
            txt = (sh.text_frame.text or "").strip()
            if not txt or len(txt) > 80 or txt.count("\n") > 2:
                continue
            if best is None or top < best[1]:
                best = (sh, top)
        except Exception:
            continue
    return best[0] if best else None


def _ensure_title_placeholder(slide):
    """Возвращает title-placeholder слайда. Если на слайде его НЕТ (генерация
    создала слайд без placeholder'ов) — инстанцирует из layout, КОПИРУЯ ПОЛНУЮ
    геометрию (off+ext, напр. 35,38/963×54). Полный xfrm = placeholder не
    схлопнется (схлопывание было от частичного xfrm — off без ext)."""
    try:
        t = slide.shapes.title
        if t is not None:
            return t
    except Exception:
        pass
    lay_ph = None
    for ph in slide.slide_layout.placeholders:
        try:
            if ph.placeholder_format.idx == 0 or "TITLE" in str(ph.placeholder_format.type):
                lay_ph = ph
                break
        except Exception:
            continue
    if lay_ph is None:
        return None
    sp = copy.deepcopy(lay_ph._element)
    # очистить текст-промпт из layout-плейсхолдера
    txBody = sp.find(qn("p:txBody"))
    if txBody is not None:
        for p_el in txBody.findall(qn("a:p")):
            for r in p_el.findall(qn("a:r")):
                p_el.remove(r)
    slide.shapes._spTree.append(sp)
    try:
        return slide.shapes.title
    except Exception:
        return None


def normalize_header_to_placeholder(slide, dark=False):
    """ВАРИАНТ A (исправлено 2026-05-29): заголовок контент-слайда → в штатный
    TITLE-placeholder. БЕЗОПАСНО:

    1) Заголовок УЖЕ в placeholder → НЕ ТРОГАЕМ (любое касание placeholder'а в
       PowerPoint схлопывает его в столбик; LibreOffice это не показывает).
    2) Заголовок воткнут ТЕКСТБОКСОМ → переносим текст в title-placeholder
       (инстанцируя его из layout с ПОЛНОЙ геометрией, если на слайде его нет),
       удаляем текстбокс. Текст вписываем без явной смены позиции/размера —
       геометрия из layout (единая). Цвет/вес правит общий пасс.
    Титульные/divider/вертикальные/огромные заголовки — не трогаем.
    """
    if not _is_content_layout(slide):
        return False
    try:
        title_ph = slide.shapes.title
    except Exception:
        title_ph = None
    # 1) Заголовок уже в placeholder → не трогаем.
    if title_ph is not None and title_ph.has_text_frame and (title_ph.text or "").strip():
        return False
    # 2) Заголовок в текстбоксе → перенести в placeholder.
    header_tb = _find_header_textbox(slide)
    if header_tb is None:
        return False
    header_text = (header_tb.text_frame.text or "").strip()
    ph = title_ph if title_ph is not None else _ensure_title_placeholder(slide)
    if ph is None:
        # placeholder создать не удалось — fallback: выровнять текстбокс по позиции
        header_tb.left = Emu(HEADER_X * EMU)
        header_tb.top = Emu(HEADER_Y * EMU)
        return True
    # вписать текст в placeholder (без явной геометрии — наследуется из layout)
    ph.text_frame.text = header_text
    # удалить текстбокс-заголовок
    try:
        header_tb._element.getparent().remove(header_tb._element)
    except Exception:
        pass
    return True


def enforce_canonical_slide(slide, dark=False, min_pt=12, bump_from=None, bump_to=None,
                            normalize_header=False):
    """Приводит слайд к canonical. Возвращает счётчики изменений."""
    stats = {"green_text": 0, "white_on_light": 0, "bold": 0, "size_min": 0,
             "size_bump": 0, "header_norm": 0}
    if normalize_header:
        if normalize_header_to_placeholder(slide, dark=dark):
            stats["header_norm"] = 1
    for shape, runs in _iter_text_shapes(slide.shapes):
        fill_dark, fill_green = _fill_info(shape)
        dark_ctx = fill_dark or dark
        # Bump до 16pt — только в ШИРОКИХ боксах (body-блоки, где есть место).
        # Узкие диаграмм-ноды/мелкие лейблы не бампим, чтобы текст не вылезал.
        shape_w = getattr(shape, "width", None)
        wide = (shape_w is not None) and (shape_w / EMU >= BUMP_MIN_WIDTH_PX)
        for run in runs:
            rPr = run._r.find(qn("a:rPr"))
            if rPr is None:
                rPr = run._r.makeelement(qn("a:rPr"), {})
                run._r.insert(0, rPr)
            # --- цвет ---
            col = _run_color(rPr)
            if col == "green":
                _set_run_color(rPr, GRAPHITE_HEX); stats["green_text"] += 1
            elif col == "white" and not dark_ctx:
                _set_run_color(rPr, GRAPHITE_HEX); stats["white_on_light"] += 1
            # --- жирность ---
            if rPr.get("b") == "1":
                _set_semibold(rPr); stats["bold"] += 1
            # --- размер ---
            sz = rPr.get("sz")
            if sz is not None:
                pt = int(sz) / 100.0
                if pt < min_pt:
                    rPr.set("sz", str(int(min_pt * 100))); stats["size_min"] += 1
                elif wide and bump_to and bump_from and bump_from <= pt < bump_to:
                    rPr.set("sz", str(int(bump_to * 100))); stats["size_bump"] += 1
    return stats


def slide_is_dark(slide):
    """Эвристика тёмного слайда: фон слайда/лейаута = tx1/dk* (scheme) или тёмный srgb."""
    for el in (slide._element, slide.slide_layout._element):
        bg = el.find(".//" + qn("p:bg"))
        if bg is not None:
            sc = bg.find(".//" + qn("a:schemeClr"))
            sr = bg.find(".//" + qn("a:srgbClr"))
            if sc is not None and sc.get("val") in DARK_SCHEMES:
                return True
            if sr is not None and (_lum((sr.get("val") or "").upper()) or 255) < 128:
                return True
    return False


def enforce_canonical_pptx(in_path, out_path, min_pt=12, bump_from=None, bump_to=None,
                           dark_slides=None, normalize_header=True):
    """Применяет enforce к каждому слайду .pptx. dark_slides — set индексов (1-based)
    тёмных слайдов (если None — определяем эвристикой по фону)."""
    prs = Presentation(in_path)
    total = {"green_text": 0, "white_on_light": 0, "bold": 0, "size_min": 0,
             "size_bump": 0, "header_norm": 0}
    for i, slide in enumerate(prs.slides, 1):
        dark = (i in dark_slides) if dark_slides is not None else slide_is_dark(slide)
        st = enforce_canonical_slide(slide, dark=dark, min_pt=min_pt,
                                     bump_from=bump_from, bump_to=bump_to,
                                     normalize_header=normalize_header)
        for k in total:
            total[k] += st[k]
    prs.save(out_path)
    return total


def main():
    if len(sys.argv) < 3:
        print("Usage: enforce_canonical.py in.pptx out.pptx [--min 12] [--bump-from 11 --bump-to 16]",
              file=sys.stderr)
        sys.exit(1)
    args = sys.argv[3:]
    min_pt = 12; bump_from = None; bump_to = None
    if "--min" in args:
        min_pt = float(args[args.index("--min") + 1])
    if "--bump-to" in args:
        bump_to = float(args[args.index("--bump-to") + 1])
        bump_from = float(args[args.index("--bump-from") + 1]) if "--bump-from" in args else min_pt
    # Опасные правки выключены по умолчанию: --normalize-header чтобы включить.
    normalize_header = "--normalize-header" in args
    total = enforce_canonical_pptx(sys.argv[1], sys.argv[2], min_pt=min_pt,
                                   bump_from=bump_from, bump_to=bump_to,
                                   normalize_header=normalize_header)
    print(f"enforce_canonical: {total}", file=sys.stderr)
    print(f"Saved {sys.argv[2]}")


if __name__ == "__main__":
    main()
