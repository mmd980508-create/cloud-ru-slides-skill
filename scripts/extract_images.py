#!/usr/bin/env python3
"""
extract_images.py — извлекает все изображения из .pptx в файлы.

Usage:
    python3 extract_images.py <input.pptx> <output_dir> [manifest.json]

Output:
- output_dir/slide{N}_img{M}.{ext} — файлы изображений
- manifest.json (optional) — JSON с метаданными:
    {
      "source": "...",
      "images": [
        {
          "slide_num": 2,
          "image_num": 1,
          "file": "slide2_img1.png",
          "name": "Picture 5",
          "left_emu": ..., "top_emu": ..., "width_emu": ..., "height_emu": ...,
          "left_px": ..., "top_px": ..., "width_px": ..., "height_px": ...,
          "ext": "png", "size_bytes": ...
        }
      ]
    }
"""
import sys
import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_PX = 9525  # at 96 DPI


def extract(input_path, output_dir, manifest_path=None):
    p = Presentation(input_path)
    os.makedirs(output_dir, exist_ok=True)

    manifest = {
        "source": input_path,
        "slide_count": len(p.slides),
        "images": [],
    }

    for snum, slide in enumerate(p.slides, start=1):
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            img_idx += 1

            # Получить blob через shape.image
            try:
                image = shape.image
                blob = image.blob
                ext = image.ext  # e.g. "png", "jpeg"
            except Exception as e:
                print(f"WARN: slide {snum} pic {img_idx}: {e}", file=sys.stderr)
                continue

            filename = f"slide{snum}_img{img_idx}.{ext}"
            file_path = os.path.join(output_dir, filename)
            with open(file_path, "wb") as f:
                f.write(blob)

            manifest["images"].append({
                "slide_num": snum,
                "image_num": img_idx,
                "file": filename,
                "name": shape.name,
                "left_emu": shape.left, "top_emu": shape.top,
                "width_emu": shape.width, "height_emu": shape.height,
                "left_px": int(shape.left / EMU_PER_PX) if shape.left else None,
                "top_px": int(shape.top / EMU_PER_PX) if shape.top else None,
                "width_px": int(shape.width / EMU_PER_PX) if shape.width else None,
                "height_px": int(shape.height / EMU_PER_PX) if shape.height else None,
                "ext": ext,
                "size_bytes": len(blob),
            })

    if manifest_path:
        with open(manifest_path, "w", encoding="utf-8") as f:
            json.dump(manifest, f, ensure_ascii=False, indent=2)

    return manifest


def main():
    if len(sys.argv) < 3:
        print("Usage: extract_images.py <input.pptx> <output_dir> [manifest.json]", file=sys.stderr)
        sys.exit(1)
    input_path = sys.argv[1]
    output_dir = sys.argv[2]
    manifest_path = sys.argv[3] if len(sys.argv) > 3 else None

    m = extract(input_path, output_dir, manifest_path)
    print(f"Extracted {len(m['images'])} images from {m['slide_count']} slides → {output_dir}",
          file=sys.stderr)
    if manifest_path:
        print(f"Manifest → {manifest_path}", file=sys.stderr)


if __name__ == "__main__":
    main()
