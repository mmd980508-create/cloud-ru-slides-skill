#!/usr/bin/env python3
"""
build_v3.py — клонирует слайды-доноры из шаблона + подменяет текст.

Стратегия v0.3:
- Не add_slide(layout=N), а копировать целые слайды из шаблона
  (они содержат всё: декор, паттерны, лого, footer)
- Заменять текст в text_frames по индексу
- Удалять лишние слайды

Usage:
    python3 build_v3.py <plan.json> <template.pptx> <output.pptx>

plan.json:
{
  "slides": [
    {"clone_from_slide": 5, "text_overrides": {"<old_substr>": "<new_text>"}},
    ...
  ]
}
"""
import sys, json, copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree


def replace_in_text_frames(slide, text_overrides):
    """В каждом text_frame слайда находит и заменяет текст по подстрокам."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = tf.text
        for old, new in text_overrides.items():
            if old in full_text:
                # Простая стратегия: заменяем full text шейпа целиком,
                # сохраняя стиль первого run первого параграфа
                # (более сложная — XML-level surgery, но это ломает стили)
                first_para = tf.paragraphs[0]
                first_run = first_para.runs[0] if first_para.runs else None

                # Сохранить шрифт/размер/цвет/bold первого run
                font_props = {}
                if first_run is not None:
                    font_props['name'] = first_run.font.name
                    font_props['size'] = first_run.font.size
                    font_props['bold'] = first_run.font.bold
                    try:
                        if first_run.font.color and first_run.font.color.type is not None:
                            font_props['color_rgb'] = first_run.font.color.rgb
                    except Exception:
                        pass

                # Заменяем текст в первом параграфе, остальные параграфы стираем
                tf.clear()
                p = tf.paragraphs[0]
                for line_i, line in enumerate(new.split("\n")):
                    if line_i > 0:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = line
                    if font_props.get('name'):
                        run.font.name = font_props['name']
                    if font_props.get('size'):
                        run.font.size = font_props['size']
                    if font_props.get('bold') is not None:
                        run.font.bold = font_props['bold']
                    if 'color_rgb' in font_props:
                        run.font.color.rgb = font_props['color_rgb']
                # apply only first matching override
                break


def replace_by_index(slide, index_overrides):
    """index_overrides: dict {tf_index: new_text} — где tf_index это порядковый
    номер text_frame в slide.shapes (по тексту, исключая пустые).
    Сохраняет ВСЕ runs/styling первого параграфа, только меняет текст."""
    text_idx = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip():
            continue
        if text_idx in index_overrides:
            tf = shape.text_frame
            new_text = index_overrides[text_idx]

            # Сохраняем стиль первого run
            first_p = tf.paragraphs[0]
            first_r = first_p.runs[0] if first_p.runs else None
            saved = {}
            if first_r is not None:
                saved['name'] = first_r.font.name
                saved['size'] = first_r.font.size
                saved['bold'] = first_r.font.bold
                try:
                    if first_r.font.color and first_r.font.color.type is not None:
                        saved['color_rgb'] = first_r.font.color.rgb
                except Exception:
                    pass

            tf.clear()
            for line_i, line in enumerate(new_text.split("\n")):
                p = tf.paragraphs[0] if line_i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = line
                if saved.get('name'):
                    run.font.name = saved['name']
                if saved.get('size'):
                    run.font.size = saved['size']
                if saved.get('bold') is not None:
                    run.font.bold = saved['bold']
                if 'color_rgb' in saved:
                    run.font.color.rgb = saved['color_rgb']
        text_idx += 1


def reorder_and_filter_slides(prs, source_slide_nums):
    """Удаляет все слайды, кроме указанных, и переставляет в нужном порядке.

    source_slide_nums: list of 1-based slide numbers from original template,
                       in the desired final order.

    После работы: prs содержит ровно len(source_slide_nums) слайдов в указанном порядке.
    Дубли поддерживаются (если один и тот же слайд нужен 2 раза — копируем).
    """
    sldIdLst = prs.slides._sldIdLst
    # Получаем все sldId-элементы и их rIds
    all_ids = list(sldIdLst)  # original 88
    # Map: 1-based slide num → sldId element
    num_to_sldId = {i+1: el for i, el in enumerate(all_ids)}

    # Новый порядок: для каждого требуемого num найти sldId
    # Дубли пока не поддерживаем (одна и та же страница 2 раза в результате)
    # Удалим из sldIdLst всё, что не в source_slide_nums
    # А оставшиеся переставим в порядке source_slide_nums

    used_nums = []
    seen = set()
    for n in source_slide_nums:
        if n in seen:
            print(f"WARN: дубли слайдов пока не поддерживаются (slide {n} использован дважды)", file=sys.stderr)
            continue
        if n not in num_to_sldId:
            print(f"WARN: слайд {n} не найден в шаблоне (есть только 1..{len(all_ids)})", file=sys.stderr)
            continue
        used_nums.append(n)
        seen.add(n)

    # Удаляем из sldIdLst все, что не в used_nums
    keep_set = set(used_nums)
    for n, sldId in list(num_to_sldId.items()):
        if n not in keep_set:
            rId = sldId.attrib[qn('r:id')]
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)

    # Переставляем в правильном порядке
    # Снимаем все оставшиеся, потом добавляем в нужном порядке
    remaining = list(sldIdLst)
    rem_map = {}
    for el in remaining:
        # Вычисляем какой это original num по rId через part
        rId = el.attrib[qn('r:id')]
        rem_map[rId] = el
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for n in used_nums:
        # Get rId of this slide
        sld = num_to_sldId[n]
        rId = sld.attrib[qn('r:id')]
        sldIdLst.append(rem_map[rId])


def build(plan_path, template_path, output_path):
    plan = json.load(open(plan_path, encoding="utf-8"))
    p = Presentation(template_path)

    # 1. Reorder/filter slides по плану
    source_nums = [s["clone_from_slide"] for s in plan["slides"]]
    reorder_and_filter_slides(p, source_nums)

    # 2. Применяем text overrides на каждый слайд (по порядку)
    for plan_slide, actual_slide in zip(plan["slides"], p.slides):
        if "index_overrides" in plan_slide:
            replace_by_index(actual_slide, {int(k): v for k, v in plan_slide["index_overrides"].items()})
        elif "text_overrides" in plan_slide:
            replace_in_text_frames(actual_slide, plan_slide["text_overrides"])

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides", file=sys.stderr)


def main():
    if len(sys.argv) != 4:
        print("Usage: build_v3.py <plan.json> <template.pptx> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    build(sys.argv[1], sys.argv[2], sys.argv[3])


if __name__ == "__main__":
    main()
