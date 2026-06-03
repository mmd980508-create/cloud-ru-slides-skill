#!/usr/bin/env python3
"""
Point 6 — реальный KPI-слайд: трекинг крупных цифр 0 (было) vs -400 (стало,
брендбук §6). Две страницы одинакового KPI через kpi_renderer.
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "..", "scripts"))
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu                        # noqa: E402
import kpi_renderer                              # noqa: E402

EMU = 9525
CFG = {"title": "Что получилось", "numbers": [
    {"value": "12", "desc": "AI-инициатив", "accent": False, "pct": False},
    {"value": "84", "desc": "вовлечённость", "accent": False, "pct": True},
    {"value": "30", "desc": "инструментов", "accent": True, "pct": False},
]}


def build():
    prs = Presentation()
    prs.slide_width = Emu(1280 * EMU); prs.slide_height = Emu(720 * EMU)

    # Слайд 1 — БЫЛО (spc 0)
    kpi_renderer.BIG_NUMBER_SPC = 0
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    kpi_renderer.render_kpi(s1, dict(CFG, title="БЫЛО · трекинг 0"))

    # Слайд 2 — СТАЛО (spc -400, брендбук §6)
    kpi_renderer.BIG_NUMBER_SPC = -400
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    kpi_renderer.render_kpi(s2, dict(CFG, title="СТАЛО · трекинг -400 (брендбук §6)"))

    out = os.path.join(HERE, "kpi_demo.pptx")
    prs.save(out); print("saved:", out)
    return out


if __name__ == "__main__":
    build()
