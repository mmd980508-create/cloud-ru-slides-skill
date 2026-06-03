#!/usr/bin/env python3
"""
Point 5 test — намеренно «слоповый» слайд (AI-tells), чтобы показать, что
brand_guardian.check_slop их ловит: card-in-card, радужные акценты, мутная
иерархия (много кеглей). Печатает находки линтера.
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
SCRIPTS = os.path.join(HERE, "..", "..", "scripts")
sys.path.insert(0, SCRIPTS)

from design_tokens import load_tokens          # noqa: E402
import brand_guardian                           # noqa: E402
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu, Pt                    # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE           # noqa: E402
from pptx.oxml.ns import qn                      # noqa: E402

T = load_tokens()
EMU = T.EMU


def emu(px):
    return Emu(int(round(px * EMU)))


def rect(slide, x, y, w, h, hexcolor):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)
    from pptx.dml.color import RGBColor
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(hexcolor.lstrip("#"))
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, s, size, hexcolor="222222"):
    from pptx.dml.color import RGBColor
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = s
    r.font.name = T.family; r.font.size = Pt(size); r.font.bold = False
    r.font.color.rgb = RGBColor.from_string(hexcolor)


def build():
    prs = Presentation(); prs.slide_width = emu(1280); prs.slide_height = emu(720)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 1280, 720, "FFFFFF")

    txt(s, 35, 30, 1100, 40, "ANTI-SLOP ДЕМО · намеренно плохой слайд", 24)

    # card-in-card: серая карточка внутри серой карточки
    rect(s, 35, 110, 600, 360, "F2F2F2")
    rect(s, 75, 160, 520, 180, "D9D9D9")   # вложенная карточка
    txt(s, 95, 180, 480, 40, "Вложенная карточка (card-in-card)", 18)

    # rainbow: 4 акцентных цвета в ряд
    rect(s, 700, 160, 110, 110, "26D07C")  # green
    rect(s, 820, 160, 110, 110, "CFF500")  # yellow
    rect(s, 940, 160, 110, 110, "A068FF")  # purple
    rect(s, 1060, 160, 110, 110, "C0E0FC") # blue
    txt(s, 700, 285, 470, 30, "Радужные акценты: green+yellow+purple+blue", 14)

    # muddy hierarchy: 7 разных кеглей
    for i, sz in enumerate((12, 16, 20, 28, 36, 44, 56)):
        txt(s, 700, 330 + i * 48, 470, 50, f"кегль {sz}pt", sz)

    out = os.path.join(HERE, "point5_slop.pptx")
    prs.save(out)

    # прогон линтера
    rep = brand_guardian.validate_slide(prs.slides[0], 1)
    slop = [w for w in rep["warnings"] if w["type"] in
            ("card_in_card", "rainbow_accents", "gradient_text", "muddy_hierarchy")]
    print("saved:", out)
    print("ANTI-SLOP находки линтера:")
    for w in slop:
        print(f"  ⚠️  {w['type']}: {w['msg']}")
    return out


if __name__ == "__main__":
    build()
