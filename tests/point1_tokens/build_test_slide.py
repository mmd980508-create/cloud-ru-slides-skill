#!/usr/bin/env python3
"""
Point 1 test slide — строит слайд, беря КАЖДУЮ константу из design_tokens loader.
Если рендер на-бренд → token-контракт работает (цвета/типографика/сетка/геометрия/
размещение заголовка читаются из одного источника, а не из хардкодов).

Запуск:
    python3 build_test_slide.py            # → point1_tokens.pptx
"""
import os
import sys
import math

# импорт loader из ../../scripts
HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "..", "scripts"))
from design_tokens import load_tokens  # noqa: E402

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

T = load_tokens()
EMU = T.EMU


def emu(px):
    return Emu(int(round(px * EMU)))


def _no_line(shape):
    shape.line.fill.background()


def _fill(shape, color_name):
    shape.fill.solid()
    shape.fill.fore_color.rgb = T.rgb(color_name)
    _no_line(shape)


def rect(slide, x, y, w, h, color_name):
    """Прямоугольник с заливкой из токена. geometry.rounded.none=0 → острые углы."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    # Снять theme-<p:style> дефолтного шаблона python-pptx (тащит тень/контур/градиент).
    # geometry: плоско, без теней (BR §2). Прод-рендеры рисуют на бренд-канвасе и
    # этой проблемы не имеют — здесь чистим явно, т.к. тест-дек на default-шаблоне.
    style = sp._element.find(qn("p:style"))
    if style is not None:
        sp._element.remove(style)
    _fill(sp, color_name)
    sp.shadow.inherit = False  # geometry: без теней
    return sp


def textbox(slide, x, y, w, h, lines):
    """lines = [(text, role, color_name|None, align)]. Всё из токенов."""
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (text, role, color_name, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        spec = T.role(role)
        r.text = text.upper() if spec.get("caps") else text
        r.font.name = T.font_face(role)          # face из токена (semibold→face)
        r.font.size = Pt(spec["size"])           # размер из токена
        r.font.bold = False                       # эмфаза только через face
        col = color_name or spec.get("color") or "Black"
        r.font.color.rgb = T.rgb(col)
    return tb


def chip(slide, x, y, number):
    """Зелёный квадратный чип-маячок (components.chip): графит regular внутри.
    Размер и кегль — из токена (size=40, text_size=12). Возвращает размер."""
    c = T.comp("chip")
    size = c["size"]
    sp = rect(slide, x, y, size, size, c["fill"])
    tf = sp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number
    r.font.name = T.family            # regular (chip.weight=regular)
    r.font.size = Pt(c.get("text_size", 12))  # из токена chip.text_size (12pt)
    r.font.bold = False
    r.font.color.rgb = T.rgb(c["text"])  # графит, не белый (white-on-green запрещён)
    return size


def add_logo(slide):
    """Лого Cloud.ru в правом верхнем углу (placement.logo.content) — на КАЖДОМ
    контент-слайде, из мастера шаблона. Координаты из токена."""
    lg = T.place("logo")["content"]
    logo_png = os.path.join(HERE, "logo.png")  # извлечён из мастера шаблона
    if os.path.isfile(logo_png):
        slide.shapes.add_picture(logo_png, emu(lg["x"]), emu(lg["y"]),
                                 emu(lg["w"]), emu(lg["h"]))


def add_copyright(slide):
    """Копирайт-футер как в шаблоне (placement.copyright) — КОНСТАНТА из мастера:
    (35,696), 6pt, #BFBFBF, точный текст."""
    cp = T.place("copyright")
    tb = slide.shapes.add_textbox(emu(cp["x"]), emu(cp["y"]), emu(cp["w"]), emu(cp["h"]))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = cp["text"]
    r.font.name = T.family
    r.font.size = Pt(cp["size"])
    r.font.bold = False
    r.font.color.rgb = T.rgb(cp["color"])


def _est_lines(text, width_px, font_pt, caps=False):
    """КОНСЕРВАТИВНАЯ оценка числа строк после переноса (намеренно ЗАВЫШАЕМ, чтобы
    при реальном SB Sans Display заголовок не наехал на тело — user: «следи за этим»)."""
    fp = font_pt * 4.0 / 3.0                  # pt → px
    char_w = (0.72 if caps else 0.62) * fp    # CAPS шире; завышаем ширину символа
    per_line = max(1, int(width_px / char_w))
    return max(1, math.ceil(len(text) / per_line))


def card(slide, x, y, w, h, number, head, body, note, variant="default"):
    """Карточка #F2F2F2 в двух вариантах (components.card.variants):
      default     — чип слева сверху, заголовок ПОД чипом (мало текста);
      text_heavy  — чип в ПРАВОМ верхнем углу, заголовок СЛЕВА вровень с чипом,
                    перенос на 2-ю строку при наложении; тело — НИЖЕ заголовка
                    (высота заголовка считается динамически → нет наложения)."""
    rect(slide, x, y, w, h, T.comp("card")["fill"])
    pad = 24
    cs = T.comp("chip")["size"]
    hpt = T.size("header")
    line_px = 1.18 * (hpt * 4.0 / 3.0)
    if variant == "text_heavy":
        # чип — в правый угол; заголовок слева, ширина обрезана до чипа → авто-перенос
        chip(slide, x + w - pad - cs, y + pad, number)
        head_w = (x + w - pad - cs - 16) - (x + pad)   # стоп перед чипом (gap 16)
        n_lines = _est_lines(head.upper(), head_w, hpt, caps=True)
        head_h = int(n_lines * line_px)
        textbox(slide, x + pad, y + pad, head_w, head_h + 6,
                [(head, "header", "Black", PP_ALIGN.LEFT)])
        # тело — ниже МАКСИМУМА(высота заголовка, высота чипа) + ПОЛНАЯ запасная
        # строка (защита от расхождения метрик шрифта → гарантированно без наложения)
        body_y = y + pad + max(cs, head_h) + int(line_px) + 8
    else:
        chip(slide, x + pad, y + pad, number)
        textbox(slide, x + pad, y + pad + cs + 14, w - 2 * pad, 60,
                [(head, "header", "Black", PP_ALIGN.LEFT)])
        body_y = y + pad + cs + 78
    textbox(slide, x + pad, body_y, w - 2 * pad, h - (body_y - y) - 56,
            [(body, "body", "Black", PP_ALIGN.LEFT)])
    textbox(slide, x + pad, y + h - 44, w - 2 * pad, 28,
            [(note, "caption", "text_gray", PP_ALIGN.LEFT)])


def build():
    prs = Presentation()
    prs.slide_width = emu(T.meta["canvas"]["w"])
    prs.slide_height = emu(T.meta["canvas"]["h"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # фон белый (surface)
    rect(slide, 0, 0, T.meta["canvas"]["w"], T.meta["canvas"]["h"], "White")

    safe = T.safe
    hp = T.place("header")

    # --- Лого Cloud.ru: правый верхний угол (placement.logo.content) ---
    add_logo(slide)

    # --- Заголовок: координаты и роль ИЗ ТОКЕНА placement.header ---
    # NB: НИКАКОЙ декоративной линии под/над заголовком и НИКАКИХ зелёных линий
    # (forbidden.decor.header_rule + forbidden.green_lines, user 2026-06-02).
    textbox(slide, hp["x"], hp["y"], safe["right"] - hp["x"] - 170, 40,
            [("Токен-контракт | два варианта карточек", "header", "Black", PP_ALIGN.LEFT)])
    textbox(slide, hp["x"], hp["y"] + 38, safe["right"] - hp["x"], 24,
            [("Лого (1096,38) из мастера | чип 40×40/12pt | заголовок 20pt CAPS | сетка 35→1245",
              "caption", "text_gray", PP_ALIGN.LEFT)])

    # --- 2 карточки: вариант A (мало текста) и вариант B (много текста) ---
    gap = T.space("card_gap")
    safe_w = safe["right"] - safe["left"]
    cw = (safe_w - gap) / 2
    cy, ch = 200, 400

    card(slide, safe["left"], cy, cw, ch, "01",
         "Вариант A | мало текста",
         "Чип слева сверху, заголовок под чипом. Серая поверхность #F2F2F2, "
         "острые углы. Для коротких блоков.",
         "components.card.variants.default", variant="default")

    card(slide, safe["left"] + cw + gap, cy, cw, ch, "02",
         "Вариант B | много текста с длинным заголовком",
         "Чип уходит в правый верхний угол, заголовок встаёт слева вровень с чипом "
         "и переносится на вторую строку при наложении. Тело начинается ниже "
         "заголовка (высота считается динамически + запас) — наложения нет, и абзац "
         "помещается без уменьшения кегля ниже 16pt.",
         "components.card.variants.text_heavy", variant="text_heavy")

    # Копирайт-футер — КАК В ШАБЛОНЕ (placement.copyright), не произвольный текст.
    add_copyright(slide)

    out = os.path.join(HERE, "point1_tokens.pptx")
    prs.save(out)
    print("saved:", out)
    return out


if __name__ == "__main__":
    build()
