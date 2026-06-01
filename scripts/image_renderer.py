#!/usr/bin/env python3
"""
image_renderer.py — image-as-content rendering с auto-fit.

Зачем: image-content слайды не должны иметь donor с фиксированными PNG-заглушками.
Картинка = главный контент → auto-fit на максимальный canvas с canonical typography.

Layout (1280×720):
  Title:       (35, 38, 1209, 53) px       32pt SemiBold
  Image zone:  (60, 120, 1160, 480) px     auto-fit с aspect-ratio preserved
  Caption:     (35, 620, 1209, 60) px      14pt Regular gray (optional)
"""
import os
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image as PILImage

from kpi_renderer import (
    GRAPHITE, GREEN, WHITE, _add_text_box, clean_slide_to_blank,
    set_slide_title, BLANK_DONOR_WHITE, BLANK_DONOR_DARK
)

EMU = 9525


def render_image_native(slide, image_config, dark=False, wide_zone=False):
    """
    Render image-as-content slide with auto-fit.

    image_config = {
        "title": "Архитектура решения",
        "image_path": "path/to/image.png",
        "caption": "Источник: ..." (optional)
    }
    wide_zone: если True — расширенная image zone (для charts/wide images)
    """
    text_color = WHITE if dark else GRAPHITE

    # Title — в штатный placeholder шаблона, единый стиль (Problem #6)
    title = image_config.get("title", "")
    if title:
        set_slide_title(slide, title, dark=dark)

    # Image zone bounds (px) — wide_zone расширяет width, оставляя место для caption
    if wide_zone:
        ZONE_X = 20
        ZONE_Y = 110
        ZONE_W = 1240
        ZONE_H = 500   # уменьшено с 560 чтобы caption вмещался ниже без перекрытия
        CAPTION_Y = 620
    else:
        ZONE_X = 60
        ZONE_Y = 120
        ZONE_W = 1160
        ZONE_H = 480
        CAPTION_Y = 620

    image_path = image_config.get("image_path")
    if not image_path or not os.path.exists(image_path):
        # Placeholder text если нет картинки
        _add_text_box(slide, ZONE_X, ZONE_Y + ZONE_H // 2 - 30, ZONE_W, 60,
                      "[image not found]", font_size_pt=14, color=text_color,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return

    # Auto-fit: get image dimensions, preserve aspect ratio
    try:
        with PILImage.open(image_path) as img:
            img_w, img_h = img.size
    except Exception:
        img_w, img_h = 1, 1

    img_aspect = img_w / max(img_h, 1)
    zone_aspect = ZONE_W / ZONE_H

    # Mode: "fit" (default — preserve, может остаться empty space)
    # "fill" — растянуть на entire zone (crop overflow)
    mode = image_config.get("mode", "fit")

    if mode == "fill":
        # Fill: всегда занять всю zone (нет empty space)
        # Для квадратных image это даст stretching unless image close to zone aspect
        final_w = ZONE_W
        final_h = ZONE_H
    else:
        # Fit (default): preserve aspect, может оставить empty space
        if img_aspect > zone_aspect:
            final_w = ZONE_W
            final_h = int(ZONE_W / img_aspect)
        else:
            final_h = ZONE_H
            final_w = int(ZONE_H * img_aspect)

    # Center within zone
    final_x = ZONE_X + (ZONE_W - final_w) // 2
    final_y = ZONE_Y + (ZONE_H - final_h) // 2

    slide.shapes.add_picture(
        image_path,
        Emu(final_x * EMU), Emu(final_y * EMU),
        width=Emu(final_w * EMU), height=Emu(final_h * EMU)
    )

    # Caption — позиция зависит от zone (wide → ниже укороченного chart)
    # Для wide_zone (charts) — серая плашка под caption (как в эталоне slide 49)
    caption = image_config.get("caption", "")
    if caption:
        if wide_zone:
            # Серый bg-rectangle под caption (canonical chart стиль)
            from pptx.enum.shapes import MSO_SHAPE
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(0 * EMU), Emu((CAPTION_Y - 10) * EMU),
                Emu(1280 * EMU), Emu(80 * EMU)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
            bg.line.fill.background()  # без обводки
        _add_text_box(slide, 35, CAPTION_Y, 1209, 60, caption,
                      font_size_pt=14 if wide_zone else 12,
                      color=text_color,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# Test standalone
if __name__ == "__main__":
    from pptx import Presentation
    from template_path import resolve_template
    p = Presentation(resolve_template())
    slide = list(p.slides)[BLANK_DONOR_WHITE - 1]
    clean_slide_to_blank(slide)

    render_image_native(slide, {
        "title": "АРХИТЕКТУРА РЕШЕНИЯ",
        "image_path": "pptx-skill/output/draft_images_test_part/slide8_img1.png",
        "caption": "Схема корпоративной аттестованной инфраструктуры"
    })

    out = "pptx-skill/output/image_renderer_test.pptx"
    p.save(out)
    print(f"Saved {out}")
