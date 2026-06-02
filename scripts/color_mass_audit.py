#!/usr/bin/env python3
"""
color_mass_audit.py — измеряет распределение цвета ПО МАССЕ (площади) в .pptx.
Грунтует бренд-правило: нейтрали (белый + F2F2F2) доминируют, зелёный — главный
акцент, дополнительные цвета (Yellow/Purple/Blue) — каждый МЕНЬШЕ зелёного.

Метод: для каждого слайда total = W*H. Для каждой фигуры с solidFill(srgbClr)
area = w*h, бакетит к ближайшему палитровому цвету. Картинки → бакет image.
Фон слайда (если задан srgb/scheme) учитывается; остаток площади → background≈white.
Текст не считаем (масса тонкая). Приближённо (наложения клампятся), но показывает
порядок величин.

CLI:  python3 color_mass_audit.py file1.pptx [file2.pptx ...]
"""
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.oxml.ns import qn

BUCKETS = {
    "White":    "#FFFFFF",
    "Gray(F2)": "#F2F2F2",
    "Graphite": "#222222",
    "Green":    "#26D07C",
    "Yellow":   "#CFF500",
    "Purple":   "#A068FF",
    "Blue":     "#C0E0FC",
}
ADDITIONAL = ("Yellow", "Purple", "Blue")


def _rgb(h):
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _nearest_bucket(hex6):
    r, g, b = _rgb(hex6)
    best, bestd = "Other", 9e18
    for name, thex in BUCKETS.items():
        tr, tg, tb = _rgb(thex)
        d = (r - tr) ** 2 + (g - tg) ** 2 + (b - tb) ** 2
        if d < bestd:
            best, bestd = name, d
    # «Other» если далеко от всех (порог большой → почти всё попадает в палитру)
    return best if bestd < 9000 else "Other"


def _shape_fill_hex(shape):
    """→ hex solidFill srgbClr фигуры или None."""
    el = shape._element
    spPr = el.find(qn("p:spPr"))
    if spPr is None:
        return None
    sf = spPr.find(qn("a:solidFill"))
    if sf is None:
        return None
    srgb = sf.find(qn("a:srgbClr"))
    if srgb is not None:
        return "#" + srgb.get("val")
    return None


def _iter(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:
            try:
                yield from _iter(sh.shapes)
            except Exception:
                pass


def _run_color_hex(run):
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        return None
    sf = rPr.find(qn("a:solidFill"))
    if sf is None:
        return None
    srgb = sf.find(qn("a:srgbClr"))
    return "#" + srgb.get("val") if srgb is not None else None


def audit(path, max_slides=None):
    prs = Presentation(path)
    W = prs.slide_width / 9525.0
    H = prs.slide_height / 9525.0
    slide_area = W * H
    area = {k: 0.0 for k in list(BUCKETS) + ["Other", "image"]}
    # «ink» текста: прокси-площадь = len(text) * size_pt (для green-vs-additional)
    ink = {k: 0.0 for k in list(BUCKETS) + ["Other"]}
    n = 0
    for i, slide in enumerate(prs.slides):
        if max_slides and i >= max_slides:
            break
        n += 1
        covered = 0.0
        for sh in _iter(slide.shapes):
            try:
                a = (sh.width / 9525.0) * (sh.height / 9525.0)
            except Exception:
                a = 0.0
            if sh.shape_type == 13:  # picture
                area["image"] += a
                covered += a
                continue
            hx = _shape_fill_hex(sh)
            if hx:
                area[_nearest_bucket(hx)] += a
                covered += a
            # текст-цвет (ink-прокси)
            if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        ch = _run_color_hex(r)
                        if not ch:
                            continue
                        sz = r.font.size.pt if r.font.size else 16
                        ink[_nearest_bucket(ch)] += len(r.text or "") * sz
        bg = max(0.0, slide_area - covered)
        area["White"] += bg
    total = sum(area.values()) or 1.0
    inktot = sum(ink.values()) or 1.0
    return (n,
            {k: 100.0 * v / total for k, v in area.items()},
            {k: 100.0 * v / inktot for k, v in ink.items()})


def main():
    max_slides = None
    args = [a for a in sys.argv[1:]]
    if args and args[0].startswith("--max="):
        max_slides = int(args[0].split("=")[1]); args = args[1:]
    for path in args:
        try:
            n, pct, ink = audit(path, max_slides=max_slides)
        except Exception as e:
            print(f"\n{os.path.basename(path)}: ERROR {e}")
            continue
        name = os.path.basename(path)
        print(f"\n=== {name[:58]} ({n} слайдов) ===")
        print("  [МАССА поверхности: фон+заливки+картинки]")
        for k in ["White", "Gray(F2)", "Graphite", "Green", "Yellow", "Purple", "Blue", "image", "Other"]:
            print(f"    {k:9} {pct[k]:5.1f}%  {'#' * int(pct[k] / 2)}")
        print("  [АКЦЕНТ-ink текста: green vs дополнительные]")
        for k in ["Graphite", "Green", "Yellow", "Purple", "Blue", "Other"]:
            print(f"    {k:9} {ink[k]:5.1f}%  {'#' * int(ink[k] / 2)}")
        neutral = pct["White"] + pct["Gray(F2)"]
        addl_ink = max(ink["Yellow"], ink["Purple"], ink["Blue"])
        verdict = "OK: additional<green" if addl_ink <= ink["Green"] + 0.01 else "WARN: additional>green"
        print(f"  -- нейтрали(масса)={neutral:.1f}% | green(ink)={ink['Green']:.1f}% "
              f"max(additional ink)={addl_ink:.1f}% → {verdict}")


if __name__ == "__main__":
    main()
