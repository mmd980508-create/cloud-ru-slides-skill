#!/usr/bin/env python3
"""
render_slides.py — конвертирует .pptx в серию PNG через LibreOffice.

Usage:
    python3 render_slides.py <input.pptx> <output_dir/>

Workflow:
    1. soffice --headless --convert-to pdf input.pptx → input.pdf
    2. pdftoppm input.pdf <output_dir>/slide -png -r 96 → slide-01.png ...
"""
import sys, os, subprocess, shutil, tempfile


# Начертания SemiBold, которые LibreOffice НЕ умеет показывать жирным
# (схлопывает в Regular). Для превью подменяем их на обычный face + bold-флаг,
# чтобы эмфаза была видна в PNG. На сам .pptx это не влияет — подмена идёт
# только во временной копии (Problem #3, 2026-05-29; preview-only).
PREVIEW_SEMIBOLD_FACES = ("SB Sans Display Semibold", "SB Sans Display SemiBold")
PREVIEW_BASE_FACE = "SB Sans Display"


def _iter_runs(shapes):
    """Рекурсивно обходит все run'ы в shapes (включая группы и таблицы)."""
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                yield from _iter_runs(shape.shapes)
            except Exception:
                pass
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            yield r
            continue
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    yield r


def _make_preview_copy(input_pptx, tmpdir):
    """Создаёт временную копию .pptx, где SemiBold-face заменён на обычный
    face + bold-флаг — чтобы LibreOffice отрисовал эмфазу жирным в превью.
    Возвращает путь к копии (или исходник, если python-pptx недоступен)."""
    try:
        from pptx import Presentation
    except ImportError:
        return input_pptx
    try:
        prs = Presentation(input_pptx)
        changed = 0
        for slide in prs.slides:
            for run in _iter_runs(slide.shapes):
                if run.font.name in PREVIEW_SEMIBOLD_FACES:
                    run.font.name = PREVIEW_BASE_FACE
                    run.font.bold = True
                    changed += 1
        if changed == 0:
            return input_pptx
        out = os.path.join(tmpdir, "_preview_" + os.path.basename(input_pptx))
        prs.save(out)
        print(f"[preview] SemiBold→bold для превью: {changed} run(ов)", file=sys.stderr)
        return out
    except Exception as e:
        print(f"[preview] подмена шрифта пропущена: {e}", file=sys.stderr)
        return input_pptx


def find_soffice():
    home = os.path.expanduser("~")
    candidates = [
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        f"{home}/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        # Linux
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        # Windows
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # PATH lookup (любая ОС)
        "soffice",
        "soffice.exe",
    ]
    for c in candidates:
        if os.path.isfile(c) or shutil.which(c):
            return c
    return None


def pdf_to_pngs(pdf_path, output_dir, dpi=96):
    """PDF → PNG series. Uses pdftoppm if available, else PyMuPDF."""
    prefix = os.path.join(output_dir, "slide")
    if shutil.which("pdftoppm"):
        result = subprocess.run([
            "pdftoppm", "-png", "-r", str(dpi), pdf_path, prefix
        ], capture_output=True, text=True, timeout=300)
        if result.returncode != 0:
            raise RuntimeError(f"pdftoppm failed: {result.stderr}")
        return
    try:
        import fitz  # PyMuPDF
    except ImportError as e:
        raise RuntimeError(
            "Ни pdftoppm, ни PyMuPDF не доступны. "
            "Поставь: pip3 install --user pymupdf"
        ) from e
    doc = fitz.open(pdf_path)
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    width = len(str(len(doc)))
    for i, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=mat, alpha=False)
        out = f"{prefix}-{str(i).zfill(max(2, width))}.png"
        pix.save(out)
    doc.close()


def render(input_pptx, output_dir, dpi=96):
    soffice = find_soffice()
    if not soffice:
        print("ERROR: soffice/LibreOffice не найден. Установи: brew install --cask libreoffice", file=sys.stderr)
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmpdir:
        # Preview-only: SemiBold-face → regular+bold, чтобы LibreOffice показал
        # эмфазу жирным. Оригинал input_pptx не меняется.
        render_src = _make_preview_copy(input_pptx, tmpdir)

        # Step 1: pptx → pdf
        print(f"[1/2] Конвертирую {input_pptx} → PDF…", file=sys.stderr)
        result = subprocess.run([
            soffice, "--headless", "--convert-to", "pdf",
            "--outdir", tmpdir, render_src
        ], capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            print(f"soffice failed: {result.stderr}", file=sys.stderr)
            sys.exit(1)

        pdf_name = os.path.splitext(os.path.basename(render_src))[0] + ".pdf"
        pdf_path = os.path.join(tmpdir, pdf_name)
        if not os.path.exists(pdf_path):
            print(f"ERROR: PDF не создан: {pdf_path}", file=sys.stderr)
            sys.exit(1)

        # Step 2: pdf → png series (pdftoppm если есть, иначе PyMuPDF)
        print(f"[2/2] PDF → PNG ({dpi} dpi)…", file=sys.stderr)
        try:
            pdf_to_pngs(pdf_path, output_dir, dpi=dpi)
        except RuntimeError as e:
            print(str(e), file=sys.stderr)
            sys.exit(1)

        files = sorted(f for f in os.listdir(output_dir) if f.startswith("slide") and f.endswith(".png"))
        print(f"OK: {len(files)} слайдов в {output_dir}/", file=sys.stderr)
        for f in files:
            print(f"  {f}")


def main():
    if len(sys.argv) != 3:
        print("Usage: render_slides.py <input.pptx> <output_dir/>", file=sys.stderr)
        sys.exit(1)
    render(sys.argv[1], sys.argv[2])


if __name__ == "__main__":
    main()
