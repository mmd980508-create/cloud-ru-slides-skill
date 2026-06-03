#!/usr/bin/env python3
"""
Point 7 — KPI: голый вариант (без акцента) vs card-вариант (акцентная карточка
зелёная). Цифры Regular, графит; на зелёной карточке — тоже графит.
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "..", "scripts"))
from pptx import Presentation                   # noqa: E402
from pptx.util import Emu                        # noqa: E402
import kpi_renderer                              # noqa: E402

EMU = 9525
NUMS = [
    {"value": "12", "desc": "AI-инициатив", "accent": False, "pct": False},
    {"value": "84", "desc": "вовлечённость", "accent": False, "pct": True},
    {"value": "30", "desc": "инструментов", "accent": True, "pct": False},
]


def build():
    prs = Presentation()
    prs.slide_width = Emu(1280 * EMU); prs.slide_height = Emu(720 * EMU)

    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    kpi_renderer.render_kpi(s1, {"title": "Голые цифры · без акцента (Regular)",
                                 "numbers": NUMS})

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    kpi_renderer.render_kpi(s2, {"title": "В карточках · акцентная карточка зелёная",
                                 "layout": "cards", "numbers": NUMS})

    out = os.path.join(HERE, "kpi_cards_demo.pptx")
    prs.save(out); print("saved:", out)
    return out


if __name__ == "__main__":
    build()
