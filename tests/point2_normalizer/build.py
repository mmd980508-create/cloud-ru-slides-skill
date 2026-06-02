#!/usr/bin/env python3
"""
Point 2 test — нормализатор «чужого скина» к Cloud.ru.
Строит ЗАВЕДОМО НЕ-брендовый слайд (скругления, тени, градиент, Arial italic,
произвольные цвета), затем прогоняет brand_normalizer → бренд. Рендерит ОБА:
  before.pptx — импортированный референс (чужой скин)
  after.pptx  — нормализовано к токенам Cloud.ru
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", "..", "scripts"))

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

import brand_normalizer

EMU = 9525


def emu(px):
    return Emu(int(round(px * EMU)))


def add_shadow(shape):
    """Явный drop-shadow (чтобы было что убирать нормализатору)."""
    spPr = shape._element.find(qn("p:spPr"))
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"),
                         {"blurRad": "50000", "dist": "25000",
                          "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "000000"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": "45000"}))
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def rrect(slide, x, y, w, h, fill_hex, grad_to=None):
    """Скруглённый прямоугольник + тень + (опц.) градиент — чужой скин."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                emu(x), emu(y), emu(w), emu(h))
    if grad_to:
        try:
            sp.fill.gradient()
            sp.fill.gradient_stops[0].color.rgb = RGBColor.from_string(fill_hex)
            sp.fill.gradient_stops[1].color.rgb = RGBColor.from_string(grad_to)
        except Exception:
            sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    sp.line.fill.background()
    add_shadow(sp)
    return sp


def foreign_text(slide, x, y, w, h, lines):
    """lines = [(text, pt, hex, italic, bold)] — Arial, часто italic (чужой шрифт)."""
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, pt, hexc, ital, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"            # чужой шрифт
        r.font.size = Pt(pt)
        r.font.italic = ital
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(hexc)
    return tb


def build_before():
    prs = Presentation()
    prs.slide_width = emu(1280)
    prs.slide_height = emu(720)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # заголовок — Arial italic, тёмно-синий (не бренд)
    foreign_text(s, 35, 38, 1100, 48,
                 [("БЫЛО · импортированный референс (чужой скин)", 22, "1B3A6B", True, True)])

    cards = [
        (35,  "1FB6A6", None,     "Облачная платформа", "FFFFFF"),   # бирюзовый
        (450, "FF7A00", "FF3D81", "Аналитика данных",   "FFFFFF"),   # оранж→розовый градиент
        (865, "7B5BD6", None,     "DevOps сервисы",      "FFFFFF"),   # фиолетовый
    ]
    for x, fill, grad, title, txt in cards:
        rrect(s, x, 150, 380, 360, fill, grad_to=grad)
        foreign_text(s, x + 28, 180, 324, 300, [
            (title, 24, txt, True, True),
            ("Скруглённые углы, тень, градиент и Arial italic — "
             "всё это чужой скин, который нужно снять.", 14, txt, True, False),
        ])

    # маленький «чип» — тоже скруглённый, произвольный цвет
    rrect(s, 35, 540, 70, 70, "E8456B")

    out = os.path.join(HERE, "before.pptx")
    prs.save(out)
    return out


def main():
    before = build_before()
    after = os.path.join(HERE, "after.pptx")
    stats = brand_normalizer.normalize_pptx(before, after)
    print("before:", before)
    print("after :", after)
    print("normalized stats:", stats)


if __name__ == "__main__":
    main()
