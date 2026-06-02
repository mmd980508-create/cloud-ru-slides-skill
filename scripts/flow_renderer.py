#!/usr/bin/env python3
"""
flow_renderer.py — native PowerPoint flow-diagrams / schemas / process maps.

Зачем: до v1.6 скилл умел charts (chart_pptx_native) и KPI (kpi_native),
но не умел собирать схемы блок-стрелка-блок (SGR, pipeline, process map).
Альтернатива была — рисовать в Figma и вставлять PNG, что не редактируется.

Решение: модуль рисует схему из примитивов PowerPoint shapes на blank canvas.
Все элементы редактируемые: блок можно подвинуть, текст переписать, стрелку
перенаправить — обычными средствами PowerPoint.

Используется build_v9 при slide_type == "flow_diagram_native".

Стиль (canonical, согласован с эталонными слайдами Cloud.ru;
правило 2026-05-06 — feedback_flow_diagram_composition.md):
  - Блоки: серые #F2F2F2, текст SB Sans Display #222222.
  - Текст в блоках: align=LEFT, vertical_anchor=TOP. Поля: 12px со всех сторон,
    низ 16px.
  - Стрелки: единая толщина 1pt, цвет #434343 (тёмно-серый, не #222222 графит).
    Наконечник: открытая галочка (type='arrow'), размер 8 в PPT UI (w=lg,
    len=med). Только горизонтальные/вертикальные — диагонали запрещены.
  - Пунктирные рамки для группировки фаз: 1pt #888888 dash.
  - Заголовок: 20pt SemiBold CAPS, top-left (35, 60).
  - Decor (по бренду): зелёные стрелки ↗ (нативные фигуры-группы, 1pt) — bottom corner.
  - Safe-area: SAFE_TOP=140, SAFE_BOTTOM=660, SAFE_LEFT=35, SAFE_RIGHT=1245
    (совпадает с направляющими PowerPoint и TITLE-placeholder @35 — заголовок и
    фреймы выровнены по ОДНИМ границам). Все координаты блоков лежат внутри.

Координаты — в пикселях (slide 1280×720). EMU = px × 9525.

Config schema (передаётся через plan.json):
  {
    "slide_type": "flow_diagram_native",
    "dark": false,
    "flow": {
      "header": "Заголовок схемы",            # обязательное
      "subtitle": "...",                       # опц., 13pt под header'ом
      "subtitle_url": "https://...",           # опц., 12pt серая ссылка
      "blocks": [
        {
          "id": "b1",                          # опц., для arrows by ref
          "x": 175, "y": 180, "w": 235, "h": 50,
          "lines": ["Title", "subtitle text"],
          "font_sizes": [16, 16],              # опц., стандарт 16pt (мин 12)
          "bolds": [true, false],              # опц.
          "caps_first": false,                 # опц.
          "fill": "gray|green|dark|white",     # опц., default "gray"
          "align": "left|center|right",        # опц., default "left" (canonical)
          "vanchor": "top|middle"              # опц., default "top" (canonical)
        }
      ],
      "groups": [                              # опц., dashed рамки + лейблы
        {
          "label": "Phase 1",
          "x": 167, "y": 154, "w": 251, "h": 86,
          "label_pos": "top|bottom"            # default "top"
        }
      ],
      "arrows": [
        # вариант 1 — явные координаты
        {"x1": 410, "y1": 205, "x2": 432, "y2": 205,
         "with_head": true, "dashed": false},
        # вариант 2 — между блоками по id
        {"from": "b1", "to": "b2", "side": "right"}
      ],
      "labels": [                              # опц., свободные подписи
        {"x": 35, "y": 122, "w": 600, "h": 22,
         "text": "SGR · Schema-Guided Reasoning",
         "font_size": 11, "bold": false,
         "align": "left", "caps": false}
      ],
      "decor": {                               # опц., зелёные стрелки ↗ (нативные фигуры, 1pt)
        "enabled": true,
        "count": 4,
        "x_start": 950, "y_start": 625,
        "size": 38, "gap": 12
      }
    }
  }
"""
import os
import io
import json

from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree


EMU = 9525
SLIDE_W_PX = 1280
SLIDE_H_PX = 720


# ============================================================================
# Палитра (из brand/palette.json — единый источник истины)
# ============================================================================
def _load_palette():
    here = os.path.dirname(os.path.abspath(__file__))
    for path in (
        os.path.join(here, "..", "brand", "palette.json"),
        os.path.join(os.getcwd(), "brand", "palette.json"),
        os.path.join(os.getcwd(), "pptx-skill", "brand", "palette.json"),
    ):
        try:
            if os.path.isfile(path):
                with open(path, encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            continue
    return None


_PAL = _load_palette()


def _hex(hx):
    h = hx.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _from_palette(section, name, fallback):
    if _PAL and name in _PAL.get(section, {}):
        return _hex(_PAL[section][name])
    return _hex(fallback)


GRAPHITE = _from_palette("base", "Black", "#222222")
WHITE = _from_palette("base", "White", "#FFFFFF")
GRAY = _from_palette("base", "Gray", "#F2F2F2")
GREEN = _from_palette("base", "Green", "#26D07C")
# Тёмная плашка = #222222 (canonical 2026-05-29, user-decision). НЕ #343F48 —
# Graphite-Iron не из палитры Cloud.ru.
DARK_FILL = GRAPHITE
DASH_GRAY = RGBColor(0x88, 0x88, 0x88)
SEPARATOR_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
# Тело текста в эталонных деках — средний серый (заголовок графит, тело светлее).
# 3-уровневая иерархия (правило A-6). Сверено по референсам 2026-06-01.
TEXT_GRAY = RGBColor(0x5C, 0x5C, 0x5C)

# Canonical для flow-схем (правило 2026-05-06):
ARROW_COLOR = RGBColor(0x43, 0x43, 0x43)   # единый цвет всех стрелок
ARROW_WIDTH_PT = 1.0                        # единая толщина всех стрелок
# Микроправило (2026-05-29): резерв (px) под наконечник на финальном сегменте
# ортогональной стрелки-ветки — чтобы окончание стрелки хорошо просматривалось.
ARROW_ENTRY_RESERVE = 22
# Зазор между колонками грида при наличии branching-стрелок — больше обычного,
# чтобы веткам/наконечникам хватало места.
GRID_GAP_DEFAULT = 24
GRID_GAP_BRANCHING = 44
# Базовый отступ МЕЖДУ фреймами в пресетах-архетипах (правило 2026-06-01):
# не более 10px, базово 4px — как в эталонных деках (плотный ритм).
PRESET_GAP = 4


FONT = "SB Sans Display"
# Полужирное начертание = ОТДЕЛЬНЫЙ font face, встроенный в шаблон (НЕ bold-флаг).
# Canonical 2026-05-29 (Problem #3): Bold запрещён — эмфаза только через SemiBold.
FONT_SEMIBOLD = "SB Sans Display Semibold"


def _set_weight(font, semibold):
    """Эмфаза через начертание SemiBold, а не bold-флаг (Problem #3).

    semibold=True → face «SB Sans Display Semibold» (встроен в шаблон), bold=False.
    semibold=False → обычный «SB Sans Display».
    Жирный (bold=True) НЕ используется нигде.
    """
    font.name = FONT_SEMIBOLD if semibold else FONT
    font.bold = False


# Safe-area для 1280×720 Cloud.ru slide.
# Все координаты flow-схем должны рассчитываться относительно этих констант.
SAFE_TOP = 140       # под заголовком 20pt
SAFE_BOTTOM = 660    # выше footer copyright
# SAFE_LEFT/RIGHT = направляющие PowerPoint + левый край TITLE-placeholder (@35).
# Раньше было 30/1250 — фреймы выходили за поля и не совпадали с заголовком.
SAFE_LEFT = 35
SAFE_RIGHT = 1245
SAFE_W = SAFE_RIGHT - SAFE_LEFT   # 1210
SAFE_H = SAFE_BOTTOM - SAFE_TOP    # 520


# ============================================================================
# Низкоуровневые примитивы (можно импортировать отдельно для кастомных схем)
# ============================================================================
def px(v):
    """px → EMU."""
    return Emu(int(v) * EMU)


def _no_effects(shape):
    """Снять любые эффекты (тени/glow/reflection) — и явные, и тематические
    (effectRef). Бренд: плоско, эффектов нет вообще. См. effects_util."""
    from effects_util import strip_effects
    strip_effects(shape._element)


def _resolve_fill(fill_name):
    """Имя fill → (rgb_color, text_color)."""
    if fill_name in (None, "gray"):
        return GRAY, GRAPHITE
    if fill_name == "white":
        return WHITE, GRAPHITE
    if fill_name == "green":
        # Canonical 2026-05-29 (Problem #2): на зелёной плашке текст #222222,
        # НЕ белый. White-on-green плохо читается — текст всегда графит.
        return GREEN, GRAPHITE
    if fill_name == "dark":
        return DARK_FILL, WHITE
    # raw hex
    if isinstance(fill_name, str) and fill_name.startswith("#"):
        return _hex(fill_name), GRAPHITE
    return GRAY, GRAPHITE


# Маркеры для auto-detect bullets (canonical 2026-05-26, исправлено 2026-05-29):
# Если строка начинается с одного из этих префиксов — превратить в native PP bullet.
# Используем штатный PowerPoint «Заполненные квадратные маркеры» = шрифт Wingdings,
# символ "§" (0xA7) → плоский залитый квадрат. Это НЕ символ ▪ U+25AA: тот в части
# рендереров подхватывает emoji-вариант и даёт «эффект объёма». Wingdings § —
# обычный ASCII-символ, emoji-фолбэка не бывает, всегда плоский квадрат.
BULLET_PREFIXES = ("▪ ", "■ ", "• ", "- ", "* ")
BULLET_FONT = "Wingdings"
BULLET_CHAR = "§"  # § в Wingdings = заполненный квадрат (штатный PP bullet)
BULLET_COLOR = "434343"  # графит, дефолтный цвет маркера
BULLET_INDENT_EMU = 228600  # ~24px hanging indent для маркера


def _detect_bullet(line):
    """Возвращает (is_bullet, clean_text). is_bullet=True если строка
    начинается с одного из BULLET_PREFIXES — префикс удаляется."""
    for prefix in BULLET_PREFIXES:
        if line.startswith(prefix):
            return True, line[len(prefix):]
    return False, line


def _apply_bullet_to_paragraph(p, char=BULLET_CHAR, indent_emu=BULLET_INDENT_EMU):
    """Добавить штатный PowerPoint bullet «Заполненный квадрат» к параграфу.

    Это настоящий PP-буллет (Wingdings §, цвет #434343) — пользователь может в PP
    кликнуть «Маркеры» → выбрать другой стиль. Без эффектов объёма: § — плоский
    залитый квадрат, без emoji-фолбэка.

    Порядок дочерних элементов pPr по схеме OOXML: buClr → buFont → buChar.
    """
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    # marL и indent — hanging indent для маркера
    pPr.set("marL", str(indent_emu))
    pPr.set("indent", str(-indent_emu))
    # Удалить существующие bullet-элементы
    for tag in ("buClr", "buSzPct", "buNone", "buChar", "buAutoNum", "buFont"):
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    # buClr — цвет маркера #434343 (должен идти ПЕРЕД buFont/buChar)
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    srgb = etree.SubElement(buClr, qn("a:srgbClr"))
    srgb.set("val", BULLET_COLOR)
    # buFont — Wingdings (штатный набор PP filled square)
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", BULLET_FONT)
    buFont.set("pitchFamily", "2")
    buFont.set("charset", "2")
    # buChar — § = заполненный квадрат
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)


def _apply_no_bullet_to_paragraph(p):
    """Явно отключить bullet — для строк без маркера в списке смешанного типа."""
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    for tag in ("buClr", "buSzPct", "buChar", "buAutoNum", "buFont"):
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    # Удалить marL/indent — текст идёт от левого края
    for attr in ("marL", "indent"):
        if pPr.get(attr):
            del pPr.attrib[attr]
    # Явно <a:buNone/>
    has_buNone = pPr.find(qn("a:buNone")) is not None
    if not has_buNone:
        etree.SubElement(pPr, qn("a:buNone"))


def add_block(slide, x, y, w, h, lines,
              font_sizes=None, bolds=None,
              caps_first=False, fill="gray", align="left",
              vanchor="top",
              text_color=None, text_colors=None):
    """Прямоугольник с многострочным контентом.

    Canonical (правило 2026-05-06):
    - Текст ВСЕГДА выравнивается по ЛЕВОМУ (align='left') и ВЕРХНЕМУ (vanchor='top') краю.
    - Поля одинаковые 12px, нижнее 16px.

    lines: list[str].
    font_sizes / bolds: list of same len. По умолчанию [16, 16, ...] (стандарт 16pt,
        мин 12pt — Problem #5) / [True, False, ...].
    caps_first: первая строка → CAPS.
    fill: "gray" (default) / "white" / "green" / "dark" / hex string.
    align: "left" (default, canonical) / "center" / "right".
    vanchor: "top" (default, canonical) / "middle".
    text_color: явный override (RGBColor); иначе из fill.
    text_colors: список цветов по строкам (3-уровневая иерархия — заголовок графит,
        тело TEXT_GRAY). Имеет приоритет над text_color для соответствующей строки.
    """
    rgb_fill, default_text = _resolve_fill(fill)
    text = text_color if text_color is not None else default_text

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_fill
    shape.line.fill.background()
    _no_effects(shape)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP if vanchor == "top" else MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(12 * EMU)
    tf.margin_right = Emu(12 * EMU)
    tf.margin_top = Emu(12 * EMU)
    tf.margin_bottom = Emu(16 * EMU)

    if font_sizes is None:
        # Стандарт 16pt (Problem #5, 2026-05-29). Для плотных схем LLM задаёт
        # меньше через font_sizes, но не ниже 12pt (кроме крайней перегрузки).
        font_sizes = [16] + [16] * max(0, len(lines) - 1)
    if bolds is None:
        bolds = [True] + [False] * max(0, len(lines) - 1)

    align_enum = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_enum
        # Auto-detect bullet — если строка начинается с ▪/■/•/-/* → native PP bullet
        is_bullet, clean_line = _detect_bullet(line)
        if is_bullet:
            _apply_bullet_to_paragraph(p)
        run = p.add_run()
        run.text = clean_line.upper() if (caps_first and i == 0) else clean_line
        run.font.size = Pt(font_sizes[i] if i < len(font_sizes) else font_sizes[-1])
        _set_weight(run.font, bolds[i] if i < len(bolds) else bolds[-1])
        if text_colors is not None and i < len(text_colors) and text_colors[i] is not None:
            run.font.color.rgb = text_colors[i]
        else:
            run.font.color.rgb = text
    return shape


def add_label(slide, x, y, w, h, text,
              font_size=16, bold=False,
              align="left", anchor="top", caps=False,
              color=None):
    """Свободная текстовая подпись (без фона)."""
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor == "middle" else MSO_ANCHOR.TOP
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    align_enum = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_enum
        # Auto-detect bullet
        is_bullet, clean_line = _detect_bullet(line)
        if is_bullet:
            _apply_bullet_to_paragraph(p)
        run = p.add_run()
        run.text = clean_line.upper() if caps else clean_line
        run.font.size = Pt(font_size)
        _set_weight(run.font, bold)
        run.font.color.rgb = color if color is not None else GRAPHITE
    return box


def add_arrow(slide, x1, y1, x2, y2,
              with_head=True, w_pt=None, dashed=False, color=None):
    """Прямая стрелка/коннектор.

    Canonical (правило 2026-05-06):
    - Только горизонтальные или вертикальные. Диагонали запрещены (валидация ниже).
    - Толщина: единая ARROW_WIDTH_PT = 1.0 (w_pt override только для исключений).
    - Цвет: единый ARROW_COLOR = #434343.
    - Голова: открытая галочка (type='arrow') размер 8 (w=lg, len=med).
    """
    if x1 != x2 and y1 != y2:
        raise ValueError(
            f"add_arrow: диагональная стрелка ({x1},{y1})→({x2},{y2}) запрещена. "
            "Используй ломаную через 90° (несколько add_arrow с промежуточными точками)."
        )
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2)
    )
    conn.line.color.rgb = color if color is not None else ARROW_COLOR
    effective_w = w_pt if w_pt is not None else ARROW_WIDTH_PT
    conn.line.width = Emu(int(effective_w * 12700))
    if dashed:
        ln = conn.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    if with_head:
        ln = conn.line._get_or_add_ln()
        for tag in ("tailEnd", "headEnd"):
            for el in ln.findall(qn(f"a:{tag}")):
                ln.remove(el)
        tail_end = etree.SubElement(ln, qn("a:tailEnd"))
        # PowerPoint UI «открытая стрелка, размер 8» = (type=arrow, w=lg, len=med).
        # (w=lg, len=lg) = размер 9 (максимум).
        tail_end.set("type", "arrow")
        tail_end.set("w", "lg")
        tail_end.set("len", "med")
    return conn


def add_orthogonal_arrow(slide, x1, y1, x2, y2, color=None, w_pt=None):
    """Ортогональная стрелка-ветка (Z-маршрут) для связей между ячейками РАЗНЫХ
    строк И колонок — вместо запрещённой диагонали (canonical: только прямые углы).

    Маршрут: горизонтальный выход → вертикаль в середине зазора → горизонтальный
    вход в цель. Голова — только на последнем сегменте (одна стрелка = одна голова).
    """
    if y1 == y2 or x1 == x2:
        return add_arrow(slide, x1, y1, x2, y2, with_head=True, color=color, w_pt=w_pt)
    # Микроправило (2026-05-29): резервируем место под наконечник на ФИНАЛЬНОМ
    # сегменте (вход в бокс), чтобы окончание стрелки хорошо просматривалось.
    # Вертикальную «шину» ставим ближе к источнику → вход в цель длиннее.
    entry = ARROW_ENTRY_RESERVE
    if x2 > x1:
        jx = max(x1 + 8, x2 - entry)
    else:
        jx = min(x1 - 8, x2 + entry)
    add_arrow(slide, x1, y1, jx, y1, with_head=False, color=color, w_pt=w_pt)
    add_arrow(slide, jx, y1, jx, y2, with_head=False, color=color, w_pt=w_pt)
    return add_arrow(slide, jx, y2, x2, y2, with_head=True, color=color, w_pt=w_pt)


def add_dashed_rect(slide, x, y, w, h, color=None, w_pt=1.0):
    """Пунктирная рамка (для группировки phase-блоков)."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    rect.fill.background()
    rect.line.color.rgb = color if color is not None else DASH_GRAY
    rect.line.width = Emu(int(w_pt * 12700))
    ln = rect.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")
    _no_effects(rect)
    return rect


def snap_panel_to_safe(x, y, w, h, tol=24):
    """Притягивает края панели к направляющим safe-area, если они РЯДОМ (в пределах
    tol px). Чинит два дефекта hand-coded панелей: (1) выход за поля и (2) пустой
    зазор справа («фреймы не доведены до конца»). Не трогает края, которые явно
    внутри (намеренный отступ) — двигаем только то, что метило в границу.

    Возвращает (x, y, w, h). Координаты — px.
    """
    left, right = x, x + w
    top, bottom = y, y + h
    if abs(left - SAFE_LEFT) <= tol:
        left = SAFE_LEFT
    if abs(right - SAFE_RIGHT) <= tol:
        right = SAFE_RIGHT
    if abs(top - SAFE_TOP) <= tol:
        top = SAFE_TOP
    if abs(bottom - SAFE_BOTTOM) <= tol:
        bottom = SAFE_BOTTOM
    return left, top, right - left, bottom - top


def add_filled_panel(slide, x, y, w, h, label=None, dark=False, fill=None):
    """Залитая серая панель-секция для группировки в НАГРУЖЕННЫХ схемах — вместо
    пунктирной рамки (правило 2026-05-29: меньше пунктира, читается чище).

    Рисуется ФОНОМ (до блоков). Внутри — section-title в левом-верхнем углу;
    сами блоки внутри делать `fill="white"`, чтобы карточки выделялись на сером.
    """
    rgb = fill if fill is not None else GRAY
    if isinstance(rgb, str):
        rgb = _hex(rgb)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = rgb
    rect.line.fill.background()
    _no_effects(rect)
    if label:
        add_label(slide, x + 16, y + 12, w - 32, 22, label,
                  font_size=13, bold=True, caps=True,
                  color=WHITE if dark else GRAPHITE)
    return rect


def _draw_check(slide, x, y, size, color=None, w_pt=2.0):
    """Иконка-галочка как в эталоне (Чек-лист успеха): КОНТУРНЫЙ КРУЖОК + галочка
    внутри. Кружок — обводка цветом col без заливки, галочка — две линии-сегмента
    со скруглёнными концами. Всё нативные фигуры (редактируемые). Сверено по
    референсу 2026-06-01: иконка графит на зелёном чипе, не сплошной глиф.
    Координаты — относительно бокса (x,y,size)."""
    col = color if color is not None else GRAPHITE
    # Кружок-обводка занимает ~76% бокса, центрирован.
    d = size * 0.76
    cx = x + (size - d) / 2.0
    cy = y + (size - d) / 2.0
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx), px(cy), px(d), px(d))
    ring.fill.background()
    ring.line.color.rgb = col
    ring.line.width = Emu(int(w_pt * 12700))
    _no_effects(ring)
    # Галочка ✓ внутри кружка (координаты относительно всего бокса).
    pts_rel = [(0.33, 0.52), (0.45, 0.64), (0.68, 0.38)]
    abspts = [(x + rx * size, y + ry * size) for rx, ry in pts_rel]
    for (x1, y1), (x2, y2) in zip(abspts, abspts[1:]):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
        conn.line.color.rgb = col
        conn.line.width = Emu(int(w_pt * 12700))
        ln = conn.line._get_or_add_ln()
        ln.set("cap", "rnd")


def add_green_chip(slide, x, y, size=50, number=None, check=False,
                   font_size=22, fill=None, text_color=None, bold=False):
    """Фирменный зелёный квадрат-чип ~50×50px — повторяющийся мотив Cloud.ru
    (design-principles-from-decks.md, правило A-10).

    Внутри по центру — цифра (01–05) ИЛИ контурная галочка. Ставится слева от
    заголовка / строки / карточки. Дозированный зелёный акцент.

    ВАЖНО (canonical): контент на зелёном — ГРАФИТ #222222, НЕ белый (white-on-green
    запрещён). Цифры — REGULAR-начертанием (не bold), как в эталоне.

    number: строка/число ("01" или 1). check=True → контурная галочка вместо цифры.
    """
    col = text_color if text_color is not None else GRAPHITE
    chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  px(x), px(y), px(size), px(size))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill if fill is not None else GREEN
    chip.line.fill.background()
    _no_effects(chip)
    if check:
        _draw_check(slide, x, y, size, color=col)
        return chip
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    if number is not None:
        run.text = str(number)
    run.font.size = Pt(font_size)
    _set_weight(run.font, bold)
    run.font.color.rgb = col
    return chip


def add_header(slide, text, dark=False):
    """Заголовок слайда — в штатный TITLE-placeholder шаблона (canonical позиция/
    размер 35,38 / 963×54 / 20pt SemiBold CAPS). Problem #6, 2026-05-29.
    Делегирует в общий set_slide_title (с fallback на textbox)."""
    from kpi_renderer import set_slide_title
    return set_slide_title(slide, text, dark=dark)


def add_top_separator(slide, y=110, color=None):
    """Тонкая серая линия под заголовком (визуальный разделитель)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, px(35), px(y), px(1245), px(y)
    )
    conn.line.color.rgb = color if color is not None else SEPARATOR_GRAY
    conn.line.width = Emu(int(0.5 * 12700))


# Фирменная стрелка-декор Cloud.ru (SVG) — стрелка ↗ (диагональ вверх-вправо +
# уголок-наконечник в верхнем-правом). Хранится в brand/icons/, вставляется как
# редактируемый вектор (SVG в .pptx) — Problem #4, 2026-05-29.
DECOR_ARROW_SVG = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "brand", "icons", "brand_arrow.svg"
)
# URI расширения Microsoft для встраивания SVG в blip + namespace asvg.
_SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
_ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"


# Толщина линии декор-стрелки = 1pt (как обычные линии в PowerPoint),
# user-decision 2026-05-29. SVG масштабируется, поэтому stroke-width в SVG
# вычисляется так, чтобы при размере декора отрисовалось ровно 1pt.
DECOR_STROKE_PT = 1.0
# 1pt = 12700 EMU; наш px = 9525 EMU ⇒ display_pt = size_px * 9525/12700.
# stroke_svg = target_pt * 101(viewBox) / display_pt.
_PT_PER_PX = 9525.0 / 12700.0  # = 0.75


def _brand_arrow_svg_bytes(stroke_width_svg):
    """Берёт canonical brand_arrow.svg и проставляет нужный stroke-width
    (в единицах viewBox), чтобы при текущем размере линия была 1pt."""
    import re
    with open(DECOR_ARROW_SVG, encoding="utf-8") as f:
        svg = f.read()
    svg = re.sub(r'\s*stroke-width="[^"]*"', "", svg)  # убрать старые
    svg = re.sub(r'(stroke="[^"]*")',
                 r'\1 stroke-width="%.3f"' % stroke_width_svg, svg)
    return svg.encode("utf-8")


def _svg_bytes_to_png(svg_bytes, scale=8):
    """Растрирует SVG-байты → PNG (RGBA) через PyMuPDF — fallback-картинка для
    LibreOffice/старых вьюеров. PowerPoint рисует сам SVG."""
    import fitz
    doc = fitz.open(stream=svg_bytes, filetype="svg")
    page = doc[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=True)
    data = pix.tobytes("png")
    doc.close()
    return data


def add_svg_picture(slide, svg_bytes, x, y, w, h):
    """Вставляет SVG (байты) как РЕДАКТИРУЕМЫЙ вектор в слайд.

    Технически: picture с PNG-fallback (r:embed) + расширение <asvg:svgBlip>,
    указывающее на встроенную SVG-часть. В PowerPoint объект — настоящий SVG
    (можно «Преобразовать в фигуру» / перекрасить), в превью рисуется PNG.
    """
    from pptx.opc.package import Part
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    png_bytes = _svg_bytes_to_png(svg_bytes)
    pic = slide.shapes.add_picture(io.BytesIO(png_bytes), px(x), px(y), px(w), px(h))

    package = slide.part.package
    partname = package.next_partname("/ppt/media/image%d.svg")
    svg_part = Part(partname, "image/svg+xml", package, svg_bytes)
    rId = slide.part.relate_to(svg_part, RT.IMAGE)

    blip_fill = pic._element.find(qn("p:blipFill"))
    blip = blip_fill.find(qn("a:blip"))
    ext_lst = etree.SubElement(blip, qn("a:extLst"))
    ext = etree.SubElement(ext_lst, qn("a:ext"))
    ext.set("uri", _SVG_EXT_URI)
    svg_blip = etree.SubElement(ext, "{%s}svgBlip" % _ASVG_NS)
    svg_blip.set(qn("r:embed"), rId)
    return pic


def _draw_arrow_lines(slide, x, y, size, w_pt, color):
    """Рисует фирменную стрелку ↗ нативными линиями PowerPoint, возвращает список
    линий-фигур. Геометрия = brand_arrow.svg: уголок-наконечник ↗ (верхняя грань +
    правая грань) + диагональ-древко из нижнего-левого к наконечнику.

    Цвет бренд-стрелок — зелёный ИЛИ серый, обе опции валидны (user 2026-06-02).
    Дефолт — зелёный (canon §8); передай color=ARROW_COLOR для серого варианта."""
    c = color if color is not None else GREEN
    top = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, px(x), px(y), px(x + size), px(y))
    right = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, px(x + size), px(y), px(x + size), px(y + size))
    shaft = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px(x), px(y + size), px(x + int(size * 0.9)), px(y + int(size * 0.1)))
    for ln in (top, right, shaft):
        ln.line.color.rgb = c
        ln.line.width = Emu(int(w_pt * 12700))
        _no_effects(ln)   # коннекторы тоже без эффектов (плоско, без теней)
    return [top, right, shaft]


def _next_shape_id(slide):
    """Следующий свободный shape id на слайде (для группы)."""
    ids = [1]
    for el in slide.shapes._spTree.iter(qn("p:cNvPr")):
        v = el.get("id")
        if v and v.isdigit():
            ids.append(int(v))
    return max(ids) + 1


def _group_shapes(slide, shape_list, x, y, w, h, name="BrandArrow"):
    """Оборачивает фигуры в группу (p:grpSp) с identity-трансформом — стрелка
    становится ОДНОЙ редактируемой фигурой (как результат «Преобразовать в фигуру»)."""
    spTree = slide.shapes._spTree
    grp = etree.SubElement(spTree, qn("p:grpSp"))
    nv = etree.SubElement(grp, qn("p:nvGrpSpPr"))
    cNvPr = etree.SubElement(nv, qn("p:cNvPr"))
    cNvPr.set("id", str(_next_shape_id(slide)))
    cNvPr.set("name", name)
    etree.SubElement(nv, qn("p:cNvGrpSpPr"))
    etree.SubElement(nv, qn("p:nvPr"))
    grp_pr = etree.SubElement(grp, qn("p:grpSpPr"))
    xfrm = etree.SubElement(grp_pr, qn("a:xfrm"))
    ox, oy, cx, cy = int(px(x)), int(px(y)), int(px(w)), int(px(h))
    off = etree.SubElement(xfrm, qn("a:off")); off.set("x", str(ox)); off.set("y", str(oy))
    ext = etree.SubElement(xfrm, qn("a:ext")); ext.set("cx", str(cx)); ext.set("cy", str(cy))
    ch_off = etree.SubElement(xfrm, qn("a:chOff")); ch_off.set("x", str(ox)); ch_off.set("y", str(oy))
    ch_ext = etree.SubElement(xfrm, qn("a:chExt")); ch_ext.set("cx", str(cx)); ch_ext.set("cy", str(cy))
    for shp in shape_list:
        grp.append(shp._element)  # reparent линии внутрь группы
    return grp


def add_decor_diagonals(slide,
                        count=4, x_start=20, y_start=620,
                        size=70, gap=14, w_pt=DECOR_STROKE_PT, color=None):
    """Фирменный декор Cloud.ru — ряд зелёных стрелок ↗.

    Canonical 2026-05-29 (Problem #4, обновлено): каждая стрелка = НАТИВНАЯ
    редактируемая ФИГУРА PowerPoint (группа линий), а не картинка — её сразу
    можно двигать/перекрашивать/реформировать, без «Преобразовать в фигуру».
    Древко ↗ + уголок-наконечник (геометрия brand_arrow.svg), линия 1pt.

    Ставится в bottom-left или bottom-right (свободный угол).
    """
    for i in range(count):
        x = x_start + i * (size + gap)
        y = y_start
        lines = _draw_arrow_lines(slide, x, y, size, w_pt, color)
        _group_shapes(slide, lines, x, y, size, size, name="BrandArrow")


# ============================================================================
# Высокоуровневая функция — сборка схемы по конфигу
# ============================================================================
def _block_anchor(block, side):
    """Точка на границе блока по стороне (right|left|top|bottom|center)."""
    x, y, w, h = block["x"], block["y"], block["w"], block["h"]
    if side == "right":
        return x + w, y + h // 2
    if side == "left":
        return x, y + h // 2
    if side == "top":
        return x + w // 2, y
    if side == "bottom":
        return x + w // 2, y + h
    # center
    return x + w // 2, y + h // 2


def _wrapped_lines(text, usable_w_px, char_w_px):
    """Сколько визуальных строк займёт логическая строка при переносе по ширине."""
    if not text:
        return 1
    cpl = max(1, int(usable_w_px / char_w_px))
    return max(1, -(-len(str(text)) // cpl))  # ceil


def compose_grid(blocks, area, cols=None, font_pt=16, gap=24,
                 pad=12, pad_bottom=16, v_center=True):
    """Грид-композиция блоков схемы (frame-to-text + сетка + единый кегль).

    Каждый блок задаётся ЛОГИЧЕСКИ: row, col, lines, fill — без пиксельных
    координат. Рендерер сам считает раскладку:
      - колонки равной ширины, единые зазоры (внутри area);
      - высота блока = под его текст (число строк с учётом переноса × межстрочный
        при ЕДИНОМ кегле + поля); высота строки = максимум по строке (выравнивание);
      - все блоки одной строки имеют общий верх и высоту; одной колонки — общий
        left/width. Кегль единый везде.
    area=(left, top, right, bottom) в px. Мутирует блоки (x,y,w,h,font_sizes).
    """
    if not blocks:
        return blocks
    left, top, right, bottom = area
    n_cols = cols or (max(b.get("col", 0) for b in blocks) + 1)
    n_rows = max(b.get("row", 0) for b in blocks) + 1
    avail_w = right - left
    col_w = int((avail_w - (n_cols - 1) * gap) / max(1, n_cols))
    font_px = font_pt * 4.0 / 3.0
    char_w = 0.58 * font_px      # консервативно (чуть шире) — чтобы не переполнить
    line_h = 1.30 * font_px
    usable_w = max(1, col_w - 2 * pad)
    # требуемая высота каждого блока (под текст)
    for b in blocks:
        lines = b.get("lines", []) or []
        total = sum(_wrapped_lines(ln, usable_w, char_w) for ln in lines) if lines else 1
        b["_need_h"] = int(total * line_h + pad + pad_bottom)
    # высота строки = максимум по строке
    row_h = [0] * n_rows
    for b in blocks:
        r = b.get("row", 0)
        row_h[r] = max(row_h[r], b["_need_h"])
    grid_h = sum(row_h) + (n_rows - 1) * gap
    avail_h = bottom - top
    start_y = top + (max(0, (avail_h - grid_h) // 2) if v_center else 0)
    row_y = []
    y = start_y
    for r in range(n_rows):
        row_y.append(y)
        y += row_h[r] + gap
    for b in blocks:
        r = b.get("row", 0)
        c = b.get("col", 0)
        b["x"] = left + c * (col_w + gap)
        b["y"] = row_y[r]
        b["w"] = col_w
        b["h"] = row_h[r]
        b["font_sizes"] = [font_pt] * len(b.get("lines", []))
    return blocks


# ============================================================================
# Пресеты дизайн-архетипов (design-principles-from-decks.md, раздел B).
# Каждый — самодостаточная композиция тела слайда внутри safe-area.
# Заголовок/subtitle рисует render_flow_diagram_slide ДО вызова пресета.
# ============================================================================
def _panel_rect(slide, x, y, w, h, dark=False, fill=None):
    """Простой залитый прямоугольник-фон (карточка/лента) без текста."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    rect.fill.solid()
    if fill is not None:
        rect.fill.fore_color.rgb = fill if not isinstance(fill, str) else _hex(fill)
    else:
        rect.fill.fore_color.rgb = DARK_FILL if dark else GRAY
    rect.line.fill.background()
    _no_effects(rect)
    return rect


def render_numbered_rows(slide, cfg, dark=False):
    """Архетип 5B — нумерованные строки. Сверено по референсу 2026-06-01
    (ПОЧЕМУ CLOUD.RU ЛУЧШЕ): чистый стиль — номер «01» графитом REGULAR + заголовок
    SemiBold графит в ОДНУ строку, тело TEXT_GRAY ниже. Без серых лент и зелёных
    чипов. Раскладка в N колонок (по умолчанию 1; 6–8 пунктов → 2 колонки).

    cfg["rows"]: [{"title","text"?,"number"?}], cfg["cols"] (default 1).
    cfg["style"]="band" → старый вариант: серые ленты во всю ширину + зелёный
        чип-номер слева (для коротких списков-чеклистов).
    """
    rows = cfg.get("rows", [])
    n = len(rows)
    if n == 0:
        return
    top = cfg.get("content_top", SAFE_TOP)
    avail = SAFE_BOTTOM - top
    title_col = WHITE if dark else GRAPHITE
    body_col = RGBColor(0xCF, 0xCF, 0xCF) if dark else TEXT_GRAY

    if cfg.get("style") == "band":
        gap = cfg.get("gap", PRESET_GAP)
        band_h = int((avail - (n - 1) * gap) / n)
        chip = min(cfg.get("chip_size", 54), band_h - 20)
        chip_x = SAFE_LEFT + cfg.get("chip_inset", 20)
        for i, row in enumerate(rows):
            y = top + i * (band_h + gap)
            _panel_rect(slide, SAFE_LEFT, y, SAFE_W, band_h, dark=dark)
            chip_y = y + (band_h - chip) // 2
            add_green_chip(slide, chip_x, chip_y, size=chip,
                           number=row.get("number", "%02d" % (i + 1)),
                           font_size=cfg.get("chip_font", 22))
            tx = chip_x + chip + cfg.get("chip_text_gap", 20)
            tw = SAFE_RIGHT - tx - 16
            title = row.get("title", "")
            text = row.get("text", "")
            lines = [l for l in (title, text) if l]
            sizes, bolds, cols_c = [], [], []
            if title:
                sizes.append(cfg.get("title_size", 24)); bolds.append(True)
                cols_c.append(title_col)
            if text:
                sizes.append(cfg.get("text_size", 16)); bolds.append(False)
                cols_c.append(body_col)
            add_block(slide, tx, y, tw, band_h, lines,
                      font_sizes=sizes, bolds=bolds,
                      fill=("dark" if dark else "gray"),
                      vanchor="middle", text_colors=cols_c)
        return

    # Чистый стиль (default) — текст с графитовым номером, без плашек.
    cols = cfg.get("cols", 1)
    per_col = -(-n // cols)
    col_gap = cfg.get("col_gap", 44)
    cw = int((SAFE_W - (cols - 1) * col_gap) / cols)
    slot_h = int(avail / per_col)
    num_w = cfg.get("number_w", 44)
    num_size = cfg.get("number_size", 17)
    title_size = cfg.get("title_size", 17)
    text_size = cfg.get("text_size", 14)
    for i, row in enumerate(rows):
        col = i // per_col
        rr = i % per_col
        x = SAFE_LEFT + col * (cw + col_gap)
        y = top + rr * slot_h
        title = row.get("title", "")
        text = row.get("text", "")
        number = str(row.get("number", "%02d" % (i + 1)))
        # Номер графитом REGULAR.
        add_label(slide, x, y, num_w, slot_h, number,
                  font_size=num_size, bold=False, anchor="top", color=title_col)
        tx = x + num_w
        tw = cw - num_w
        if title:
            add_label(slide, tx, y, tw, int(1.4 * title_size * 4 / 3), title,
                      font_size=title_size, bold=True, anchor="top", color=title_col)
        if text:
            ty = y + int(1.5 * title_size * 4 / 3)
            add_label(slide, tx, ty, tw, slot_h - (ty - y), text,
                      font_size=text_size, bold=False, anchor="top", color=body_col)


def render_card_grid(slide, cfg, dark=False):
    """Архетип 4 — сетка карточек: карточки в сетке cols×rows, заголовок (SemiBold,
    графит) + тело (TEXT_GRAY). Сверено по референсам 2026-06-01 (Evolution AI
    Factory, AI Workflows): 3-уровневая иерархия — заголовок графит, тело серое.

    Чип-номер/галочка ОПЦИОНАЛЕН (card["check"] или card["number"]): тогда чип
    слева, заголовок В ТУ ЖЕ СТРОКУ справа от чипа, тело — ниже под заголовком.
    Без чипа — заголовок сверху, тело под ним (как в большинстве эталонных карточек).

    cfg["cards"]: [{"title","text","check"?,"number"?}], cfg["cols"] (default 2).
    """
    cards = cfg.get("cards", [])
    n = len(cards)
    if n == 0:
        return
    cols = cfg.get("cols", 2)
    n_rows = -(-n // cols)
    gap = cfg.get("gap", PRESET_GAP)          # ≤10px, базово 4px
    top = cfg.get("content_top", SAFE_TOP)
    avail_h = SAFE_BOTTOM - top
    cw = int((SAFE_W - (cols - 1) * gap) / cols)
    ch = int((avail_h - (n_rows - 1) * gap) / n_rows)
    pad = cfg.get("pad", 22)
    chip = cfg.get("chip_size", 44)
    title_size = cfg.get("title_size", 20)
    text_size = cfg.get("text_size", 15)
    title_col = WHITE if dark else GRAPHITE
    # На тёмной плашке тело — светло-серое; на светлой — TEXT_GRAY.
    body_col = RGBColor(0xCF, 0xCF, 0xCF) if dark else TEXT_GRAY
    transparent = cfg.get("transparent", False)   # карточки без заливки (как 17.24.38)
    for i, card in enumerate(cards):
        r, c = i // cols, i % cols
        x = SAFE_LEFT + c * (cw + gap)
        y = top + r * (ch + gap)
        if not transparent:
            _panel_rect(slide, x, y, cw, ch, dark=dark)
        title = card.get("title", "")
        text = card.get("text", "")
        num = card.get("number")
        has_check = bool(card.get("check", False))
        has_chip = has_check or (num is not None)
        if has_chip:
            use_check = has_check or num is None
            add_green_chip(slide, x + pad, y + pad, size=chip,
                           number=(None if use_check else num),
                           check=use_check,
                           font_size=cfg.get("chip_font", 18))
            tx = x + pad + chip + 14
            tw = cw - (tx - x) - pad
            # Заголовок в строку с чипом — по центру высоты чипа.
            if title:
                add_label(slide, tx, y + pad, tw, chip, title,
                          font_size=title_size, bold=True,
                          anchor="middle", color=title_col)
            body_y = y + pad + chip + 12
            body_x, body_w = tx, tw
        else:
            tx = x + pad
            tw = cw - 2 * pad
            # высота заголовка под число строк (single/two-line)
            t_lines = _wrapped_lines(title, tw, 0.58 * title_size * 4.0 / 3.0)
            t_h = int(t_lines * 1.25 * title_size * 4.0 / 3.0) + 4
            if title:
                add_label(slide, tx, y + pad, tw, t_h, title,
                          font_size=title_size, bold=True,
                          anchor="top", color=title_col)
            body_y = y + pad + t_h + 6
            body_x, body_w = tx, tw
        if text:
            add_label(slide, body_x, body_y, body_w,
                      y + ch - body_y - pad, text,
                      font_size=text_size, bold=False,
                      anchor="top", color=body_col)


def render_numbered_columns(slide, cfg, dark=False):
    """Архетип 5A — открытые нумерованные колонки: без заливки, в каждой колонке
    зелёная черта сверху, заголовок + текст, внизу крупный зелёный номер (~40pt).

    cfg["columns"]: [{"title","text","number"}]
    """
    cols_data = cfg.get("columns", [])
    n = len(cols_data)
    if n == 0:
        return
    gap = cfg.get("gap", 32)
    top = cfg.get("content_top", SAFE_TOP)
    cw = int((SAFE_W - (n - 1) * gap) / n)
    num_h = 90
    num_y = SAFE_BOTTOM - num_h
    title_col = WHITE if dark else GRAPHITE
    body_col = RGBColor(0xCF, 0xCF, 0xCF) if dark else TEXT_GRAY
    num_color = _hex(cfg["number_color"]) if cfg.get("number_color") else GREEN
    for i, col in enumerate(cols_data):
        x = SAFE_LEFT + i * (cw + gap)
        # Черта сверху ОТКЛЮЧЕНА по умолчанию (сверено по референсам: открытые
        # колонки без декоративной линии). Если включить cfg["rule"]=True — зелёная
        # (canon §5A): она привязана к колонке и держит композицию — допустимо.
        if cfg.get("rule", False):
            rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          px(x), px(top), px(cw), px(3))
            rule.fill.solid(); rule.fill.fore_color.rgb = GREEN
            rule.line.fill.background(); _no_effects(rule)
        title = col.get("title", "")
        text = col.get("text", "")
        lines = [l for l in (title, text) if l]
        sizes, bolds, cols_c = [], [], []
        if title:
            sizes.append(cfg.get("title_size", 20)); bolds.append(True)
            cols_c.append(title_col)
        if text:
            sizes.append(cfg.get("text_size", 15)); bolds.append(False)
            cols_c.append(body_col)
        add_block(slide, x, top + 16, cw, num_y - (top + 16) - 12, lines,
                  font_sizes=sizes, bolds=bolds,
                  fill=("dark" if dark else "white"),
                  vanchor="top", text_colors=cols_c)
        # Крупный номер внизу — REGULAR-начертание (сверено: не bold), крупнее.
        add_label(slide, x, num_y, cw, num_h,
                  str(col.get("number", "%02d" % (i + 1))),
                  font_size=cfg.get("number_size", 56), bold=False,
                  align="left", anchor="middle", color=num_color)


def render_hero_statement(slide, cfg, dark=False):
    """Архетип 6 — hero-утверждение: крупный текст (~54pt CAPS) на зелёной плашке
    + 2 смещённые контурные рамки (декор) + поддерживающий текст снизу-справа.

    cfg["statement"] (обязат.), cfg["support"] (опц.)
    """
    statement = cfg.get("statement", "")
    support = cfg.get("support", "")
    bx, by = SAFE_LEFT, cfg.get("content_top", SAFE_TOP) + 30
    bw = int(SAFE_W * 0.72)
    # Плашка «обнимает» текст: высота считается под число строк. Если не влезает
    # по высоте — уменьшаем кегль (до 32pt), чтобы текст не вытекал за зелёный блок.
    margin = 32
    inner_w = bw - 2 * margin
    block_bottom_max = cfg.get("block_bottom_max", 545)
    fs = cfg.get("statement_size", 54)
    text_caps = statement.upper()
    while True:
        font_px = fs * 4.0 / 3.0
        line_h = 1.30 * font_px
        char_w = 0.66 * font_px           # CAPS-буквы шире обычных
        n_lines = _wrapped_lines(text_caps, inner_w, char_w)
        bh = int(n_lines * line_h + 2 * 24)   # + верх/низ поля по 24
        if by + bh <= block_bottom_max or fs <= 32:
            break
        fs -= 4
    # Смещённые контурные рамки (рисуем ПЕРЕД плашкой — выглядывают снизу-справа).
    # Зелёные outline-рамки: привязаны к hero-блоку (offset-композиция, blueprint-
    # глубина) — это КОМПОЗИЦИЯ, не «недолиния» → зелёный допустим (user 2026-06-02).
    # (Можно переключить на серый — обе опции валидны.)
    for off in (24, 12):
        fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    px(bx + off), px(by + off), px(bw), px(bh))
        fr.fill.background()
        fr.line.color.rgb = GREEN
        fr.line.width = Emu(int(1.0 * 12700))
        _no_effects(fr)
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(bx), px(by), px(bw), px(bh))
    block.fill.solid(); block.fill.fore_color.rgb = GREEN
    block.line.fill.background(); _no_effects(block)
    tf = block.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(32 * EMU)
    tf.margin_right = Emu(32 * EMU)
    tf.margin_top = Emu(24 * EMU)
    tf.margin_bottom = Emu(24 * EMU)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = statement.upper()
    run.font.size = Pt(fs)
    _set_weight(run.font, True)
    run.font.color.rgb = GRAPHITE
    if support:
        sx = bx + bw - 360
        add_label(slide, sx, by + bh + 36, SAFE_RIGHT - sx, 90, support,
                  font_size=cfg.get("support_size", 16), bold=False,
                  align="right", anchor="top",
                  color=WHITE if dark else GRAPHITE)


_PRESETS = {
    "numbered_rows": render_numbered_rows,
    "numbered_columns": render_numbered_columns,
    "card_grid": render_card_grid,
    "hero_statement": render_hero_statement,
}


def render_flow_diagram_slide(slide, flow_config, dark=False):
    """Собирает flow-схему на готовом (очищенном) blank slide.

    Предполагается, что slide уже прошёл clean_slide_to_blank() — этим
    занимается build_v9 до вызова render_flow_diagram_slide.

    Режим grid (flow_config["grid"]=true): блоки задаются логически (row/col/lines),
    рендерер считает раскладку через compose_grid — frame-to-text + сетка + единый кегль.
    """
    if not isinstance(flow_config, dict):
        raise ValueError("flow_config должен быть dict")

    # 1. Header (заголовок ВСЕГДА в safe space через TITLE-placeholder @35,38).
    header = flow_config.get("header", "")
    if header:
        add_header(slide, header, dark=dark)
        # Дивайдер под заголовком БОЛЬШЕ НЕ рисуем по умолчанию (правило 2026-06-01):
        # он не нужен между заголовком и контентом. Включается ТОЛЬКО явным флагом
        # flow_config["top_separator"]=true — например, чтобы отделять абзацы
        # друг от друга в определённых типах слайда.
        if flow_config.get("top_separator"):
            add_top_separator(slide)

    # 1a. Subtitle / URL
    subtitle = flow_config.get("subtitle")
    if subtitle:
        add_label(slide, 35, 120, 1000, 24, subtitle,
                  font_size=13, bold=False, align="left",
                  color=WHITE if dark else GRAPHITE)
    subtitle_url = flow_config.get("subtitle_url")
    if subtitle_url:
        add_label(slide, 35, 144, 1000, 20, subtitle_url,
                  font_size=12, bold=False, align="left",
                  color=DASH_GRAY)

    # 1b. Отступ тела под подзаголовок: при наличии subtitle/URL контент пресета
    #     стартует ниже, чтобы не прижимался к подзаголовку (правило 2026-06-01).
    #     Пользователь может задать content_top явно — тогда не трогаем.
    if "content_top" not in flow_config:
        if subtitle_url:
            flow_config["content_top"] = 190
        elif subtitle:
            flow_config["content_top"] = 172
        else:
            flow_config["content_top"] = SAFE_TOP

    # 1c. Preset-архетипы (готовые композиции). Если задан preset — рисуем тело
    #     слайда соответствующей функцией и выходим (header/subtitle уже выше).
    preset = flow_config.get("preset")
    if preset:
        fn = _PRESETS.get(preset)
        if fn is None:
            raise ValueError(
                "Неизвестный preset '%s'. Доступны: %s"
                % (preset, ", ".join(sorted(_PRESETS)))
            )
        fn(slide, flow_config, dark=dark)
        decor = flow_config.get("decor")
        if decor and decor.get("enabled"):
            add_decor_diagonals(
                slide,
                count=decor.get("count", 4),
                x_start=decor.get("x_start", 20),
                y_start=decor.get("y_start", 620),
                size=decor.get("size", 70),
                gap=decor.get("gap", 14),
                w_pt=decor.get("w_pt", 1.4),
            )
        return

    # 2. Blocks — собираем по id для arrow refs
    blocks = flow_config.get("blocks", [])
    # Grid-режим: считаем раскладку (frame-to-text + сетка + единый кегль).
    if flow_config.get("grid"):
        grid_top = 170 if (subtitle or subtitle_url) else 150
        # Микроправило: если есть branching-стрелки (ветка между разными строкой И
        # колонкой) — зазор колонок больше, чтобы веткам/наконечникам хватало места.
        def _is_branch(a):
            f, t = a.get("from"), a.get("to")
            return (isinstance(f, (list, tuple)) and isinstance(t, (list, tuple))
                    and f[0] != t[0] and f[1] != t[1])
        has_branch = any(_is_branch(a) for a in flow_config.get("arrows", []))
        default_gap = GRID_GAP_BRANCHING if has_branch else GRID_GAP_DEFAULT
        compose_grid(
            blocks, (SAFE_LEFT, grid_top, SAFE_RIGHT, SAFE_BOTTOM),
            cols=flow_config.get("cols"),
            font_pt=flow_config.get("font_size", 16),
            gap=flow_config.get("gap", default_gap),
        )
        for b in blocks:
            if "id" not in b and "row" in b and "col" in b:
                b["id"] = "%s,%s" % (b["row"], b["col"])

    # 1d. Залитые панели-секции (style="panel") — ФОНОМ, до блоков (нагруженные схемы).
    #     Притягиваем края к направляющим safe-area (фреймы до конца, без выхода за поля).
    for grp in flow_config.get("groups", []):
        if grp.get("style") == "panel":
            sx, sy, sw, sh = snap_panel_to_safe(
                grp["x"], grp["y"], grp["w"], grp["h"])
            add_filled_panel(slide, sx, sy, sw, sh,
                             label=grp.get("label"), dark=dark, fill=grp.get("fill"))

    blocks_by_id = {}
    for blk in blocks:
        # Canonical default: align=left, vanchor=top (правило 2026-05-06).
        shape = add_block(
            slide,
            blk["x"], blk["y"], blk["w"], blk["h"],
            blk.get("lines", []),
            font_sizes=blk.get("font_sizes"),
            bolds=blk.get("bolds"),
            caps_first=blk.get("caps_first", False),
            fill=blk.get("fill", "gray"),
            align=blk.get("align", "left"),
            vanchor=blk.get("vanchor", "top"),
        )
        if "id" in blk:
            blocks_by_id[blk["id"]] = blk

    # 3. Groups — пунктирная рамка + label (style="panel" уже нарисованы фоном).
    for grp in flow_config.get("groups", []):
        if grp.get("style") == "panel":
            continue
        gx, gy, gw, gh = snap_panel_to_safe(
            grp["x"], grp["y"], grp["w"], grp["h"])
        add_dashed_rect(slide, gx, gy, gw, gh)
        label_text = grp.get("label")
        if label_text:
            label_pos = grp.get("label_pos", "top")
            label_y = gy + 2 if label_pos == "top" else gy + gh - 20
            add_label(slide, gx, label_y, gw, 18, label_text,
                      font_size=12, bold=True, align="center")

    # 4. Arrows
    for arr in flow_config.get("arrows", []):
        if "from" in arr and "to" in arr:
            fr, to = arr["from"], arr["to"]
            if isinstance(fr, (list, tuple)):
                fr = "%s,%s" % (fr[0], fr[1])   # ссылка на ячейку [row,col]
            if isinstance(to, (list, tuple)):
                to = "%s,%s" % (to[0], to[1])
            src = blocks_by_id.get(fr)
            dst = blocks_by_id.get(to)
            if src is None or dst is None:
                continue
            # Авто-определение стороны выхода стрелки по позициям ячеек грида
            # (чтобы вертикальные связи не рисовались как горизонтальные и не «пропадали»).
            side = arr.get("side")
            if side is None and "row" in src and "row" in dst:
                if src["row"] == dst["row"]:
                    side = "right" if dst["col"] >= src["col"] else "left"
                elif src["col"] == dst["col"]:
                    side = "bottom" if dst["row"] >= src["row"] else "top"
            if side is None:
                side = "right"
            x1, y1 = _block_anchor(src, side)
            # По умолчанию входим в противоположную сторону
            opposite = {"right": "left", "left": "right",
                        "top": "bottom", "bottom": "top"}.get(side, "left")
            x2, y2 = _block_anchor(dst, arr.get("to_side", opposite))
        else:
            x1, y1 = arr["x1"], arr["y1"]
            x2, y2 = arr["x2"], arr["y2"]
        # Диагональ (ветка между разными строкой И колонкой) → ортогональный
        # Z-маршрут вместо запрещённой диагонали. Иначе — прямая.
        if x1 != x2 and y1 != y2:
            add_orthogonal_arrow(slide, x1, y1, x2, y2, w_pt=arr.get("w_pt"))
        else:
            add_arrow(
                slide, x1, y1, x2, y2,
                with_head=arr.get("with_head", True),
                w_pt=arr.get("w_pt"),
                dashed=arr.get("dashed", False),
            )

    # 5. Labels (произвольные)
    for lab in flow_config.get("labels", []):
        add_label(
            slide,
            lab["x"], lab["y"], lab["w"], lab["h"],
            lab["text"],
            font_size=lab.get("font_size", 11),
            bold=lab.get("bold", False),
            align=lab.get("align", "left"),
            caps=lab.get("caps", False),
        )

    # 6. Decor (зелёные уголки)
    decor = flow_config.get("decor")
    if decor and decor.get("enabled"):
        add_decor_diagonals(
            slide,
            count=decor.get("count", 4),
            x_start=decor.get("x_start", 20),
            y_start=decor.get("y_start", 620),
            size=decor.get("size", 70),
            gap=decor.get("gap", 14),
            w_pt=decor.get("w_pt", 1.4),
        )


# ============================================================================
# CLI для standalone-теста
# ============================================================================
def _cli_main():
    """python3 flow_renderer.py <flow_config.json> <template.pptx> <out.pptx>"""
    import sys
    from pptx import Presentation

    if len(sys.argv) < 4:
        print("Usage: flow_renderer.py <flow_config.json> <template.pptx> <out.pptx>",
              file=sys.stderr)
        sys.exit(1)

    cfg_path, tpl_path, out_path = sys.argv[1], sys.argv[2], sys.argv[3]
    with open(cfg_path, encoding="utf-8") as f:
        cfg = json.load(f)

    # Импорт через relative path (для standalone)
    here = os.path.dirname(os.path.abspath(__file__))
    import sys as _sys
    _sys.path.insert(0, here)
    from kpi_renderer import clean_slide_to_blank, BLANK_DONOR_WHITE, BLANK_DONOR_DARK
    from build_v9 import clone_slide

    dark = cfg.get("dark", False)
    flow = cfg.get("flow", cfg)  # допустим, если передан "flat" config

    p = Presentation(tpl_path)
    originals = list(p.slides)
    blank_idx = BLANK_DONOR_DARK if dark else BLANK_DONOR_WHITE
    new_slide = clone_slide(p, originals[blank_idx - 1])

    # удалить оригиналы
    sldIdLst = p.slides._sldIdLst
    for sid in list(sldIdLst)[:len(originals)]:
        rid = sid.attrib[qn("r:id")]
        try:
            p.part.drop_rel(rid)
        except Exception:
            pass
        sldIdLst.remove(sid)

    clean_slide_to_blank(new_slide)
    render_flow_diagram_slide(new_slide, flow, dark=dark)
    p.save(out_path)
    print(f"Saved {out_path}", file=sys.stderr)


if __name__ == "__main__":
    _cli_main()
