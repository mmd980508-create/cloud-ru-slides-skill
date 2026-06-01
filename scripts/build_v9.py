#!/usr/bin/env python3
"""
build_v7.py — build_v6 + поддержка ДУБЛЕЙ donor.

Ключевое улучшение vs v6:
- Если donor 13 нужен 3 раза → клонируется 3 раза (не один и не игнорируется)
- Через XML deepcopy slide part-а в presentation
- Сохраняет правильный порядок слайдов

Plan:
{
  "slides": [
    {"clone_from_slide": 13, "slots": {...}},
    {"clone_from_slide": 12, "slots": {...}},
    {"clone_from_slide": 13, "slots": {...}}  ← ДУБЛЬ donor 13!
  ]
}

Usage:
    python3 build_v7.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]
"""
import sys
import json
import os
import copy
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from build_v5 import (
    load_donor_map, get_text_frame_by_shape_idx, replace_text_with_style,
    clear_text_frame
)
from kpi_renderer import (
    render_kpi, clean_slide_to_blank,
    BLANK_DONOR_WHITE, BLANK_DONOR_DARK
)
from image_renderer import render_image_native
try:
    from chart_engine import render_chart
    CHART_AVAILABLE = True
except ImportError:
    CHART_AVAILABLE = False

try:
    from chart_native_pptx import render_chart_pptx_slide
    CHART_NATIVE_PPTX_AVAILABLE = True
except ImportError:
    CHART_NATIVE_PPTX_AVAILABLE = False

try:
    from flow_renderer import render_flow_diagram_slide
    FLOW_RENDERER_AVAILABLE = True
except ImportError:
    FLOW_RENDERER_AVAILABLE = False

try:
    from table_renderer import render_table_native
    TABLE_RENDERER_AVAILABLE = True
except ImportError:
    TABLE_RENDERER_AVAILABLE = False

EMU_PER_PX = 9525


def clone_slide(prs, src_slide):
    """Глубоко копирует slide-part и регистрирует его в presentation.
    Возвращает новый Slide (последний в prs.slides)."""
    from pptx.opc.constants import CONTENT_TYPE as CT
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.packuri import PackURI
    from pptx.parts.slide import SlidePart

    src_part = src_slide.part
    src_xml = src_part.blob

    # Подбираем уникальное имя slideN.xml
    package = prs.part.package
    existing_partnames = {str(p.partname) for p in package.iter_parts()}
    next_idx = 1
    while f"/ppt/slides/slide{next_idx}.xml" in existing_partnames:
        next_idx += 1
    new_partname = PackURI(f"/ppt/slides/slide{next_idx}.xml")

    # Создаём новый part — используем тот же content_type как у src
    new_part = SlidePart.load(
        partname=new_partname,
        content_type=src_part.content_type,
        blob=src_xml,
        package=package,
    )

    # Копируем relationships от source slide в new slide
    for rel in src_part.rels.values():
        if rel.is_external:
            new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_part.relate_to(rel.target_part, rel.reltype)

    # Регистрируем slide в presentation через relationship
    rId = prs.part.relate_to(new_part, RT.SLIDE)

    # Добавляем sldId в sldIdLst
    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(el.attrib["id"]) for el in sldIdLst if "id" in el.attrib]
    next_id = max(existing_ids) + 1 if existing_ids else 256
    new_sldId = etree.SubElement(sldIdLst, qn("p:sldId"))
    new_sldId.set("id", str(next_id))
    new_sldId.set(qn("r:id"), rId)

    return prs.slides[-1]


def build(plan_path, template_path, output_path, donor_map_path):
    plan = json.load(open(plan_path, encoding="utf-8"))
    p = Presentation(template_path)
    donors = load_donor_map(donor_map_path)

    # === STEP 1: Собираем нужные donor_nums и клонируем все слайды plan-а ===
    # Сохраняем reference на оригинальные donor slides ДО любых модификаций
    original_slides = list(p.slides)
    donor_originals = {}  # {donor_num: original_slide}
    for ps in plan["slides"]:
        n = ps.get("clone_from_slide")
        if n and n not in donor_originals:
            if 1 <= n <= len(original_slides):
                donor_originals[n] = original_slides[n - 1]
            else:
                print(f"WARN: donor {n} вне диапазона (1..{len(original_slides)})", file=sys.stderr)

    # Клонируем КАЖДЫЙ слайд из plan (включая дубли) → новые slides в конце
    # Для slide_type=="kpi_native": используем blank donor (slide 30/22 шаблона)
    cloned_for_plan = []
    for ps in plan["slides"]:
        slide_type = ps.get("slide_type")
        if slide_type in ("kpi_native", "image_native", "chart_native", "chart_pptx_native", "flow_diagram_native", "table_native"):
            dark = ps.get("dark", False)
            blank_idx = (BLANK_DONOR_DARK if dark else BLANK_DONOR_WHITE)
            if 1 <= blank_idx <= len(original_slides):
                new_slide = clone_slide(p, original_slides[blank_idx - 1])
                cloned_for_plan.append(new_slide)
                continue
        n = ps.get("clone_from_slide")
        if not n or n not in donor_originals:
            cloned_for_plan.append(None)
            continue
        new_slide = clone_slide(p, donor_originals[n])
        cloned_for_plan.append(new_slide)

    # === STEP 2: Удаляем все ОРИГИНАЛЬНЫЕ слайды (101+ template слайдов), оставляем только клоны ===
    sldIdLst = p.slides._sldIdLst
    n_originals = len(original_slides)
    # Первые n_originals элементов — это оригиналы. Удаляем их.
    all_sldIds = list(sldIdLst)
    for sldId in all_sldIds[:n_originals]:
        rId = sldId.attrib[qn('r:id')]
        try:
            p.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)

    # === STEP 3: Заполняем text + pictures для каждого clone ===
    pictures_inserted = 0
    for plan_slide, actual in zip(plan["slides"], cloned_for_plan):
        if actual is None:
            continue

        # === NATIVE RENDERS: build shapes from scratch on clean canvas ===
        slide_type = plan_slide.get("slide_type")
        if slide_type == "kpi_native":
            kpi_config = plan_slide.get("kpi", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_kpi(actual, kpi_config, dark=dark)
            continue
        if slide_type == "image_native":
            image_config = plan_slide.get("image", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_image_native(actual, image_config, dark=dark)
            continue
        if slide_type == "chart_pptx_native":
            if not CHART_NATIVE_PPTX_AVAILABLE:
                print("WARN: chart_native_pptx модуль недоступен — chart_pptx_native пропущен",
                      file=sys.stderr)
                continue
            chart_config = plan_slide.get("chart", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_chart_pptx_slide(actual, chart_config, dark=dark)
            continue
        if slide_type == "flow_diagram_native":
            if not FLOW_RENDERER_AVAILABLE:
                print("WARN: flow_renderer модуль недоступен — flow_diagram_native пропущен",
                      file=sys.stderr)
                continue
            flow_config = plan_slide.get("flow", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_flow_diagram_slide(actual, flow_config, dark=dark)
            continue
        if slide_type == "table_native":
            if not TABLE_RENDERER_AVAILABLE:
                print("WARN: table_renderer модуль недоступен — table_native пропущен",
                      file=sys.stderr)
                continue
            table_config = plan_slide.get("table", {})
            dark = plan_slide.get("dark", False)
            clean_slide_to_blank(actual)
            render_table_native(actual, table_config, dark=dark)
            continue
        if slide_type == "chart_native":
            if not CHART_AVAILABLE:
                print("WARN: matplotlib не установлен — chart_native пропущен", file=sys.stderr)
                continue
            chart_config = plan_slide.get("chart", {})
            dark = plan_slide.get("dark", False)
            # Render chart to PNG
            chart_png = plan_slide.get("chart_output_png",
                                        f"pptx-skill/output/_chart_slide_{id(plan_slide)}.png")
            render_chart(chart_config, chart_png, dpi=150)
            # Pass to image_native renderer (wide_zone for charts)
            clean_slide_to_blank(actual)
            render_image_native(actual, {
                "title": chart_config.get("slide_title", chart_config.get("title", "")),
                "image_path": chart_png,
                "caption": chart_config.get("caption", "")
            }, dark=dark, wide_zone=True)
            continue

        src_num = plan_slide.get("clone_from_slide")
        if src_num is None:
            continue
        donor_def = donors.get(src_num)

        # === STEP 3a: PRE-CLEANUP (PNG-stripping) ===
        # Источники remove_idx:
        #   1. donor_def.remove_before_fill — всегда удалять
        #   2. plan_slide.remove_shapes — ad-hoc per slide
        #   3. donor_def.remove_if_not_used — удалять если slot пустой
        #      (формат: {slot_name: [shape_idx]} в slot.shape_idx_when_unused — упрощённо мапим)
        #   4. donor_def.remove_if_user_provides_table — удалять если plan имеет table_data
        if donor_def is not None:
            remove_idx_list = list(donor_def.get("remove_before_fill", []))
            remove_idx_list += list(plan_slide.get("remove_shapes", []))

            # remove_if_user_provides_table: например donor 53 имеет PNG-таблицу-заглушку
            if plan_slide.get("table_data"):
                remove_idx_list += list(donor_def.get("remove_if_user_provides_table", []))

            # remove_if_not_used: парсим формат {slot_name: shape_idx_to_strip}
            # Если slot не указан в plan_slide.slots — добавить shape_idx в remove
            remove_when_unused = donor_def.get("remove_if_not_used", {}) or {}
            slots_filled_now = plan_slide.get("slots", {}) or {}
            if isinstance(remove_when_unused, dict):
                for slot_name, idx_to_strip in remove_when_unused.items():
                    if slot_name not in slots_filled_now:
                        if isinstance(idx_to_strip, list):
                            remove_idx_list += idx_to_strip
                        else:
                            remove_idx_list.append(idx_to_strip)

            # WARN если donor_type=fixed_png_content и нет ни remove_before_fill, ни overrides
            dtype = donor_def.get("donor_type")
            if dtype == "fixed_png_content" and not remove_idx_list:
                print(
                    f"WARN: donor {src_num} is 'fixed_png_content' но без remove_before_fill — "
                    f"PNG-заглушка может перекрыть контент",
                    file=sys.stderr,
                )

            if remove_idx_list:
                spTree = actual.shapes._spTree
                shape_elements = list(spTree)
                content_tags = ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp')
                content_shapes = [el for el in shape_elements
                                  if el.tag.split('}')[-1] in content_tags]
                for idx in sorted(set(remove_idx_list), reverse=True):
                    if 0 <= idx < len(content_shapes):
                        spTree.remove(content_shapes[idx])

        # TEXT slots
        if donor_def is not None:
            slot_defs = donor_def.get("slots", {})
            slots_filled = plan_slide.get("slots", {})
            styles_override = plan_slide.get("slot_styles_override", {})

            for slot_name, new_text in slots_filled.items():
                if slot_name not in slot_defs:
                    print(f"WARN: slot '{slot_name}' undefined for donor {src_num}", file=sys.stderr)
                    continue
                shape_idx = slot_defs[slot_name]["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is None:
                    continue
                override = styles_override.get(slot_name)
                replace_text_with_style(tf, new_text, override)

            # Очистить незаполненные обязательные слоты
            for slot_name, slot_def in slot_defs.items():
                if slot_name in slots_filled:
                    continue
                if slot_def.get("optional"):
                    continue
                shape_idx = slot_def["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is not None:
                    clear_text_frame(tf)

        # PICTURES (вставляются ПОВЕРХ donor shapes)
        for pic in plan_slide.get("pictures", []):
            file_path = pic.get("file")
            if not file_path or not os.path.exists(file_path):
                print(f"WARN: image not found: {file_path}", file=sys.stderr)
                continue
            try:
                actual.shapes.add_picture(
                    file_path,
                    Emu(pic.get("left_px", 0) * EMU_PER_PX),
                    Emu(pic.get("top_px", 0) * EMU_PER_PX),
                    Emu(pic.get("width_px", 100) * EMU_PER_PX),
                    Emu(pic.get("height_px", 100) * EMU_PER_PX),
                )
                pictures_inserted += 1
            except Exception as e:
                print(f"WARN: insert_picture failed: {e}", file=sys.stderr)

        # TABLES (v8: fill_existing если donor имеет встроенную таблицу с брендовым стилем!)
        table_data = plan_slide.get("table_data")
        if table_data:
            try:
                # Найти существующую таблицу в donor (если есть)
                existing_table = None
                for sh in actual.shapes:
                    if sh.has_table:
                        existing_table = sh.table
                        break

                rows_needed = len(table_data)
                cols_needed = max(len(r) for r in table_data) if rows_needed else 1

                if existing_table:
                    # Donor уже имеет таблицу с брендовым стилем — заполняем её!
                    table_rows = len(existing_table.rows)
                    table_cols = len(existing_table.columns)
                    for r_idx, row_data in enumerate(table_data):
                        if r_idx >= table_rows:
                            break
                        for c_idx in range(table_cols):
                            cell = existing_table.cell(r_idx, c_idx)
                            if c_idx < len(row_data):
                                cell.text = str(row_data[c_idx])
                            else:
                                # Лишние колонки очищаем
                                cell.text = ""
                    # Очистить лишние строки если наш data короче
                    for r_idx in range(rows_needed, table_rows):
                        for c_idx in range(table_cols):
                            existing_table.cell(r_idx, c_idx).text = ""
                else:
                    # Donor не имеет таблицы — добавляем новую
                    left = Emu(35 * EMU_PER_PX)
                    top = Emu(120 * EMU_PER_PX)
                    width = Emu(1210 * EMU_PER_PX)
                    height = Emu(min(550, rows_needed * 50) * EMU_PER_PX)
                    tbl_shape = actual.shapes.add_table(rows_needed, cols_needed, left, top, width, height)
                    tbl = tbl_shape.table
                    for r_idx, row_data in enumerate(table_data):
                        for c_idx, cell_text in enumerate(row_data):
                            if c_idx >= cols_needed:
                                continue
                            tbl.cell(r_idx, c_idx).text = str(cell_text)
            except Exception as e:
                print(f"WARN: table fill failed: {e}", file=sys.stderr)

    # === FINAL: canonical enforcement над ВСЕМИ слайдами (для clone-based, где
    # native-фиксы не действуют). БЕЗОПАСНОЕ:
    #   - цвет: зелёный/белый текст → #222222 (кроме тёмного фона) [Problem #2]
    #   - вес: bold → SemiBold [Problem #3]
    #   - размер <12 → 12
    #   - заголовок контент-слайда → штатный TITLE-placeholder (35,38)/20pt
    #     SemiBold CAPS, СЕМАНТИЧЕСКИ (Вариант A) — не «угадывая по позиции»;
    #     титульные/divider и вертикальные/огромные заголовки не трогаются.
    # Bump до 16pt НЕ включаем — он даёт overflow на плотных/код-боксах. ===
    try:
        from enforce_canonical import enforce_canonical_slide, slide_is_dark
        enf_total = {}
        for slide in p.slides:
            st = enforce_canonical_slide(
                slide, dark=slide_is_dark(slide),
                min_pt=12, bump_from=None, bump_to=None, normalize_header=True)
            for k, v in st.items():
                enf_total[k] = enf_total.get(k, 0) + v
        if any(enf_total.values()):
            print(f"enforce_canonical: {enf_total}", file=sys.stderr)
    except Exception as e:
        print(f"WARN: enforce_canonical pass skipped: {e}", file=sys.stderr)

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides, {pictures_inserted} pictures inserted",
          file=sys.stderr)


def main():
    if len(sys.argv) < 4:
        print("Usage: build_v7.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]",
              file=sys.stderr)
        sys.exit(1)
    plan_p = sys.argv[1]
    tpl_p = sys.argv[2]
    out_p = sys.argv[3]
    donor_p = sys.argv[4] if len(sys.argv) > 4 else "pptx-skill/brand/donor-slot-map.yaml"
    build(plan_p, tpl_p, out_p, donor_p)


if __name__ == "__main__":
    main()
