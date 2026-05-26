#!/usr/bin/env python3
"""
build_v6.py — build_v5 + поддержка вставки изображений.

Plan расширен:
{
  "slides": [
    {
      "clone_from_slide": 34,
      "slots": {...},
      "slot_styles_override": {...},
      "pictures": [
        {
          "file": "path/to/image.png",
          "left_px": 35, "top_px": 272,
          "width_px": 100, "height_px": 100
        },
        ...
      ]
    }
  ]
}

При build_v6:
1. Сначала clone slide + replace text (как в v5)
2. Потом INSERT pictures через slide.shapes.add_picture()
3. Pictures появляются ПОВЕРХ существующих shapes донора

Usage:
    python3 build_v6.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]
"""
import sys
import json
import os
from pptx import Presentation
from pptx.util import Emu

# Reuse build_v5 functions
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from build_v5 import (
    load_donor_map, get_text_frame_by_shape_idx, replace_text_with_style,
    clear_text_frame, reorder_and_filter_slides
)

EMU_PER_PX = 9525  # at 96 DPI


def insert_picture(slide, file_path, left_px, top_px, width_px, height_px):
    """Вставить картинку на слайд по pixel-координатам."""
    if not os.path.exists(file_path):
        print(f"WARN: image file not found: {file_path}", file=sys.stderr)
        return False
    try:
        slide.shapes.add_picture(
            file_path,
            Emu(left_px * EMU_PER_PX),
            Emu(top_px * EMU_PER_PX),
            Emu(width_px * EMU_PER_PX),
            Emu(height_px * EMU_PER_PX),
        )
        return True
    except Exception as e:
        print(f"WARN: insert_picture failed for {file_path}: {e}", file=sys.stderr)
        return False


def build(plan_path, template_path, output_path, donor_map_path):
    plan = json.load(open(plan_path, encoding="utf-8"))
    p = Presentation(template_path)
    donors = load_donor_map(donor_map_path)

    cloned_nums = [s.get("clone_from_slide") for s in plan["slides"]]
    reorder_and_filter_slides(p, cloned_nums)

    # Список cloned slides в порядке plan (поддержка повторов donor!)
    cloned_slides_ordered = list(p.slides)

    pictures_inserted = 0
    plan_idx = 0
    for plan_slide in plan["slides"]:
        src_num = plan_slide.get("clone_from_slide")
        if not src_num:
            continue
        actual = cloned_slides_ordered[plan_idx]
        plan_idx += 1
        donor_def = donors.get(src_num)

        # === TEXT (как в v5) ===
        if donor_def is not None:
            slot_defs = donor_def.get("slots", {})
            slots_filled = plan_slide.get("slots", {})
            styles_override = plan_slide.get("slot_styles_override", {})

            for slot_name, new_text in slots_filled.items():
                if slot_name not in slot_defs:
                    print(f"WARN: slot '{slot_name}' undefined for donor {src_num}", file=sys.stderr)
                    continue
                shape_idx = slot_defs[slot_name]["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is None:
                    continue
                override = styles_override.get(slot_name)
                replace_text_with_style(tf, new_text, override)

            # Очистить незаполненные обязательные слоты
            for slot_name, slot_def in slot_defs.items():
                if slot_name in slots_filled:
                    continue
                if slot_def.get("optional"):
                    continue
                shape_idx = slot_def["shape_idx"]
                tf = get_text_frame_by_shape_idx(actual, shape_idx)
                if tf is not None:
                    clear_text_frame(tf)

        # === PICTURES (новое в v6) ===
        for pic in plan_slide.get("pictures", []):
            ok = insert_picture(
                actual,
                pic["file"],
                pic.get("left_px", 0),
                pic.get("top_px", 0),
                pic.get("width_px", 100),
                pic.get("height_px", 100),
            )
            if ok:
                pictures_inserted += 1

    p.save(output_path)
    print(f"Saved {output_path}: {len(p.slides)} slides, {pictures_inserted} pictures inserted",
          file=sys.stderr)


def main():
    if len(sys.argv) < 4:
        print("Usage: build_v6.py <plan.json> <template.pptx> <output.pptx> [donor-slot-map.yaml]",
              file=sys.stderr)
        sys.exit(1)
    plan_p = sys.argv[1]
    tpl_p = sys.argv[2]
    out_p = sys.argv[3]
    donor_p = sys.argv[4] if len(sys.argv) > 4 else "pptx-skill/brand/donor-slot-map.yaml"
    build(plan_p, tpl_p, out_p, donor_p)


if __name__ == "__main__":
    main()
